import streamlit as st
import requests
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
from datetime import date, timedelta
import time
import json
import re
import urllib.parse

# ── Config ─────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="Dashboard Market Gamer", layout="wide", page_icon="🎮")

TN_TOKEN = st.secrets["TN_TOKEN"]
TN_STORE_ID = st.secrets["TN_STORE_ID"]
SHEET_ID = st.secrets.get("SHEET_ID", "1wY2KjSC8SX-nMQD7J43xrdSY0SgG8fJeL9d5I_02DdE")
GCP_CREDS = st.secrets.get("gcp_service_account", {})
ANTHROPIC_KEY = st.secrets.get("ANTHROPIC_KEY", "")
MP_ACCESS_TOKEN = st.secrets.get("MP_ACCESS_TOKEN", "")

# ── Design tokens ───────────────────────────────────────────────────────────────
MG_BG       = "#0a0a0b"
MG_SURF     = "#131316"
MG_SURF2    = "#1a1a1e"
MG_BORDER   = "#26272b"
MG_TEXT     = "#f7f2f2"
MG_MUTED    = "#8a8b90"
MG_DIM      = "#5a5b60"
MG_RED      = "#ff3130"
MG_RED_DIM  = "#cc1f1e"

COLORES = [MG_RED, "#009EE3", "#fbbf24", "#4ade80", "#a78bfa", "#f472b6"]

# ── Global CSS ──────────────────────────────────────────────────────────────────
st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Hanken+Grotesk:wght@400;500;600;700&family=Space+Mono:wght@400;700&display=swap');

/* ── Base fonts ── */
html, body {{ font-family: 'Hanken Grotesk', system-ui, sans-serif; }}
p, h1, h2, h3, h4, h5, h6, label, input, textarea,
[data-testid="stMarkdownContainer"],
[data-testid="stMarkdownContainer"] *,
[data-testid="stMetricLabel"] p,
[data-testid="stMetricValue"],
[data-testid="stSidebar"] label,
[data-testid="stCaption"] p,
[data-baseweb] label, [data-baseweb] input {{
    font-family: 'Hanken Grotesk', system-ui, sans-serif !important;
}}

/* ── Surfaces ── */
.stApp,
[data-testid="stAppViewContainer"],
section[data-testid="stMain"] {{
    background-color: {MG_BG} !important;
}}
[data-testid="stSidebar"] {{
    background-color: {MG_BG} !important;
    border-right: 1px solid {MG_BORDER} !important;
}}
[data-testid="stMainBlockContainer"] {{
    padding-top: 1.5rem !important;
}}

/* ── Sidebar nav ── */
[data-testid="stSidebar"] label p,
[data-testid="stSidebar"] .stCaption p,
[data-testid="stSidebar"] .stRadio label p {{
    color: {MG_MUTED} !important;
    font-size: 0.85rem !important;
}}
[data-testid="stSidebar"] [data-testid="stRadio"] div[role="radio"][aria-checked="true"] label p {{
    color: {MG_TEXT} !important;
    font-weight: 600 !important;
}}
[data-testid="stSidebar"] [data-testid="stRadio"] div[role="radio"][aria-checked="true"] {{
    border-left: 2px solid {MG_RED} !important;
    padding-left: 6px !important;
}}

/* ── Dividers ── */
hr {{ border-color: {MG_BORDER} !important; }}

/* ── Streamlit metrics ── */
[data-testid="stMetricLabel"] p {{
    color: {MG_MUTED} !important;
    font-size: 0.62rem !important;
    text-transform: uppercase !important;
    letter-spacing: 0.08em !important;
    font-weight: 600 !important;
    font-family: 'Space Mono', ui-monospace, monospace !important;
}}
[data-testid="stMetricValue"] {{
    color: {MG_TEXT} !important;
    font-size: 1.4rem !important;
    font-weight: 700 !important;
    line-height: 1.2 !important;
    white-space: nowrap !important;
    overflow: visible !important;
}}
[data-testid="stMetricDelta"] {{
    font-size: 0.75rem !important;
}}

/* ── Headers ── */
[data-testid="stMarkdownContainer"] h1,
[data-testid="stMarkdownContainer"] h2,
[data-testid="stMarkdownContainer"] h3 {{
    color: {MG_TEXT} !important;
}}

/* ── Expanders ── */
[data-testid="stExpander"] {{
    background-color: {MG_SURF} !important;
    border: 1px solid {MG_BORDER} !important;
    border-radius: 6px !important;
}}
[data-testid="stExpander"] summary {{
    color: {MG_MUTED} !important;
    font-size: 0.85rem !important;
}}
[data-testid="stExpander"] summary:hover {{
    color: {MG_TEXT} !important;
}}
[data-testid="stExpander"] summary svg {{
    fill: {MG_MUTED} !important;
}}

/* ── Buttons ── */
.stButton > button[kind="primary"],
.stButton > button[data-testid="baseButton-primary"] {{
    background-color: {MG_RED} !important;
    color: {MG_TEXT} !important;
    border: none !important;
    font-weight: 700 !important;
    border-radius: 4px !important;
    font-family: 'Hanken Grotesk', sans-serif !important;
    letter-spacing: 0.02em !important;
}}
.stButton > button[kind="primary"]:hover {{
    background-color: {MG_RED_DIM} !important;
}}
.stButton > button:not([kind="primary"]),
.stButton > button[data-testid="baseButton-secondary"] {{
    background-color: {MG_SURF} !important;
    color: {MG_TEXT} !important;
    border: 1px solid {MG_BORDER} !important;
    border-radius: 4px !important;
    font-family: 'Hanken Grotesk', sans-serif !important;
}}
.stButton > button:not([kind="primary"]):hover {{
    border-color: {MG_MUTED} !important;
    color: {MG_TEXT} !important;
}}

/* ── Inputs / selectbox / number_input ── */
[data-testid="stTextInput"] input,
[data-testid="stNumberInput"] input,
[data-testid="stSelectbox"] > div > div {{
    background-color: {MG_SURF} !important;
    color: {MG_TEXT} !important;
    border: 1px solid {MG_BORDER} !important;
    border-radius: 4px !important;
}}
[data-testid="stTextInput"] input:focus,
[data-testid="stNumberInput"] input:focus {{
    border-color: {MG_RED} !important;
    box-shadow: 0 0 0 1px {MG_RED} !important;
}}
[data-testid="stSelectbox"] > div > div:focus-within {{
    border-color: {MG_RED} !important;
    box-shadow: 0 0 0 1px {MG_RED} !important;
}}
[data-baseweb="select"] > div {{
    background-color: {MG_SURF} !important;
    border-color: {MG_BORDER} !important;
}}
[data-baseweb="select"] > div:hover {{
    border-color: {MG_MUTED} !important;
}}
[data-baseweb="popover"] ul {{
    background-color: {MG_SURF} !important;
    border: 1px solid {MG_BORDER} !important;
}}
[data-baseweb="popover"] li {{
    color: {MG_TEXT} !important;
}}
[data-baseweb="popover"] li:hover {{
    background-color: {MG_SURF2} !important;
}}
[data-baseweb="date-input"] input {{
    background-color: {MG_SURF} !important;
    color: {MG_TEXT} !important;
    border-color: {MG_BORDER} !important;
}}

/* ── Dataframes / tables ── */
[data-testid="stDataFrame"] table,
[data-testid="stDataFrameGlideDataEditor"] {{
    background-color: {MG_SURF} !important;
    color: {MG_TEXT} !important;
}}
[data-testid="stDataFrame"] th {{
    background-color: {MG_BG} !important;
    color: {MG_MUTED} !important;
    font-size: 0.68rem !important;
    text-transform: uppercase !important;
    letter-spacing: 0.08em !important;
    font-weight: 600 !important;
    font-family: 'Space Mono', ui-monospace, monospace !important;
    border-bottom: 1px solid {MG_BORDER} !important;
}}
[data-testid="stDataFrame"] td {{
    border-color: {MG_BORDER} !important;
    color: {MG_TEXT} !important;
}}
[data-testid="stDataFrame"] tr:hover td {{
    background-color: {MG_SURF2} !important;
}}

/* ── Progress bar ── */
[data-testid="stProgressBar"] > div > div {{
    background-color: {MG_RED} !important;
}}
[data-testid="stProgressBar"] > div {{
    background-color: {MG_SURF2} !important;
}}

/* ── Spinner ── */
[data-testid="stSpinner"] svg circle {{
    stroke: {MG_RED} !important;
}}

/* ── Alert boxes ── */
[data-testid="stAlert"] {{
    border-radius: 4px !important;
}}
[data-testid="stAlert"][kind="success"] {{
    background-color: #001a0d !important;
    border-left: 3px solid #4ade80 !important;
}}
[data-testid="stAlert"][kind="warning"] {{
    background-color: #1a1200 !important;
    border-left: 3px solid #fbbf24 !important;
}}
[data-testid="stAlert"][kind="error"] {{
    background-color: #1a0000 !important;
    border-left: 3px solid {MG_RED} !important;
}}
[data-testid="stAlert"][kind="info"] {{
    background-color: #00101a !important;
    border-left: 3px solid #38bdf8 !important;
}}

/* ── Toggle ── */
[data-testid="stToggle"] svg {{
    color: {MG_RED} !important;
}}

/* ── Pills / multiselect ── */
[data-testid="stPills"] button[aria-pressed="true"] {{
    background-color: {MG_RED} !important;
    border-color: {MG_RED} !important;
    color: {MG_TEXT} !important;
}}
[data-testid="stPills"] button {{
    background-color: {MG_SURF} !important;
    border: 1px solid {MG_BORDER} !important;
    color: {MG_MUTED} !important;
    border-radius: 4px !important;
}}
[data-testid="stPills"] button:hover {{
    border-color: {MG_MUTED} !important;
    color: {MG_TEXT} !important;
}}

/* ── Chat messages ── */
[data-testid="stChatMessage"] {{
    background-color: {MG_SURF} !important;
    border: 1px solid {MG_BORDER} !important;
    border-radius: 6px !important;
}}
[data-testid="stChatInputTextArea"] textarea {{
    background-color: {MG_SURF} !important;
    color: {MG_TEXT} !important;
    border-color: {MG_BORDER} !important;
}}
[data-testid="stChatInputTextArea"] textarea:focus {{
    border-color: {MG_RED} !important;
}}

/* ── Captions ── */
[data-testid="stCaption"] p {{
    color: {MG_DIM} !important;
    font-size: 0.72rem !important;
}}

/* ── Info/warning/success text colors ── */
.stSuccess, .stWarning, .stError, .stInfo {{
    border-radius: 4px !important;
}}
</style>
""", unsafe_allow_html=True)

# ── Plotly dark template ────────────────────────────────────────────────────────
_mg_template = go.layout.Template()
_mg_template.layout = go.Layout(
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor=MG_SURF,
    font=dict(family="Hanken Grotesk, system-ui, sans-serif", color=MG_TEXT, size=12),
    title=dict(font=dict(size=13, color=MG_TEXT, family="Hanken Grotesk, sans-serif"), x=0, pad=dict(l=0)),
    xaxis=dict(gridcolor=MG_BORDER, linecolor=MG_BORDER, tickfont=dict(color=MG_MUTED, size=11), title_font=dict(color=MG_MUTED)),
    yaxis=dict(gridcolor=MG_BORDER, linecolor=MG_BORDER, tickfont=dict(color=MG_MUTED, size=11), title_font=dict(color=MG_MUTED)),
    legend=dict(bgcolor="rgba(0,0,0,0)", font=dict(color=MG_MUTED, size=11), borderwidth=0),
    margin=dict(l=8, r=8, t=36, b=8),
    colorway=COLORES,
)
pio.templates["mg_dark"] = _mg_template
pio.templates.default = "mg_dark"

# ── Page eyebrow ────────────────────────────────────────────────────────────────
st.markdown(
    f'<p style="font-size:0.68rem;color:{MG_DIM};margin:0 0 0.75rem 0;'
    f'letter-spacing:0.12em;font-family:\'Space Mono\',monospace;font-weight:700;">'
    f'MARKET GAMER · DASHBOARD</p>',
    unsafe_allow_html=True,
)

# ── Helpers ────────────────────────────────────────────────────────────────────
def fmt(n):
    try:
        return f"${float(n):,.0f}".replace(",", ".")
    except:
        return "-"

def fmt_pct(n):
    return f"{n:.2f}%"

def _normalizar(s):
    # Colapsar espacios y eliminar espacios entre letras/números del modelo
    # ej: "RG 477M 12+256GB" → "rg477m 12+256gb"  (espacios dentro del modelo se quitan)
    s = re.sub(r'\s+', ' ', str(s).strip().lower())
    # Quitar espacios entre partes alfanuméricas del modelo (no entre modelo y variante)
    s = re.sub(r'([a-z\d])\s+([a-z\d])', r'\1\2', s)
    return s

_NOISE_WORDS = ('almacenamiento', 'transparente', 'ram', 'negro', 'blanco',
                'azul', 'rojo', 'naranja', 'verde', 'gris', 'violeta',
                'purpura', 'rosa', 'dorado', 'plateado')

def _norm_compact(s):
    """Normalización agresiva: solo alfanumérico, sin colores ni 'GB/RAM/etc'.
    Ej: 'ANBERNIC RG 477 V (12GB RAM + 256GB Almacenamiento)' → 'anbernicrg477v12256'
        'Anbernic RG 477 V 12+256'                             → 'anbernicrg477v12256'
    """
    s = re.sub(r'[^a-z0-9]', '', str(s).lower())
    s = s.replace('gb', '').replace('tb', '')
    for w in _NOISE_WORDS:
        s = s.replace(w, '')
    return s

# ── Brand catalog (fuente: catalogo_market_gamer.csv + modelos TN) ─────────────
BRAND_CATALOG = {
    # ── Anbernic (bare model keys) ──────────────────────────────────────────────
    "rgnano":"Anbernic","rg28xx":"Anbernic","rg34xx":"Anbernic","rg34xxsp":"Anbernic",
    "rg35xx":"Anbernic","rg35xx2024":"Anbernic","rg35xxplus":"Anbernic",
    "rg35xxh":"Anbernic","rg35xxsp":"Anbernic","rg35xxpro":"Anbernic",
    "rg40xxh":"Anbernic","rg40xxv":"Anbernic",
    "rgcube":"Anbernic","rgcubexx":"Anbernic",
    "rg351p":"Anbernic","rg353p":"Anbernic","rg353v":"Anbernic",
    "rg353vs":"Anbernic","rg353m":"Anbernic","rg353ps":"Anbernic",
    "rgarcd":"Anbernic","rgarcs":"Anbernic","rg503":"Anbernic",
    "rg405m":"Anbernic","rg405v":"Anbernic","rg505":"Anbernic",
    "rgds":"Anbernic","rgp01":"Anbernic","rgg01":"Anbernic",
    "rg556":"Anbernic","rg557":"Anbernic",
    "rg406v":"Anbernic","rg406h":"Anbernic",
    "rgslide":"Anbernic","rg476h":"Anbernic",
    "rg477m":"Anbernic","rg477v":"Anbernic",
    "rgvita":"Anbernic","rgvitapro":"Anbernic",
    "win600":"Anbernic","k101plus":"Anbernic",
    # ── Anbernic (con prefijo de marca — para TN y proveedores que lo incluyen) ─
    "anbernicrgnano":"Anbernic","anbernicrg28xx":"Anbernic",
    "anbernicrg34xx":"Anbernic","anbernicrg34xxsp":"Anbernic",
    "anbernicrg35xx":"Anbernic","anbernicrg35xx2024":"Anbernic",
    "anbernicrg35xxplus":"Anbernic","anbernicrg35xxh":"Anbernic",
    "anbernicrg35xxsp":"Anbernic","anbernicrg35xxpro":"Anbernic",
    "anbernicrg40xxh":"Anbernic","anbernicrg40xxv":"Anbernic",
    "anbernicrgcube":"Anbernic","anbernicrgcubexx":"Anbernic",
    "anbernicrg351p":"Anbernic","anbernicrg353p":"Anbernic",
    "anbernicrg353v":"Anbernic","anbernicrg353vs":"Anbernic",
    "anbernicrg353m":"Anbernic","anbernicrg353ps":"Anbernic",
    "anbernicrgarcd":"Anbernic","anbernicrgarcs":"Anbernic","anbernicrg503":"Anbernic",
    "anbernicrg405m":"Anbernic","anbernicrg405v":"Anbernic","anbernicrg505":"Anbernic",
    "anbernicrgds":"Anbernic","anbernicrgp01":"Anbernic",
    "anbernicrg556":"Anbernic","anbernicrg557":"Anbernic",
    "anbernicrg406v":"Anbernic","anbernicrg406h":"Anbernic",
    "anbernicrgslide":"Anbernic","anbernicrg476h":"Anbernic",
    "anbernicrg477m":"Anbernic","anbernicrg477v":"Anbernic",
    "anbernicrgvita":"Anbernic","anbernicrgvitapro":"Anbernic",
    # ── Powkiddy (bare) ─────────────────────────────────────────────────────────
    "rgb10x":"Powkiddy","rgb20s":"Powkiddy","rgb20pro":"Powkiddy",
    "rgb20sx":"Powkiddy","rgb30":"Powkiddy",
    "rgb10max3":"Powkiddy","rgb10max3pro":"Powkiddy",
    "v10":"Powkiddy","v20":"Powkiddy","v90":"Powkiddy","v90s":"Powkiddy",
    "x28":"Powkiddy","x35h":"Powkiddy","x35s":"Powkiddy",
    "x51":"Powkiddy","x55":"Powkiddy","x70":"Powkiddy","x39pro":"Powkiddy",
    "q20mini":"Powkiddy","q90":"Powkiddy",
    # ── Powkiddy (con prefijo) ───────────────────────────────────────────────────
    "powkiddyrgb10x":"Powkiddy","powkiddyrgb20s":"Powkiddy","powkiddyrgb20pro":"Powkiddy",
    "powkiddyrgb20sx":"Powkiddy","powkiddyrgb30":"Powkiddy",
    "powkiddyrgb10max3":"Powkiddy","powkiddyrgb10max3pro":"Powkiddy",
    "powkiddyv10":"Powkiddy","powkiddyv20":"Powkiddy",
    "powkiddyv90":"Powkiddy","powkiddyv90s":"Powkiddy",
    "powkiddyx28":"Powkiddy","powkiddyx35h":"Powkiddy","powkiddyx35s":"Powkiddy",
    "powkiddyx51":"Powkiddy","powkiddyx55":"Powkiddy","powkiddyx70":"Powkiddy",
    "powkiddyx39pro":"Powkiddy","powkiddyq20mini":"Powkiddy","powkiddyq90":"Powkiddy",
    # ── Miyoo (bare) ────────────────────────────────────────────────────────────
    "miyoominiv3":"Miyoo","miyoominiv4":"Miyoo",
    "miyoominiplus":"Miyoo","miniplus":"Miyoo",
    "miyoominiflip":"Miyoo","miniflip":"Miyoo",
    "miyooflipv2":"Miyoo","flipv2":"Miyoo",
    "miyooa30":"Miyoo",
    # ── Trimui (siempre con prefijo — "brick" solo es demasiado genérico) ────────
    "trimuimodels":"Trimui","trimuismart":"Trimui",
    "trimuismartpro":"Trimui","trimuismartpros":"Trimui",
    "trimuibrick":"Trimui","trimuibrickhammer":"Trimui",
    # ── Retroid (bare + con prefijo) ─────────────────────────────────────────────
    "pocket4":"Retroid","pocket4pro":"Retroid","pocket5":"Retroid",
    "pocket6":"Retroid","pocketg2":"Retroid","pocketminiv2":"Retroid",
    "pocketclassic":"Retroid","pocketflip2":"Retroid","flip2":"Retroid",
    "retroidpocket4":"Retroid","retroidpocket4pro":"Retroid","retroidpocket5":"Retroid",
    "retroidpocket6":"Retroid","retroidpocketg2":"Retroid","retroidpocketminiv2":"Retroid",
    "retroidpocketclassic":"Retroid","retroidpocketflip2":"Retroid",
    "retroidflip2":"Retroid","retroidg2":"Retroid",
}

def _get_brand(product_name: str) -> str:
    """Detecta la marca de un producto por nombre normalizado."""
    cleaned = re.sub(r'\s*\d+\+\d+\s*gb.*', '', product_name, flags=re.IGNORECASE).strip()
    norm = re.sub(r'[^a-z0-9]', '', cleaned.lower())
    return BRAND_CATALOG.get(norm, "—")

# ── Proveedores helpers ────────────────────────────────────────────────────────
from datetime import date as _date

def compute_catalog_diff(old_products: dict, new_products: dict) -> list:
    """Compara dos catálogos y devuelve lista de cambios (cambiado/nuevo/eliminado)."""
    diff = []
    old_keys = set(old_products.keys())
    new_keys = set(new_products.keys())
    for name in old_keys & new_keys:
        try:
            old_price = float(old_products[name].get("precio_usd", 0))
            new_price = float(new_products[name].get("precio_usd", 0))
        except (TypeError, ValueError):
            continue
        if abs(old_price - new_price) > 0.01:
            diff.append({
                "Producto": name,
                "Antes (USD)": old_price,
                "Después (USD)": new_price,
                "tipo": "cambiado",
            })
    for name in new_keys - old_keys:
        diff.append({
            "Producto": name,
            "Antes (USD)": None,
            "Después (USD)": float(new_products[name].get("precio_usd", 0)),
            "tipo": "nuevo",
        })
    for name in old_keys - new_keys:
        diff.append({
            "Producto": name,
            "Antes (USD)": float(old_products[name].get("precio_usd", 0)),
            "Después (USD)": None,
            "tipo": "eliminado",
        })
    return diff


def parse_pptx_catalog(file_bytes: bytes) -> list:
    """
    Extrae filas (Producto, FOB USD, Storage) de un PowerPoint de lista de precios.
    Soporta dos formatos:
    - Tablas (columnas: nombre, precio, storage)
    - Text boxes agrupados (ej. Powkiddy): "MODELO  NNpcs  STORAGE  $PRECIO"
    """
    import io, re as _re
    from pptx import Presentation
    rows = []
    seen = set()
    NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    PRODUCT_RE = _re.compile(
        r'([A-Za-z][A-Za-z0-9 ]*?)\s+\d+\s*pcs\s+([^$]*?)\$\s*(\d+\.?\d*)',
        _re.IGNORECASE,
    )

    def _slide_xml_text(slide):
        parts = []
        for t in slide._element.iter(f'{NS}t'):
            if t.text and t.text.strip():
                parts.append(t.text.strip())
        return ' '.join(parts)

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception:
        return []

    for slide in prs.slides:
        # ── Estrategia 1: tablas explícitas ──────────────────────────────────
        found_table = False
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            found_table = True
            table = shape.table
            for row_idx in range(len(table.rows)):
                cells = [cell.text.strip() for cell in table.rows[row_idx].cells]
                for i, cell in enumerate(cells):
                    try:
                        val = float(cell.replace("$", "").replace(",", "").replace("USD", "").strip())
                        if not (1 < val < 5000):
                            continue
                        nombre_parts = [
                            c for c in cells[:i]
                            if c and len(c) > 2 and not c.replace(".", "").isdigit()
                        ]
                        if not nombre_parts:
                            continue
                        nombre = nombre_parts[0]
                        if nombre in seen:
                            break
                        seen.add(nombre)
                        storage = "—"
                        for c in cells:
                            if any(s in c.upper() for s in ["64G", "128G", "256G", "32G", "16G", "512G", "1TB"]):
                                storage = c.strip()
                                break
                        rows.append({"Producto": nombre, "FOB (USD)": val, "Marca": "—", "Pantalla": "—", "CPU": "—", "Storage": storage})
                        break
                    except ValueError:
                        continue

        if found_table:
            continue

        # ── Estrategia 2: text boxes agrupados (formato Powkiddy/Trimui) ────────
        # Estructura: "MODELO NNpcs STORAGE $PRECIO / / MODELO2 NNpcs STORAGE $PRECIO2"
        # Algunos slides no usan "/" como separador (ej. "X28 50pcs $93 X55 50pcs $48").
        # Algoritmo por marcador "NNpcs":
        #   - Nombre: en el segmento ANTERIOR al pcs, tomar texto después del último precio.
        #     Si no hay precio previo, tomar el último word-group (split en /  o espacios).
        #   - Precio: primer $XX en el segmento POSTERIOR al pcs (hasta el siguiente pcs).
        slide_text = _slide_xml_text(slide)
        PCS_ITER = list(_re.finditer(r'\d+\s*pcs', slide_text, _re.IGNORECASE))
        NAME_MAP = {"HAMMER": "Brick Hammer", "SMART PRO S": "Smart Pro S",
                    "SMART PRO": "Smart Pro", "SMART": "Smart"}
        TRIMUI_MODELS = {"Smart", "Smart Pro", "Smart Pro S", "Brick", "Brick Hammer"}

        for idx, pcs_m in enumerate(PCS_ITER):
            # Segmento anterior: desde el fin del pcs previo (o inicio) hasta este pcs
            prev_end = PCS_ITER[idx - 1].end() if idx > 0 else 0
            seg_before = slide_text[prev_end:pcs_m.start()]

            # Extraer nombre: texto después del último precio en seg_before
            last_price = list(_re.finditer(r'\$\s*\d+\.?\d*', seg_before))
            if last_price:
                name_raw = seg_before[last_price[-1].end():].strip()
            else:
                parts = _re.split(r'\s*/\s*|\s{3,}', seg_before.strip())
                name_raw = parts[-1].strip() if parts else ''

            # Limpiar el nombre de slashes, storage y noise
            name_raw = _re.sub(r'^[\s/]+', '', name_raw)  # slashes/espacios al inicio
            name_raw = _re.sub(r'\b\d+\s*GB\S*\b', '', name_raw, flags=_re.IGNORECASE).strip()
            name_raw = _re.sub(r'\s+', ' ', name_raw).strip()

            if not name_raw or len(name_raw) < 2:
                continue
            if _re.fullmatch(r'[\d/\s.]+', name_raw):
                continue

            # Segmento posterior: desde este pcs hasta el siguiente (para encontrar el precio)
            seg_end = PCS_ITER[idx + 1].start() if idx + 1 < len(PCS_ITER) else len(slide_text)
            seg_after = slide_text[pcs_m.end():seg_end]
            # Buscar precio con $ (ej. "$45") o sin $ si el catálogo lo omitió (ej. "28")
            price_m = _re.search(r'\$\s*(\d+\.?\d*)', seg_after)
            if not price_m:
                # Fallback: número suelto entre 5 y 999 que no sea parte de storage/dimensiones
                price_m = _re.search(
                    r'(?<![GB\d])(?<!\d)\b(\d{2,3}(?:\.\d)?)\b(?!\s*(?:GB|g\b|mm|kg|pcs|\d))',
                    seg_after, _re.IGNORECASE
                )
            if not price_m:
                continue  # producto genuinamente sin precio en este catálogo
            try:
                val = float(price_m.group(1))
            except ValueError:
                continue
            if not (1 < val < 5000):
                continue

            name = NAME_MAP.get(name_raw.upper(), name_raw)
            brand = "Trimui" if name in TRIMUI_MODELS else "Powkiddy"
            full_name = f"{brand} {name}"

            if full_name in seen:
                continue
            seen.add(full_name)

            storage_m = _re.search(r'(\d+\s*GB\S*)', seg_after, _re.IGNORECASE)
            storage = storage_m.group(1) if storage_m else "—"
            rows.append({"Producto": full_name, "FOB (USD)": val, "Marca": brand,
                         "Pantalla": "—", "CPU": "—", "Storage": storage})

    return rows


def fuzzy_group_products(suppliers: dict, threshold: float = 0.80) -> pd.DataFrame:
    """
    Agrupa productos similares de distintos proveedores usando difflib.
    Devuelve DataFrame con Producto | Proveedor1 | Proveedor2 | ...

    Nota: usa un algoritmo greedy — cada producto se asigna al primer grupo
    con similitud >= threshold. El orden de iteración afecta los resultados.
    Para los 3 proveedores de Market Gamer (catálogos pequeños), esto es suficiente.
    """
    from difflib import SequenceMatcher

    sup_names = list(suppliers.keys())
    all_items = []
    for sup_name in sup_names:
        for prod_name, prod_info in suppliers[sup_name].get("productos", {}).items():
            all_items.append({
                "supplier": sup_name,
                "name": prod_name,
                "price": float(prod_info.get("precio_usd", 0)),
            })

    if not all_items:
        return pd.DataFrame()

    # Agrupar por similitud de nombre
    groups: list[list] = []
    for item in all_items:
        matched = None
        for group in groups:
            rep = group[0]["name"]
            if SequenceMatcher(None, _normalizar(item["name"]), _normalizar(rep)).ratio() >= threshold:
                matched = group
                break
        if matched is None:
            groups.append([item])
        else:
            matched.append(item)

    rows = []
    for group in groups:
        canonical = group[0]["name"]
        row: dict = {"Producto": canonical}
        for sup in sup_names:
            match = next((x for x in group if x["supplier"] == sup), None)
            row[sup] = match["price"] if match else None
        rows.append(row)

    return pd.DataFrame(rows).sort_values("Producto")


@st.cache_data(ttl=300)
def get_stock_for_planner() -> dict:
    """
    Devuelve {nombre_producto: stock_total} consultando la API de TN.
    Stock None = sin límite configurado.
    TTL 5 min.
    """
    products = get_tn_products()
    stock_dict: dict = {}
    for p in products:
        nombre_raw = p.get("name", {})
        nombre = (
            nombre_raw.get("es", "") if isinstance(nombre_raw, dict) else str(nombre_raw)
        ).strip()
        if not nombre:
            continue
        # Si cualquier variante no tiene límite de stock, tratamos el producto como sin límite.
        # Así el planificador no lo marca como urgente.
        total = 0
        unlimited = False
        for v in p.get("variants", []):
            s = v.get("stock", None)
            if s is None:
                unlimited = True
                break
            total += int(s)
        stock_dict[nombre] = None if unlimited else total  # None = sin límite
    return stock_dict

# ── Google Sheets ──────────────────────────────────────────────────────────────
@st.cache_resource
def get_gsheet_client():
    try:
        import gspread
        from google.oauth2.service_account import Credentials
        scopes = [
            "https://spreadsheets.google.com/feeds",
            "https://www.googleapis.com/auth/drive",
        ]
        creds = Credentials.from_service_account_info(dict(GCP_CREDS), scopes=scopes)
        return gspread.authorize(creds)
    except Exception:
        return None

def gs_read(sheet_name):
    try:
        gc = get_gsheet_client()
        if not gc or not SHEET_ID:
            return {}
        ws = gc.open_by_key(SHEET_ID).worksheet(sheet_name)
        data = ws.get_all_values()
        if len(data) >= 2 and data[1]:
            return json.loads(data[1][0])
        return {}
    except Exception:
        return {}

def gs_write(sheet_name, data_dict):
    try:
        gc = get_gsheet_client()
        if not gc or not SHEET_ID:
            return False
        sh = gc.open_by_key(SHEET_ID)
        try:
            ws = sh.worksheet(sheet_name)
        except Exception:
            ws = sh.add_worksheet(sheet_name, rows=10, cols=2)
        ws.clear()
        ws.update("A1", [["key"], [json.dumps(data_dict)]])
        return True
    except Exception:
        return False

def gs_append_snapshot(stock_map):
    """Guarda el snapshot de stock de hoy en HistorialStock (idempotente por fecha)."""
    from velocidad_restock import merge_snapshot, recortar_historial
    try:
        hoy = date.today().isoformat()
        hist = gs_read("HistorialStock") or {}
        hist = merge_snapshot(hist, hoy, stock_map)
        hist = recortar_historial(hist, max_dias=180)
        return gs_write("HistorialStock", hist)
    except Exception:
        return False

# ── Google Analytics 4 ─────────────────────────────────────────────────────────
@st.cache_resource
def get_ga4_client():
    try:
        from google.analytics.data_v1beta import BetaAnalyticsDataClient
        from google.oauth2 import service_account
        creds = service_account.Credentials.from_service_account_info(
            dict(GCP_CREDS),
            scopes=["https://www.googleapis.com/auth/analytics.readonly"],
        )
        return BetaAnalyticsDataClient(credentials=creds)
    except Exception as e:
        st.warning(f"⚠️ Error conectando a GA4: {e}")
        return None

def _ga4_date_ranges(periodo, DateRange):
    """Devuelve (rango_actual, rango_previo, label_actual, label_previo) según período."""
    hoy = date.today()
    ayer = hoy - timedelta(days=1)
    if periodo == "7d":
        return (
            DateRange(start_date="7daysAgo", end_date="yesterday"),
            DateRange(start_date="14daysAgo", end_date="8daysAgo"),
            "7d", "7d previos",
        )
    if periodo == "28d":
        return (
            DateRange(start_date="28daysAgo", end_date="yesterday"),
            DateRange(start_date="56daysAgo", end_date="29daysAgo"),
            "28d", "28d previos",
        )
    if periodo == "MoM":
        inicio_mes = hoy.replace(day=1)
        dias_transcurridos = (ayer - inicio_mes).days  # 0..30
        inicio_mes_ant = (inicio_mes - timedelta(days=1)).replace(day=1)
        fin_mes_ant_equiv = inicio_mes_ant + timedelta(days=dias_transcurridos)
        return (
            DateRange(start_date=inicio_mes.isoformat(), end_date=ayer.isoformat()),
            DateRange(start_date=inicio_mes_ant.isoformat(), end_date=fin_mes_ant_equiv.isoformat()),
            f"Mes actual ({inicio_mes.strftime('%d/%m')}–{ayer.strftime('%d/%m')})",
            f"Mes anterior (mismos días)",
        )
    if periodo == "YoY":
        inicio = hoy - timedelta(days=28)
        inicio_yoy = inicio.replace(year=inicio.year - 1)
        ayer_yoy = ayer.replace(year=ayer.year - 1)
        return (
            DateRange(start_date=inicio.isoformat(), end_date=ayer.isoformat()),
            DateRange(start_date=inicio_yoy.isoformat(), end_date=ayer_yoy.isoformat()),
            "Últimos 28d", f"Mismos 28d {inicio_yoy.year}",
        )
    # fallback
    return _ga4_date_ranges("28d", DateRange)


@st.cache_data(ttl=1800)
def get_ga4_metrics(periodo="28d"):
    """KPIs, canales, páginas, productos, categorías, devices, OS, resolución, landing.
    Comparativo: 7d / 28d / MoM / YoY."""
    property_id = st.secrets.get("GA4_PROPERTY_ID")
    if not property_id:
        return None
    client = get_ga4_client()
    if not client:
        return None

    try:
        from google.analytics.data_v1beta.types import (
            RunReportRequest, Dimension, Metric, DateRange, OrderBy,
            FilterExpression, Filter,
        )
    except Exception as e:
        st.warning(f"⚠️ Falta dependencia google-analytics-data: {e}")
        return None

    prop = f"properties/{property_id}"
    rango_actual, rango_previo, label_act, label_prev = _ga4_date_ranges(periodo, DateRange)

    def _val(row, idx):
        try:
            return float(row.metric_values[idx].value or 0)
        except Exception:
            return 0.0

    try:
        # KPIs: sessions, activeUsers, engagementRate, averageSessionDuration, bounceRate
        kpi_req = RunReportRequest(
            property=prop,
            metrics=[
                Metric(name="sessions"),
                Metric(name="activeUsers"),
                Metric(name="engagementRate"),
                Metric(name="averageSessionDuration"),
                Metric(name="bounceRate"),
            ],
            date_ranges=[rango_actual, rango_previo],
        )
        kpi_resp = client.run_report(kpi_req)
        # Con dos date_ranges, cada fila trae dimensión "dateRange" → "date_range_0/1"
        sesiones_act = usuarios_act = interaccion_act = duracion_act = bounce_act = 0.0
        sesiones_prev = usuarios_prev = interaccion_prev = duracion_prev = bounce_prev = 0.0
        for row in kpi_resp.rows:
            etiqueta = row.dimension_values[0].value if row.dimension_values else ""
            if etiqueta == "date_range_0":
                sesiones_act = _val(row, 0)
                usuarios_act = _val(row, 1)
                interaccion_act = _val(row, 2)
                duracion_act = _val(row, 3)
                bounce_act = _val(row, 4)
            else:
                sesiones_prev = _val(row, 0)
                usuarios_prev = _val(row, 1)
                interaccion_prev = _val(row, 2)
                duracion_prev = _val(row, 3)
                bounce_prev = _val(row, 4)

        # Checkouts iniciados: pageviews donde pagePath contiene /checkout/v3/start
        checkout_filter = FilterExpression(
            filter=Filter(
                field_name="pagePath",
                string_filter=Filter.StringFilter(
                    value="/checkout/v3/start",
                    match_type=Filter.StringFilter.MatchType.CONTAINS,
                ),
            )
        )
        checkout_req = RunReportRequest(
            property=prop,
            metrics=[Metric(name="screenPageViews")],
            date_ranges=[rango_actual, rango_previo],
            dimension_filter=checkout_filter,
        )
        checkout_resp = client.run_report(checkout_req)
        checkouts_act = checkouts_prev = 0.0
        for row in checkout_resp.rows:
            etiqueta = row.dimension_values[0].value if row.dimension_values else ""
            if etiqueta == "date_range_0":
                checkouts_act = _val(row, 0)
            else:
                checkouts_prev = _val(row, 0)

        # Sesiones por canal (con variación entre rangos)
        canal_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="sessionDefaultChannelGroup")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual, rango_previo],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="sessions"), desc=True)],
        )
        canal_resp = client.run_report(canal_req)
        canales_map = {}
        for row in canal_resp.rows:
            canal = row.dimension_values[0].value if len(row.dimension_values) > 0 else "(sin dato)"
            etiqueta = row.dimension_values[1].value if len(row.dimension_values) > 1 else ""
            sesiones = _val(row, 0)
            entry = canales_map.setdefault(canal, {"canal": canal, "sesiones_act": 0.0, "sesiones_prev": 0.0})
            if etiqueta == "date_range_0":
                entry["sesiones_act"] = sesiones
            else:
                entry["sesiones_prev"] = sesiones
        canales = sorted(canales_map.values(), key=lambda x: x["sesiones_act"], reverse=True)

        # Páginas (top 100, rango actual) — incluye engagement time por página
        # Sirve para top páginas, top categorías agrupadas y top productos
        paginas_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="pagePath")],
            metrics=[
                Metric(name="screenPageViews"),
                Metric(name="userEngagementDuration"),
            ],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)],
            limit=100,
        )
        paginas_resp = client.run_report(paginas_req)
        paginas_raw = [
            {
                "pagina": row.dimension_values[0].value,
                "vistas": _val(row, 0),
                "engagement_seg": _val(row, 1),
            }
            for row in paginas_resp.rows
        ]
        # Tiempo prom. por página en segundos = engagement total / vistas
        for p in paginas_raw:
            p["tiempo_prom_seg"] = (p["engagement_seg"] / p["vistas"]) if p["vistas"] > 0 else 0.0
        top_paginas = paginas_raw[:30]

        # Top productos: filtrar pagePath que empiece con /productos/
        top_productos = sorted(
            [p for p in paginas_raw if p["pagina"].startswith("/productos/")],
            key=lambda x: x["vistas"], reverse=True,
        )[:5]

        # Top categorías: agrupar por primer segmento del path + URL representativa
        def _bucket_categoria(path):
            if not path or path == "/":
                return "Home"
            seg = path.strip("/").split("/", 1)[0].lower()
            mapping = {
                "consolas-portatiles": "Consolas portátiles (URL custom)",
                "consolas-de-mesa": "Consolas de mesa (URL custom)",
                "accesorios": "Accesorios (URL custom)",
                "juegos": "Juegos (URL custom)",
                "productos": "Productos (fichas individuales)",
                "categoria": "Páginas de marca (/categoria/*)",
                "checkout": "Checkout",
                "cuenta": "Cuenta",
                "carrito": "Carrito",
                "blog": "Blog",
                "hotsale": "Hotsale (campaña)",
                "search": "Búsqueda interna",
                "consolas": "Consolas (URL custom)",
            }
            return mapping.get(seg, seg.replace("-", " ").title())

        cat_map = {}
        for p in paginas_raw:
            cat = _bucket_categoria(p["pagina"])
            entry = cat_map.setdefault(cat, {
                "categoria": cat, "vistas": 0.0, "engagement_seg": 0.0,
                "url_top": p["pagina"], "vistas_top": p["vistas"],
            })
            entry["vistas"] += p["vistas"]
            entry["engagement_seg"] += p["engagement_seg"]
            # Trackear la URL más vista del bucket como ejemplo
            if p["vistas"] > entry["vistas_top"]:
                entry["url_top"] = p["pagina"]
                entry["vistas_top"] = p["vistas"]
        for c in cat_map.values():
            c["tiempo_prom_seg"] = (c["engagement_seg"] / c["vistas"]) if c["vistas"] > 0 else 0.0
        top_categorias = sorted(cat_map.values(), key=lambda x: x["vistas"], reverse=True)[:10]

        # Dispositivo (mobile/desktop/tablet)
        device_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="deviceCategory")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="sessions"), desc=True)],
        )
        device_resp = client.run_report(device_req)
        dispositivos = [
            {"device": row.dimension_values[0].value or "(sin dato)", "sesiones": _val(row, 0)}
            for row in device_resp.rows
        ]

        # Marcas: cada /categoria/<marca> con vistas + engagement
        marca_filter = FilterExpression(
            filter=Filter(
                field_name="pagePath",
                string_filter=Filter.StringFilter(
                    value="/categoria/",
                    match_type=Filter.StringFilter.MatchType.BEGINS_WITH,
                ),
            )
        )
        marcas_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="pagePath")],
            metrics=[
                Metric(name="screenPageViews"),
                Metric(name="userEngagementDuration"),
            ],
            date_ranges=[rango_actual],
            dimension_filter=marca_filter,
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)],
            limit=50,
        )
        marcas_resp = client.run_report(marcas_req)
        marcas_map = {}
        for row in marcas_resp.rows:
            path = row.dimension_values[0].value or ""
            # /categoria/powkiddy/ o /categoria/powkiddy/pagina-2 → "powkiddy"
            partes = path.strip("/").split("/")
            if len(partes) < 2 or partes[0] != "categoria":
                continue
            marca = partes[1].lower()
            if not marca:
                continue
            vistas = _val(row, 0)
            engagement = _val(row, 1)
            entry = marcas_map.setdefault(marca, {"marca": marca, "vistas": 0.0, "engagement_seg": 0.0})
            entry["vistas"] += vistas
            entry["engagement_seg"] += engagement
        for m in marcas_map.values():
            m["tiempo_prom_seg"] = (m["engagement_seg"] / m["vistas"]) if m["vistas"] > 0 else 0.0
        marcas = sorted(marcas_map.values(), key=lambda x: x["vistas"], reverse=True)

        # Sistema operativo (Android/iOS/Windows/Mac)
        os_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="operatingSystem")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="sessions"), desc=True)],
        )
        os_resp = client.run_report(os_req)
        sistemas_op = [
            {"os": row.dimension_values[0].value or "(sin dato)", "sesiones": _val(row, 0)}
            for row in os_resp.rows
        ]

        # Resolución de pantalla (para decisiones de breakpoints)
        res_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="screenResolution")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="sessions"), desc=True)],
            limit=10,
        )
        res_resp = client.run_report(res_req)
        resoluciones = [
            {"resolucion": row.dimension_values[0].value or "(sin dato)", "sesiones": _val(row, 0)}
            for row in res_resp.rows
        ]

        # Landing pages: por dónde entran
        landing_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="landingPage")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="sessions"), desc=True)],
            limit=5,
        )
        landing_resp = client.run_report(landing_req)
        landing_pages = [
            {"pagina": row.dimension_values[0].value or "(sin dato)", "sesiones": _val(row, 0)}
            for row in landing_resp.rows
        ]

        # Productos vistos (todos, para tabla larga buscable)
        productos_filter = FilterExpression(
            filter=Filter(
                field_name="pagePath",
                string_filter=Filter.StringFilter(
                    value="/productos/",
                    match_type=Filter.StringFilter.MatchType.BEGINS_WITH,
                ),
            )
        )
        productos_full_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="pagePath"), Dimension(name="pageTitle")],
            metrics=[
                Metric(name="screenPageViews"),
                Metric(name="userEngagementDuration"),
            ],
            date_ranges=[rango_actual],
            dimension_filter=productos_filter,
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)],
            limit=500,
        )
        productos_full_resp = client.run_report(productos_full_req)
        productos_full = []
        for row in productos_full_resp.rows:
            path = row.dimension_values[0].value
            titulo = row.dimension_values[1].value if len(row.dimension_values) > 1 else ""
            vistas = _val(row, 0)
            engagement = _val(row, 1)
            productos_full.append({
                "path": path,
                "titulo": titulo,
                "vistas": vistas,
                "engagement_seg": engagement,
                "tiempo_prom_seg": (engagement / vistas) if vistas > 0 else 0.0,
            })

        # Checkouts por canal (atribución básica): qué canal trae a los que inician checkout
        checkouts_canal_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="sessionDefaultChannelGroup")],
            metrics=[Metric(name="screenPageViews")],
            date_ranges=[rango_actual],
            dimension_filter=checkout_filter,
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)],
        )
        checkouts_canal_resp = client.run_report(checkouts_canal_req)
        checkouts_por_canal = [
            {"canal": row.dimension_values[0].value or "(sin dato)", "checkouts": _val(row, 0)}
            for row in checkouts_canal_resp.rows
        ]

        # Tendencia diaria: sesiones + checkouts por día (rango actual)
        tendencia_sesiones_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="date")],
            metrics=[Metric(name="sessions")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(dimension=OrderBy.DimensionOrderBy(dimension_name="date"))],
        )
        tendencia_sesiones_resp = client.run_report(tendencia_sesiones_req)
        tendencia_sesiones = [
            {"fecha": row.dimension_values[0].value, "sesiones": _val(row, 0)}
            for row in tendencia_sesiones_resp.rows
        ]

        tendencia_checkouts_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="date")],
            metrics=[Metric(name="screenPageViews")],
            date_ranges=[rango_actual],
            dimension_filter=checkout_filter,
            order_bys=[OrderBy(dimension=OrderBy.DimensionOrderBy(dimension_name="date"))],
        )
        tendencia_checkouts_resp = client.run_report(tendencia_checkouts_req)
        tendencia_checkouts = {
            row.dimension_values[0].value: _val(row, 0)
            for row in tendencia_checkouts_resp.rows
        }
        # Mergear ambas series por fecha
        for d in tendencia_sesiones:
            d["checkouts"] = tendencia_checkouts.get(d["fecha"], 0.0)

        # Búsquedas internas: cualquier URL con ?q= ?_q= ?query= ?s= (TN puede usar varios formatos)
        busquedas_internas = []
        urls_con_query = []  # fallback de debug: top URLs con query string
        try:
            import re as _re
            import urllib.parse as _ulib

            # Pedimos top URLs con query string (sin filtro restrictivo) y filtramos local
            urlq_req = RunReportRequest(
                property=prop,
                dimensions=[Dimension(name="pagePathPlusQueryString")],
                metrics=[Metric(name="screenPageViews")],
                date_ranges=[rango_actual],
                order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)],
                limit=500,
            )
            urlq_resp = client.run_report(urlq_req)
            for row in urlq_resp.rows:
                path_q = row.dimension_values[0].value or ""
                if "?" not in path_q:
                    continue
                vistas = _val(row, 0)
                urls_con_query.append({"url": path_q, "vistas": vistas})
                # Detectar parámetros de búsqueda comunes (TN: _q es el más típico)
                m = _re.search(r"[?&](?:_q|q|query|s|busqueda|buscar)=([^&]+)", path_q)
                if m:
                    try:
                        termino = _ulib.unquote_plus(m.group(1)).strip()
                    except Exception:
                        termino = m.group(1)
                    if termino and len(termino) >= 2:
                        busquedas_internas.append({"termino": termino.lower(), "vistas": vistas})

            # Consolidar por término
            bmap = {}
            for b in busquedas_internas:
                e = bmap.setdefault(b["termino"], {"termino": b["termino"], "vistas": 0.0})
                e["vistas"] += b["vistas"]
            busquedas_internas = sorted(bmap.values(), key=lambda x: x["vistas"], reverse=True)[:30]

            # Quedarnos solo con top 30 URLs con query para el debug
            urls_con_query = sorted(urls_con_query, key=lambda x: x["vistas"], reverse=True)[:30]
        except Exception:
            busquedas_internas = []
            urls_con_query = []

        # Top eventos GA4 (para investigar si TN dispara add_to_cart, view_item, etc.)
        eventos_req = RunReportRequest(
            property=prop,
            dimensions=[Dimension(name="eventName")],
            metrics=[Metric(name="eventCount")],
            date_ranges=[rango_actual],
            order_bys=[OrderBy(metric=OrderBy.MetricOrderBy(metric_name="eventCount"), desc=True)],
            limit=20,
        )
        eventos_resp = client.run_report(eventos_req)
        eventos = [
            {"evento": row.dimension_values[0].value, "count": _val(row, 0)}
            for row in eventos_resp.rows
        ]

        return {
            "label_actual": label_act,
            "label_previo": label_prev,
            "kpis": {
                "sesiones": (sesiones_act, sesiones_prev),
                "usuarios": (usuarios_act, usuarios_prev),
                "interaccion": (interaccion_act, interaccion_prev),
                "checkouts": (checkouts_act, checkouts_prev),
                "duracion": (duracion_act, duracion_prev),
                "bounce": (bounce_act, bounce_prev),
            },
            "canales": canales,
            "top_paginas": top_paginas,
            "top_productos": top_productos,
            "top_categorias": top_categorias,
            "dispositivos": dispositivos,
            "sistemas_op": sistemas_op,
            "resoluciones": resoluciones,
            "landing_pages": landing_pages,
            "productos_full": productos_full,
            "marcas": marcas,
            "checkouts_por_canal": checkouts_por_canal,
            "tendencia_diaria": tendencia_sesiones,
            "busquedas_internas": busquedas_internas,
            "urls_con_query": urls_con_query,
            "eventos": eventos,
        }

    except Exception as e:
        st.warning(f"⚠️ Error consultando GA4: {e}")
        return None

# ── API Tienda Nube ────────────────────────────────────────────────────────────
def get_tn_headers():
    return {
        "Authentication": f"bearer {TN_TOKEN}",
        "User-Agent": "MarketGamerDashboard (info@marketgamer.com.ar)",
    }

def get_tn_orders(fecha_desde, fecha_hasta):
    headers = get_tn_headers()
    all_orders, seen_ids = [], set()
    combinaciones = [
        {"payment_status": "paid"},
        {"payment_status": "paid", "status": "archived"},
        {"payment_status": "paid", "status": "closed"},
    ]
    with st.spinner("Conectando con Tienda Nube..."):
        for filtros in combinaciones:
            page = 1
            while True:
                r = requests.get(
                    f"https://api.tiendanube.com/v1/{TN_STORE_ID}/orders",
                    headers=headers,
                    params={
                        "per_page": 50, "page": page,
                        "created_at_min": f"{fecha_desde}T00:00:00-03:00",
                        "created_at_max": f"{fecha_hasta}T23:59:59-03:00",
                        **filtros,
                    },
                )
                if r.status_code == 404:
                    break
                if r.status_code == 429:
                    time.sleep(3)
                    continue
                if r.status_code != 200:
                    break
                try:
                    data = r.json()
                except Exception:
                    break
                if not data:
                    break
                for o in data:
                    oid = o.get("id")
                    if oid not in seen_ids:
                        seen_ids.add(oid)
                        all_orders.append(o)
                if len(data) < 50:
                    break
                page += 1
    return all_orders

def get_tn_pagos(fecha_desde, fecha_hasta):
    headers = get_tn_headers()
    all_pagos, seen_ids = [], set()
    with st.spinner("Conectando con Pago Nube..."):
        page = 1
        while True:
            r = requests.get(
                f"https://api.tiendanube.com/v1/{TN_STORE_ID}/transactions",
                headers=headers,
                params={
                    "per_page": 50, "page": page,
                    "created_at_min": f"{fecha_desde}T00:00:00-03:00",
                    "created_at_max": f"{fecha_hasta}T23:59:59-03:00",
                },
            )
            if r.status_code in [404, 422]:
                break
            if r.status_code == 429:
                time.sleep(3)
                continue
            if r.status_code != 200:
                break
            try:
                data = r.json()
            except Exception:
                break
            if not data:
                break
            for p in data:
                pid = p.get("id")
                if pid not in seen_ids:
                    seen_ids.add(pid)
                    all_pagos.append(p)
            if len(data) < 50:
                break
            page += 1
    return all_pagos

def get_tn_products():
    headers = get_tn_headers()
    all_products = []
    page = 1
    while True:
        r = requests.get(
            f"https://api.tiendanube.com/v1/{TN_STORE_ID}/products",
            headers=headers,
            params={"per_page": 50, "page": page},
        )
        if r.status_code != 200:
            break
        try:
            data = r.json()
        except Exception:
            break
        if not data:
            break
        all_products.extend(data)
        if len(data) < 50:
            break
        page += 1
    return all_products

# ── Tasas Pago Nube ────────────────────────────────────────────────────────────
PROC_BASE = 0.0329
IVA_FACTOR = 1.2600
PROC_EFECTIVO = PROC_BASE * IVA_FACTOR

CUOTAS_BASE = {
    1: 0.0, 2: 0.0606, 3: 0.0798, 6: 0.1552,
    12: 0.3104, 18: 0.4346, 24: 0.5432,
}

def tasa_pago_nube(metodo, cuotas):  # noqa: keep for compatibilidad
    metodo = str(metodo).lower()
    if any(x in metodo for x in ["transfer", "wire", "account_money"]):
        return 0.0099 * IVA_FACTOR
    if any(x in metodo for x in ["debit", "debito", "modo"]):
        return PROC_EFECTIVO
    cuotas = int(cuotas or 1)
    opciones = sorted(CUOTAS_BASE.keys())
    cuotas_key = min(opciones, key=lambda x: abs(x - cuotas))
    costo_cuotas = CUOTAS_BASE.get(cuotas_key, 0.0) * IVA_FACTOR
    return PROC_EFECTIVO + costo_cuotas

# ── Tasas Mercado Pago ──────────────────────────────────────────────────────────
# Fuente: MP → Costos y cuotas → Checkout Pro (abr-2026)
# Fórmula: (financing_fee × 1.21 IVA) + 3.87% comisión base
COSTOS_MP_DEFAULTS = {
    "Transferencia": 0.0,
    "Contado":       3.87,
    "2 cuotas":      8.71,   # 4.00% × 1.21 + 3.87
    "3 cuotas":      11.13,  # 6.00% × 1.21 + 3.87
    "6 cuotas":      15.97,  # 10.00% × 1.21 + 3.87
    "9 cuotas":      22.62,  # 15.50% × 1.21 + 3.87
    "12 cuotas":     28.07,  # 20.00% × 1.21 + 3.87
}

def _es_gateway_mp(gateway):
    """True si la orden pasó por Mercado Pago (no Pago Nube)."""
    gw = str(gateway).lower().replace(" ", "").replace("-", "").replace("_", "")
    return "mercadopago" in gw or gw in {"mp", "mercadopagov1", "mercadopagocheckoutpro"}

def _es_convenir(gateway, metodo):
    """True si la orden fue creada con 'pago a convenir' (link MP generado por fuera de TN)."""
    gw = str(gateway).lower().strip()
    mt = str(metodo).lower().strip()
    if "convenir" in gw or "convenir" in mt:
        return True
    # Gateway vacío + método no reconocido → probablemente a convenir
    _metodos_conocidos = {"credit_card", "debit_card", "account_money"}
    _frags_conocidos   = ["transfer", "wire", "pago", "mercado", "credit", "debit"]
    if (not gw or gw in {"", "none", "null", "other", "offline", "manual"}) \
       and mt not in _metodos_conocidos \
       and not any(f in mt for f in _frags_conocidos):
        return True
    return False

def match_mp_with_tn(df_tn, mp_payments_raw):
    """
    Cruza TODAS las órdenes TN (excepto PN explícito) con pagos aprobados de MP.
    Estrategia: monto exacto (±$1) + fecha ±1 día.

    Por qué match a TODO no-PN:
    - 'Convenir' = link MP manual (pago externo a TN) → siempre matchear
    - 'MP' = gateway MP nativo → matchear para usar fee REAL en lugar del estimado
    - 'PN' = Pago Nube (con comisión real ya reportada por TN) → no tocar

    Devuelve (df_actualizado, n_matched, n_sin_match_pendientes).
    """
    if not mp_payments_raw or df_tn.empty:
        return df_tn, 0, 0

    # Construir índice MP: {(monto_redondeado, fecha): datos_fee}
    # Excluir liquidaciones externas (Dlocal/PN) del índice
    mp_index = {}
    for p in mp_payments_raw:
        # NO filtramos liquidaciones del índice. Razón: el CUIT del payer en
        # bank_transfers entrantes es el de MG (no el de TN), así que no es
        # confiable. Las liquidaciones de TN son sumas grandes que no van a
        # coincidir con monto exacto de ninguna orden individual, así que no
        # generan falsos matches. Las transferencias de clientes reales sí
        # coinciden y deben matchear (caso 449: Sabrina, $180k).
        bruto = float(p.get("transaction_amount", 0))
        try:
            fecha_mp = pd.to_datetime(p.get("date_approved")).date()
        except Exception:
            continue
        # Usar bruto - neto: captura TODOS los fees sin importar tipo
        neto        = float(p.get("transaction_details", {}).get("net_received_amount", 0))
        costo_total = round(bruto - neto, 2)
        cuotas      = int(p.get("installments", 1) or 1)
        mp_index[(round(bruto), fecha_mp)] = {
            "id_mp":        str(int(p.get("id", 0))),
            "comision_real":costo_total,
            "neto_real":    round(neto, 2),
            "cuotas_mp":    cuotas,
            "costo_pct":    round((costo_total / bruto * 100) if bruto > 0 else 0, 2),
        }

    df = df_tn.copy()
    if "ID MP" not in df.columns:
        df["ID MP"] = ""

    matched, sin_match = 0, 0
    matched_ids = set()  # cada pago MP se usa una sola vez

    # ESTRATEGIA: probar match para TODA orden, ordenadas por fecha descendente.
    # Si encuentra un pago MP con monto exacto y fecha ±1 día, usar el fee REAL.
    # Esto cubre:
    #   - Convenir (link MP manual fuera de TN)
    #   - MP nativo (reemplaza estimado por fee real)
    #   - Órdenes mal clasificadas como PN (gateway raros tipo 'other'/'offline')
    # Riesgo: orden PN con monto idéntico a pago MP del mismo día. Mitigado con
    # match exacto (sin tolerancia $) y matched_ids para evitar duplicados.

    for idx, row in df.iterrows():
        # No matchear órdenes marcadas explícitamente como Efectivo
        if row.get("Pasarela", "") == "Efectivo":
            continue
        total = round(float(row.get("Total ($)", 0)))
        try:
            fecha_tn = pd.to_datetime(row.get("Fecha")).date()
        except Exception:
            continue

        match_data = None
        for delta in [0, 1, -1]:
            key = (total, fecha_tn + timedelta(days=delta))
            if key in mp_index and mp_index[key]["id_mp"] not in matched_ids:
                match_data = mp_index[key]
                break

        if match_data:
            df.at[idx, "Comision PN ($)"]  = match_data["comision_real"]
            df.at[idx, "Neto cobrado ($)"] = match_data["neto_real"]
            df.at[idx, "Costo PN (%)"]     = match_data["costo_pct"]
            df.at[idx, "Cuotas"]           = match_data["cuotas_mp"]
            df.at[idx, "Pasarela"]         = "MP"
            df.at[idx, "ID MP"]            = match_data["id_mp"]
            matched_ids.add(match_data["id_mp"])
            # Recalcular margen con fee real
            costo_prods = float(row.get("Costo Productos ($)", 0))
            costo_envio = float(row.get("Envio costo ($)", 0))
            nuevo_margen = match_data["neto_real"] - costo_prods - costo_envio
            df.at[idx, "Margen ($)"]   = round(nuevo_margen, 2)
            total_val = float(row.get("Total ($)", 0))
            df.at[idx, "Margen (%)"]   = round(
                (nuevo_margen / total_val * 100) if total_val > 0 else 0, 2
            )
            matched += 1
        else:
            # Solo Convenir cuenta como pendiente — MP/PN ya tienen comisión válida
            if row.get("Pasarela", "") == "Convenir":
                sin_match += 1

    return df, matched, sin_match

def tasa_pasarela(gateway, metodo, cuotas):
    """Tasa de comisión real según pasarela: Mercado Pago o Pago Nube."""
    if _es_gateway_mp(gateway):
        m = str(metodo).lower()
        c = int(cuotas or 1)
        if any(x in m for x in ["transfer", "wire", "bank_transfer"]):
            clave = "Transferencia"
        elif "debit" in m or "debito" in m:
            clave = "Contado"
        elif c <= 1:
            clave = "Contado"
        elif c <= 2:
            clave = "2 cuotas"
        elif c <= 3:
            clave = "3 cuotas"
        elif c <= 6:
            clave = "6 cuotas"
        elif c <= 9:
            clave = "9 cuotas"
        else:
            clave = "12 cuotas"
        return COSTOS_MP_DEFAULTS.get(clave, 3.87) / 100
    else:
        return tasa_pago_nube(metodo, cuotas)

@st.cache_data(ttl=900, show_spinner=False)
def get_mp_payments(fecha_desde_str, fecha_hasta_str):
    """Trae pagos aprobados de MP para el período. Devuelve lista de dicts con fee real."""
    if not MP_ACCESS_TOKEN:
        return []
    headers = {"Authorization": f"Bearer {MP_ACCESS_TOKEN}"}
    all_payments, offset = [], 0
    while True:
        try:
            r = requests.get(
                "https://api.mercadopago.com/v1/payments/search",
                headers=headers,
                params={
                    "status": "approved",
                    "begin_date": f"{fecha_desde_str}T00:00:00.000-03:00",
                    "end_date":   f"{fecha_hasta_str}T23:59:59.999-03:00",
                    "limit": 100,
                    "offset": offset,
                },
                timeout=15,
            )
        except Exception:
            break
        if r.status_code != 200:
            break
        data = r.json()
        results = data.get("results", [])
        if not results:
            break
        all_payments.extend(results)
        offset += 100
        if offset >= data.get("paging", {}).get("total", 0):
            break
    return all_payments

# CUITs de procesadores de pago externos cuyas transferencias a la cuenta MP
# NO son ventas — ya fueron contadas como ventas vía Pago Nube en TN.
CUITS_LIQUIDACIONES = {
    "30718824830",  # Tienda Nube / Pago Nube (liquidaciones a la cuenta MP)
}

def _cuit_payer(p):
    """Devuelve el CUIT/CUIL del payer del pago, limpio (solo dígitos)."""
    candidatos = [
        p.get("payer", {}).get("identification", {}).get("number", ""),
        p.get("additional_info", {}).get("payer", {}).get("identification", {}).get("number", ""),
    ]
    for c in candidatos:
        digits = "".join(ch for ch in str(c) if ch.isdigit())
        if digits:
            return digits
    return ""

def _es_liquidacion_externa(p):
    """Detecta liquidación de procesador externo (PN/Dlocal) en cuenta MP.
    Criterio principal: bank_transfer cuyo CUIT del payer está en CUITS_LIQUIDACIONES.
    Fallback: bank_transfer con 'dlocal' en description/statement_descriptor.
    """
    if p.get("payment_type_id", "") != "bank_transfer":
        return False
    if _cuit_payer(p) in CUITS_LIQUIDACIONES:
        return True
    textos = [
        str(p.get("description", "")),
        str(p.get("statement_descriptor", "")),
    ]
    return any("dlocal" in t.lower() for t in textos)

_FEE_TYPE_LABELS = {
    "mercadopago_fee":         "Cargo MP",
    "financing_fee":           "Cargo financiación",
    "application_fee":         "Cargo plataforma",   # plataforma de terceros (TN)
    "discount_fee":            "Descuento aplicado",
    "fixed_fee":               "Cargo fijo",
    "iva_fee":                 "IVA",
    "absorbed_financing_fee":  "Intereses absorbidos",
    "tax_fee":                 "Impuestos",
    "iibb":                    "IIBB",
    "credit_fee":              "Crédito",
}

_TAX_TYPES = {"iva_fee", "tax_fee", "iibb", "iibb_fee", "tax", "iva", "ingresos_brutos"}

def _desglose_fees(p):
    """Devuelve dict {label: monto} con desglose de fee_details + impuestos.
    Para análisis fino: cargo MP, IVA, IIBB, plataforma, etc.
    """
    desglose = {}
    # 1) fee_details — lista de {type, amount, fee_payer}
    for f in p.get("fee_details", []) or []:
        tipo = str(f.get("type", "")).lower()
        try:
            monto = float(f.get("amount", 0) or 0)
        except Exception:
            monto = 0.0
        # Si es impuesto, agruparlo bajo 'Impuestos'
        if tipo in _TAX_TYPES:
            desglose["Impuestos"] = round(desglose.get("Impuestos", 0) + monto, 2)
        else:
            label = _FEE_TYPE_LABELS.get(tipo, tipo or "Otro")
            desglose[label] = round(desglose.get(label, 0) + monto, 2)
    # 2) charges_details — cargos extra (plataforma de terceros, etc.)
    for c in p.get("charges_details", []) or []:
        tipo = str(c.get("type", "")).lower()
        amounts = c.get("amounts", {}) or {}
        try:
            monto = float(amounts.get("original", 0) or 0)
        except Exception:
            monto = 0.0
        if tipo in _TAX_TYPES:
            desglose["Impuestos"] = round(desglose.get("Impuestos", 0) + monto, 2)
        else:
            label = _FEE_TYPE_LABELS.get(tipo, tipo or "Otro")
            desglose[label] = round(desglose.get(label, 0) + monto, 2)
    # 3) Array 'taxes' a nivel raíz (MP a veces lo expone separado)
    for t in p.get("taxes", []) or []:
        try:
            monto = float(t.get("value", 0) or 0)
        except Exception:
            monto = 0.0
        desglose["Impuestos"] = round(desglose.get("Impuestos", 0) + monto, 2)
    # 4) taxes_amount en transaction_details (campo legacy)
    taxes_extra = (
        p.get("taxes_amount", 0)
        or p.get("transaction_details", {}).get("taxes_amount", 0)
        or 0
    )
    if taxes_extra:
        desglose["Impuestos"] = round(desglose.get("Impuestos", 0) + float(taxes_extra), 2)
    return desglose

def procesar_mp_payments(payments, montos_validos_tn=None):
    """Convierte lista raw de MP en DataFrame con fee real por operación.
    Excluye liquidaciones de procesadores externos: bank_transfer cuyo monto
    NO coincide con ninguna orden TN del período (= no es venta, es liquidación).

    Args:
        montos_validos_tn: set de montos (round) de órdenes TN del período.
                           Si se pasa, los bank_transfer que no estén ahí se excluyen.
                           Si es None, usa el filtro legacy por CUIT/descripción.
    """
    filas = []
    for p in payments:
        # Filtro mejorado de liquidaciones
        if p.get("payment_type_id", "") == "bank_transfer":
            bruto_p = round(float(p.get("transaction_amount", 0)))
            if montos_validos_tn is not None:
                # Si el monto NO coincide con ninguna orden TN, es liquidación
                if bruto_p not in montos_validos_tn:
                    continue
            else:
                # Fallback al filtro viejo (por CUIT)
                if _es_liquidacion_externa(p):
                    continue
        cuotas    = p.get("installments", 1)
        bruto     = float(p.get("transaction_amount", 0))
        neto      = float(p.get("transaction_details", {}).get("net_received_amount", 0))
        # Fee real = bruto - neto (captura todos los tipos de fee de MP)
        costo_total = round(bruto - neto, 2)
        tipo_pago = p.get("payment_type_id", "")
        if tipo_pago == "bank_transfer":
            tipo_label = "Transferencia"
        elif cuotas == 1:
            tipo_label = "Contado"
        else:
            tipo_label = f"{cuotas} cuotas"
        try:
            fecha = pd.to_datetime(p.get("date_approved")).strftime("%Y-%m-%d")
        except Exception:
            fecha = ""

        # Desglose detallado de fees (impuestos, cargos, intereses absorbidos, etc.)
        desg = _desglose_fees(p)
        filas.append({
            "ID MP": str(int(p.get("id", 0))),
            "Fecha": fecha,
            "Tipo": tipo_label,
            "Cuotas": cuotas,
            "Medio": p.get("payment_method_id", "").upper(),
            "Bruto ($)": round(bruto, 2),
            "Fee total ($)": costo_total,
            "Costo %": round((costo_total / bruto * 100) if bruto > 0 else 0, 2),
            "Neto ($)": round(neto, 2),
            "Cargo MP ($)":           round(desg.get("Cargo MP", 0), 2),
            "Cargo financiación ($)": round(desg.get("Cargo financiación", 0), 2),
            "Impuestos ($)":          round(desg.get("Impuestos", 0), 2),
            "Cargo plataforma ($)":   round(desg.get("Cargo plataforma", 0), 2),
        })
    return pd.DataFrame(filas) if filas else pd.DataFrame()

# ── FOB defaults ───────────────────────────────────────────────────────────────
FOB_DEFAULTS = {
    "Anbernic RG 34XX 64GB":      {"fob_usd": 59.00,  "peso_kg": 0.3400},
    "Anbernic RG 34XX SP 64GB":   {"fob_usd": 66.00,  "peso_kg": 0.3000},
    "Anbernic RG 35XX Pro 64GB":  {"fob_usd": 49.00,  "peso_kg": 0.3450},
    "Anbernic RG 35XX SP 64GB":   {"fob_usd": 49.00,  "peso_kg": 0.3300},
    "Anbernic RG 406H":           {"fob_usd": 136.34, "peso_kg": 0.4652},
    "Anbernic RG 406V":           {"fob_usd": 139.70, "peso_kg": 0.4900},
    "Anbernic RG 40XX H 64GB":    {"fob_usd": 47.30,  "peso_kg": 0.3730},
    "Anbernic RG 477M 128GB":     {"fob_usd": 284.00, "peso_kg": 0.6400},
    "Anbernic RG 557 128GB":      {"fob_usd": 264.00, "peso_kg": 0.6070},
    "Anbernic RG Cube 128GB":     {"fob_usd": 164.00, "peso_kg": 0.4630},
    "Anbernic RG Cube XX 64GB":   {"fob_usd": 59.00,  "peso_kg": 0.4430},
    "Anbernic RG P01 Blanco":     {"fob_usd": 12.70,  "peso_kg": 0.4300},
    "Anbernic RG P01 Negro":      {"fob_usd": 12.70,  "peso_kg": 0.4300},
    "Anbernic RG Slide 128GB":    {"fob_usd": 174.00, "peso_kg": 0.6500},
    "Anbernic RG40XX H":          {"fob_usd": 50.60,  "peso_kg": 0.3600},
    "Anbernic RG40XX V":          {"fob_usd": 48.40,  "peso_kg": 0.3700},
    "Miyoo A30":                  {"fob_usd": 34.50,  "peso_kg": 0.2200},
    "Miyoo Flip":                 {"fob_usd": 60.00,  "peso_kg": 0.2800},
    "Miyoo Mini Plus":            {"fob_usd": 40.00,  "peso_kg": 0.2700},
    "Powkiddy MAX3":              {"fob_usd": 48.00,  "peso_kg": 0.3800},
    "Powkiddy MAX3 Pro":          {"fob_usd": 90.00,  "peso_kg": 0.5000},
    "Powkiddy RGB10X":            {"fob_usd": 30.00,  "peso_kg": 0.3000},
    "Powkiddy RGB20 Pro":         {"fob_usd": 45.00,  "peso_kg": 0.3800},
    "Powkiddy RGB20S":            {"fob_usd": 32.00,  "peso_kg": 0.3500},
    "Powkiddy RGB20SX":           {"fob_usd": 48.00,  "peso_kg": 0.3800},
    "Powkiddy V10":               {"fob_usd": 28.00,  "peso_kg": 0.3000},
    "Powkiddy V20 16GB":          {"fob_usd": 33.00,  "peso_kg": 0.3500},
    "Powkiddy V90S 16GB":         {"fob_usd": 33.00,  "peso_kg": 0.3000},
    "Powkiddy X35H 16GB":         {"fob_usd": 40.00,  "peso_kg": 0.3500},
    "Powkiddy X35s":              {"fob_usd": 45.00,  "peso_kg": 0.3500},
    "R36S Dual":                  {"fob_usd": 21.50,  "peso_kg": 0.3200},
    "Trimui Brick":               {"fob_usd": 51.00,  "peso_kg": 0.3350},
    "Trimui Brick Hammer":        {"fob_usd": 60.00,  "peso_kg": 0.4000},
    "Trimui Smart":               {"fob_usd": 31.50,  "peso_kg": 0.1400},
    "Trimui Smart Pro":           {"fob_usd": 54.00,  "peso_kg": 0.4050},
}

def get_fob_usd(nombre_prod, costos_gs=None):
    nombre_norm = _normalizar(nombre_prod)
    if not nombre_norm:
        return 0.0
    nombre_compact = _norm_compact(nombre_prod)
    candidatos = []
    if costos_gs:
        for k, v in costos_gs.items():
            if k.startswith("_"):
                continue
            if isinstance(v, dict):
                fob = float(v.get("fob_usd", 0) or 0)
                if fob > 0:
                    candidatos.append((_normalizar(k), _norm_compact(k), fob))
    for k, v in FOB_DEFAULTS.items():
        candidatos.append((_normalizar(k), _norm_compact(k), float(v.get("fob_usd", 0) or 0)))
    # Tier 1: match exacto
    for k_norm, k_compact, fob in candidatos:
        if k_norm == nombre_norm:
            return fob
    # Tier 2: match compacto exacto (sin colores/GB/RAM)
    for k_norm, k_compact, fob in candidatos:
        if k_compact and k_compact == nombre_compact:
            return fob
    # Tier 3: key en nombre
    for k_norm, k_compact, fob in candidatos:
        if k_norm in nombre_norm:
            return fob
    # Tier 4: key compacto en nombre compacto (el más largo gana)
    best, best_len = 0.0, 0
    for k_norm, k_compact, fob in candidatos:
        if k_compact and k_compact in nombre_compact and len(k_compact) > best_len:
            best, best_len = fob, len(k_compact)
    if best > 0:
        return best
    # Tier 5: nombre compacto en key compacto (reverso — ej: "rg35xxsp" in "rg35xxsp64")
    best, best_len = 0.0, 0
    for k_norm, k_compact, fob in candidatos:
        if k_compact and nombre_compact in k_compact and len(nombre_compact) > best_len:
            best, best_len = fob, len(nombre_compact)
    if best > 0:
        return best
    # Tier 6: nombre en key (normalización básica)
    for k_norm, k_compact, fob in candidatos:
        if nombre_norm in k_norm:
            return fob
    return 0.0

def _match_costo_entry(nombre_prod, costos_gs=None):
    """Encuentra la MEJOR entrada de costos y devuelve (fob_usd, import_usd, total_usd).
    Garantiza que FOB, import y total vienen de la MISMA entrada."""
    nombre_norm = _normalizar(nombre_prod)
    if not nombre_norm:
        return (0.0, 0.0, 0.0)
    ckg_default = float(costos_gs.get("_costo_kg_usd", 65.0) or 65.0) if costos_gs else 65.0

    def _extract(v, ckg):
        if not isinstance(v, dict):
            return (0.0, 0.0, 0.0)
        fob = float(v.get("fob_usd", 0) or 0)
        peso = float(v.get("peso_kg", 0) or 0)
        imp = float(v.get("costo_import_usd", 0) or 0)
        ct = float(v.get("costo_total_usd", 0) or 0)
        if imp <= 0:
            imp = round(peso * ckg, 2)
        if ct <= 0:
            ct = fob + imp
        return (fob, imp, ct)

    # Construir candidatos: (norm_basico, norm_compacto, data_dict, costo_kg)
    # Ordenar: entradas con FOB > 0 primero para que matcheen antes que las vacías
    candidatos = []
    if costos_gs:
        for k, v in costos_gs.items():
            if k.startswith("_") or not isinstance(v, dict):
                continue
            candidatos.append((_normalizar(k), _norm_compact(k), v, ckg_default))
    for k, v in FOB_DEFAULTS.items():
        candidatos.append((_normalizar(k), _norm_compact(k), v, 65.0))
    candidatos.sort(key=lambda c: -(float(c[2].get("fob_usd", 0) or 0) if isinstance(c[2], dict) else 0))

    nombre_compact = _norm_compact(nombre_prod)

    def _try_tiers():
        # Recorrer todos los tiers. Si un tier matchea con FOB>0, retornar.
        # Si matchea con FOB=0, guardar como fallback y seguir buscando.
        _fallback = None

        def _consider(r):
            nonlocal _fallback
            if r and r[2] > 0:
                if r[0] > 0:       # tiene FOB → retorno inmediato
                    return r
                if not _fallback:   # FOB=0, guardar como fallback
                    _fallback = r
            return None

        # Tier 1: match exacto
        for k_norm, k_compact, v, ckg in candidatos:
            if k_norm == nombre_norm:
                hit = _consider(_extract(v, ckg))
                if hit: return hit
        # Tier 2: match compacto exacto
        for k_norm, k_compact, v, ckg in candidatos:
            if k_compact and k_compact == nombre_compact:
                hit = _consider(_extract(v, ckg))
                if hit: return hit
        # Tier 3: key en nombre (básico, el más largo gana)
        best, best_len = None, 0
        for k_norm, k_compact, v, ckg in candidatos:
            if k_norm in nombre_norm and len(k_norm) > best_len:
                r = _extract(v, ckg)
                if r[2] > 0:
                    best, best_len = r, len(k_norm)
        if best:
            hit = _consider(best)
            if hit: return hit
        # Tier 4: key compacto en nombre compacto (el más largo gana)
        best, best_len = None, 0
        for k_norm, k_compact, v, ckg in candidatos:
            if k_compact and k_compact in nombre_compact and len(k_compact) > best_len:
                r = _extract(v, ckg)
                if r[2] > 0:
                    best, best_len = r, len(k_compact)
        if best:
            hit = _consider(best)
            if hit: return hit
        # Tier 5: nombre compacto en key compacto (reverso, el más largo gana)
        best, best_len = None, 0
        for k_norm, k_compact, v, ckg in candidatos:
            if k_compact and nombre_compact in k_compact and len(nombre_compact) > best_len:
                r = _extract(v, ckg)
                if r[2] > 0:
                    best, best_len = r, len(k_compact)
        if best:
            hit = _consider(best)
            if hit: return hit
        # Tier 6: nombre en key (básico)
        for k_norm, k_compact, v, ckg in candidatos:
            if nombre_norm in k_norm:
                hit = _consider(_extract(v, ckg))
                if hit: return hit

        return _fallback  # FOB=0 fallback (mejor que nada)

    return _try_tiers() or (0.0, 0.0, 0.0)

def get_costo_total_usd(nombre_prod, costos_gs=None):
    """Retorna costo total USD (FOB + import)."""
    return _match_costo_entry(nombre_prod, costos_gs)[2]

def calcular_costo_orden_ars(productos_str, cantidad, tipo_cambio_ars, costos_gs=None):
    prods = [p.strip() for p in str(productos_str).split(" / ") if p.strip()]
    if not prods:
        return 0.0
    if len(prods) == 1:
        return get_fob_usd(prods[0], costos_gs) * int(cantidad or 1) * tipo_cambio_ars
    return sum(get_fob_usd(p, costos_gs) * tipo_cambio_ars for p in prods)

def calcular_costo_total_orden_ars(productos_str, cantidad, tipo_cambio_ars, costos_gs=None):
    """Calcula costo total (FOB + import) en ARS."""
    prods = [p.strip() for p in str(productos_str).split(" / ") if p.strip()]
    if not prods:
        return 0.0
    if len(prods) == 1:
        return get_costo_total_usd(prods[0], costos_gs) * int(cantidad or 1) * tipo_cambio_ars
    return sum(get_costo_total_usd(p, costos_gs) * tipo_cambio_ars for p in prods)

def costo_final_row(row, tipo_cambio, costos_gs):
    """Costo de productos en ARS = (FOB + Import) de cada producto × cantidad × TC.
    Usa la tabla CostosConsolas como source of truth (FOB + Import).
    Solo cae al costo TN si la tabla no tiene el producto.
    """
    # 1) Intentar con CostosConsolas (FOB + Import)
    costo_calc = calcular_costo_total_orden_ars(
        row.get("Productos", ""), row.get("Cantidad", 1), tipo_cambio, costos_gs
    )
    if costo_calc > 0:
        return round(costo_calc, 0)
    # 2) Fallback: lo que TN reporta (puede ser FOB solo o desactualizado)
    costo_tn = float(row.get("Costo Productos ($)", 0) or 0)
    return round(costo_tn, 0)

def costo_total_final_row(row, tipo_cambio, costos_gs):
    """Costo total (FOB + import) para la fila."""
    return round(calcular_costo_total_orden_ars(
        row.get("Productos", ""), row.get("Cantidad", 1), tipo_cambio, costos_gs
    ), 0)

# ── Dólar blue ─────────────────────────────────────────────────────────────────
@st.cache_data(ttl=1800)
@st.cache_data(ttl=900, show_spinner=False)
def get_meta_spend(fecha_desde_str, fecha_hasta_str):
    """Trae gasto en Meta Ads para el período. Devuelve float ARS o None."""
    meta_token = st.secrets.get("META_TOKEN", "")
    meta_account = st.secrets.get("META_AD_ACCOUNT_ID", "")
    if not meta_token or not meta_account:
        return None
    # Limpiar el prefijo act_ si ya viene incluido
    account_id = meta_account.replace("act_", "")
    url = f"https://graph.facebook.com/v21.0/act_{account_id}/insights"
    params = {
        "access_token": meta_token,
        "fields": "spend,account_currency",
        "time_range": json.dumps({"since": fecha_desde_str, "until": fecha_hasta_str}),
        "level": "account",
    }
    try:
        r = requests.get(url, params=params, timeout=10)
        if r.status_code == 200:
            data = r.json().get("data", [])
            if data:
                spend = float(data[0].get("spend", 0))
                currency = data[0].get("account_currency", "USD")
                return spend, currency
        return None
    except Exception:
        return None

def get_dolar_blue():
    try:
        r = requests.get("https://api.bluelytics.com.ar/v2/latest", timeout=5)
        if r.status_code == 200:
            return float(r.json()["blue"]["value_sell"])
    except Exception:
        pass
    try:
        r = requests.get("https://dolarapi.com/v1/dolares/blue", timeout=5)
        if r.status_code == 200:
            return float(r.json().get("venta", 0))
    except Exception:
        pass
    return None

# ── Procesamiento de datos ─────────────────────────────────────────────────────
def _extraer_nombre_producto(n):
    if isinstance(n, str):
        return n
    if isinstance(n, dict):
        return n.get("es", "") or next(iter(n.values()), "")
    return ""

def procesar_orders(orders):
    filas = []
    for o in orders:
        prods = []
        costo_productos = 0.0
        items_linea = []
        for p in o.get("products", []):
            nombre = _extraer_nombre_producto(p.get("name", ""))
            prods.append(nombre)
            qty = int(p.get("quantity", 1) or 1)
            cost = float(p.get("cost", 0) or 0)
            costo_productos += cost * qty
            items_linea.append({"producto": nombre, "cantidad": qty, "costo": cost})
        productos = " / ".join(prods)
        cantidad = sum(int(p.get("quantity", 1) or 1) for p in o.get("products", []))

        pd_raw = o.get("payment_details", {})
        gateway = str(o.get("gateway", "")).lower()
        metodo = gateway
        cuotas = 1
        if isinstance(pd_raw, dict):
            metodo = pd_raw.get("method", gateway)
            cuotas = int(pd_raw.get("installments", 1) or 1)

        if metodo == "credit_card":
            label_medio = "Credito contado" if cuotas == 1 else f"Credito {cuotas} cuotas"
        elif metodo == "debit_card":
            label_medio = "Debito"
        elif any(x in str(metodo).lower() for x in ["transfer", "wire"]):
            label_medio = "Transferencia"
        elif "account_money" in str(metodo).lower():
            label_medio = "Dinero en cuenta"
        else:
            label_medio = str(metodo).replace("_", " ").title() if metodo else str(gateway)

        try:
            fecha = pd.to_datetime(o.get("created_at", "")).strftime("%Y-%m-%d")
        except Exception:
            fecha = ""

        total = float(o.get("total", 0))
        descuento = float(o.get("discount", 0) or 0)
        costo_envio_dueno = float(o.get("shipping_cost_owner", 0) or 0)
        province = str(o.get("billing_province", "")).strip()

        if _es_gateway_mp(gateway):
            pasarela = "MP"
            tasa = tasa_pasarela(gateway, metodo, cuotas)
            comision_pn = round(total * tasa, 2)
        elif _es_convenir(gateway, metodo):
            pasarela = "Convenir"   # se resolverá en match_mp_with_tn()
            tasa = 0.0
            comision_pn = 0.0
        else:
            pasarela = "PN"
            # Fee = tasa pública oficial de PN (no es estimación inventada,
            # son los rates publicados: 1.25% transferencia, 4.15% crédito, etc.)
            # La retención IIBB NO se calcula porque TN no la expone vía API.
            tasa = tasa_pasarela(gateway, metodo, cuotas)
            comision_pn = round(total * tasa, 2)
        neto = round(total - comision_pn, 2)
        margen = round(neto - costo_productos - costo_envio_dueno, 2)
        margen_pct = round((margen / total * 100) if total > 0 else 0, 2)

        filas.append({
            "Orden": o.get("number"),
            "Fecha": fecha,
            "Cliente": str(o.get("contact_name", "")),
            "Medio de Pago": label_medio,
            "Cuotas": cuotas,
            "Pasarela": pasarela,
            "Total ($)": total,
            "Descuento ($)": descuento,
            "Envio costo ($)": costo_envio_dueno,
            "Comision PN ($)": comision_pn,
            "Costo PN (%)": round(tasa * 100, 2),
            "Neto cobrado ($)": neto,
            "Costo Productos ($)": round(costo_productos, 2),
            "Margen ($)": margen,
            "Margen (%)": margen_pct,
            "Estado Envio": o.get("shipping_status", ""),
            "Productos": productos,
            "Cantidad": cantidad,
            "Canal": str(o.get("app_id", "") or "tiendanube"),
            "Estado": o.get("status", ""),
            "ID MP": "",
            "Provincia": province,
            "Gateway raw": gateway,
            "Metodo raw": str(metodo),
            "Items": items_linea,
        })
    return pd.DataFrame(filas)

# Tasas IIBB por provincia para órdenes Pago Nube (transferencia bancaria).
# Valores observados en MG. Editables vía Google Sheets si querés ajustarlas.
IIBB_DEFAULT_RATES = {
    "Capital Federal": 2.5,
    "CABA":            2.5,
    "Buenos Aires":    1.7,
    "Córdoba":         2.5,
    "Cordoba":         2.5,
    "Santa Fe":        4.5,
    "Mendoza":         3.0,
    "Salta":           1.5,
    "Tucumán":         2.5,
    "Tucuman":         2.5,
    "Entre Ríos":      2.0,
    "Entre Rios":      2.0,
    "Misiones":        3.0,
    "Neuquén":         2.5,
    "Neuquen":         2.5,
    "Río Negro":       2.0,
    "Rio Negro":       2.0,
    "Chaco":           3.5,
    "Corrientes":      2.5,
    "San Juan":        2.5,
    "San Luis":        2.5,
    "La Pampa":        2.5,
    "La Rioja":        2.5,
    "Catamarca":       2.5,
    "Santiago del Estero": 2.5,
    "Formosa":         2.5,
    "Jujuy":           2.5,
    "Chubut":          2.5,
    "Santa Cruz":      0.0,
    "Tierra del Fuego": 0.0,
}

def _iibb_rate_provincia(provincia):
    """Devuelve % de retención IIBB según provincia. 0 si no la conoce."""
    if not provincia:
        return 0.0
    # Permitir override desde Google Sheets / session state
    custom = st.session_state.get("iibb_rates") or {}
    p = str(provincia).strip()
    if p in custom:
        return float(custom[p])
    # Fallback al default
    return float(IIBB_DEFAULT_RATES.get(p, 0.0))

def _extraer_retencion_pn(p):
    """Busca retenciones de impuestos (IIBB, IVA, etc.) en el pago PN.
    TN expone esto en distintos lugares según la versión del payload.
    """
    retencion = 0.0
    # 1) Campos directos
    for campo in ["tax_amount", "withholding_amount", "retention_amount", "retencion"]:
        v = p.get(campo, 0)
        try:
            retencion += float(v or 0)
        except Exception:
            pass
    # 2) Array 'taxes' o 'withholdings'
    for arr_name in ["taxes", "withholdings", "tax_details"]:
        for t in p.get(arr_name, []) or []:
            try:
                retencion += float(t.get("amount", 0) or t.get("value", 0) or 0)
            except Exception:
                pass
    # 3) Charges con tipo de retención
    for c in p.get("charges", []) or []:
        tipo = str(c.get("type", "")).lower()
        if any(k in tipo for k in ["iibb", "ingresos", "withhold", "retencion", "retention", "tax"]):
            try:
                retencion += float(c.get("amount", 0) or 0)
            except Exception:
                pass
    # 4) Fallback: amount - fee_amount - net_amount (si los 3 están)
    if retencion == 0:
        try:
            amt = float(p.get("amount", 0) or 0)
            fee = float(p.get("fee_amount", 0) or p.get("fee", 0) or 0)
            net = float(p.get("net_amount", 0) or 0)
            diff = round(amt - fee - net, 2)
            if diff > 0:
                retencion = diff
        except Exception:
            pass
    return round(retencion, 2)

def procesar_pagos_pn(pagos):
    filas = []
    for p in pagos:
        try:
            fecha = pd.to_datetime(p.get("created_at", "")).strftime("%Y-%m-%d")
        except Exception:
            fecha = ""
        monto = float(p.get("amount", 0) or 0)
        fee = float(p.get("fee_amount", 0) or p.get("fee", 0) or 0)
        neto_api = float(p.get("net_amount", 0) or 0)
        retencion = _extraer_retencion_pn(p)
        # Costo total = fee + retención
        costo_total = round(fee + retencion, 2)
        # Neto real = monto - costo total (si el API trae net_amount usamos eso)
        neto_real = neto_api if neto_api > 0 else round(monto - costo_total, 2)
        filas.append({
            "ID": str(p.get("id", "")),
            "Fecha": fecha,
            "Estado": p.get("status", ""),
            "Método": p.get("payment_method", ""),
            "Cuotas": p.get("installments", 1),
            "Monto ($)": monto,
            "Fee ($)":          fee,
            "Retención ($)":    retencion,
            "Costo total ($)":  costo_total,
            "Neto ($)": neto_real,
            "Orden TN": str(p.get("order_id", "")),
        })
    return pd.DataFrame(filas) if filas else pd.DataFrame()

def calcular_tasas_reales(df_pagos):
    if df_pagos.empty or "Fee ($)" not in df_pagos.columns:
        return None
    df_con_fee = df_pagos[df_pagos["Fee ($)"] > 0]
    if df_con_fee.empty:
        return None
    tasas = {}
    for metodo, grupo in df_con_fee.groupby("Método"):
        monto_total = grupo["Monto ($)"].sum()
        fee_total = grupo["Fee ($)"].sum()
        if monto_total > 0:
            tasas[metodo] = round(fee_total / monto_total * 100, 2)
    return tasas if tasas else None

# ── Session state init ─────────────────────────────────────────────────────────
defaults = {
    "df_tn": None, "df_pagos": None, "orders_raw": [],
    "costos_productos": {}, "ordenes_efectivo": set(), "ids_venta_local": set(),
    "mp_raw": [], "mp_match_stats": {"matched": 0, "sin_match": 0},
}
for key, default in defaults.items():
    if key not in st.session_state:
        st.session_state[key] = default

# Session state for pauta (so it persists across tab switches)
if "pauta_manual" not in st.session_state:
    st.session_state.pauta_manual = 0
if "pct_iva" not in st.session_state:
    st.session_state.pct_iva = 10.5
if "tipo_cambio_sf" not in st.session_state:
    st.session_state.tipo_cambio_sf = None

# ── Sidebar ────────────────────────────────────────────────────────────────────
SECCIONES = [
    "📊 Dashboard",
    "🔍 Detalle y ajustes",
    "💚 Salud Financiera",
    "📦 Stock",
    "🔥 Velocidad de ventas",
    "🌐 Web / Analytics",
    "🏗️ Gastos fijos",
    "💻 Costos de consolas",
    "📐 Margen teórico",
    "📈 Margen real",
    "🏭 Proveedores",
    "💳 Estadísticas de pago",
    "🤖 Analista IA",
]

with st.sidebar:
    st.markdown("#### Market Gamer")
    seccion = st.radio("Navegación", SECCIONES, index=0, label_visibility="collapsed")
    st.divider()
    periodo = st.radio("Período", ["Este mes", "Mes anterior", "Esta semana", "Personalizado"], index=0)
    hoy = date.today()
    if periodo == "Este mes":
        fecha_desde = hoy.replace(day=1)
        fecha_hasta = hoy
    elif periodo == "Mes anterior":
        primer_dia = hoy.replace(day=1)
        fecha_hasta = primer_dia - timedelta(days=1)
        fecha_desde = fecha_hasta.replace(day=1)
    elif periodo == "Esta semana":
        fecha_desde = hoy - timedelta(days=hoy.weekday())
        fecha_hasta = hoy
    else:
        fecha_desde = st.date_input("Desde", value=hoy.replace(day=1))
        fecha_hasta = st.date_input("Hasta", value=hoy)

    if periodo != "Personalizado":
        st.caption(f"{fecha_desde.strftime('%d/%m/%Y')} → {fecha_hasta.strftime('%d/%m/%Y')}")
    buscar = st.button("Actualizar datos", use_container_width=True)

    # ── Config Salud Financiera (solo visible en esa sección) ──────────────────
    if seccion == "💚 Salud Financiera":
        st.divider()
        st.markdown(
            f'<p style="font-size:0.62rem;color:{MG_MUTED};font-family:\'Space Mono\','
            f'monospace;letter-spacing:0.08em;text-transform:uppercase;margin-bottom:0.5rem;">'
            f'⚙️ Configuración financiera</p>',
            unsafe_allow_html=True,
        )
        _dolar_raw_sb = get_dolar_blue()
        _dolar_default_sf = int(_dolar_raw_sb) if _dolar_raw_sb else 1200
        _tc_sf = st.number_input(
            f"💵 Dólar blue",
            value=st.session_state.tipo_cambio_sf or _dolar_default_sf,
            step=10,
            help="Tipo de cambio para convertir costos USD → ARS",
        )
        _iva_sf = st.slider("🧾 IVA efectivo (%)", 0.0, 21.0, float(st.session_state.pct_iva), 0.5)

        # Meta Ads — auto-fetch
        _meta_result_sb = get_meta_spend(str(fecha_desde), str(fecha_hasta))
        if _meta_result_sb:
            _meta_spend_sb, _meta_cur_sb = _meta_result_sb
            _pauta_auto = round(_meta_spend_sb) if _meta_cur_sb == "ARS" else round(_meta_spend_sb * (_tc_sf or _dolar_default_sf))
            st.success(
                f"📡 Meta: {_meta_cur_sb} {_meta_spend_sb:,.0f}"
                + (f" → ${_pauta_auto:,.0f}" if _meta_cur_sb != "ARS" else ""),
            )
            _pauta_default = int(st.session_state.pauta_manual) if st.session_state.pauta_manual else _pauta_auto
        else:
            _meta_token_check = st.secrets.get("META_TOKEN", "")
            if not _meta_token_check:
                st.caption("💡 Agregá `META_TOKEN` + `META_AD_ACCOUNT_ID` en secrets para auto-fetch.")
            else:
                st.warning("⚠️ Meta Ads no disponible — ingresá manualmente.")
            _pauta_default = st.session_state.pauta_manual

        _pauta_sf = st.number_input(
            "📣 Pauta (ARS)",
            value=int(_pauta_default),
            step=10_000,
            help="Gasto en publicidad Meta Ads del período",
        )
        if st.button("✅ Aplicar", use_container_width=True, key="btn_aplicar_sf"):
            st.session_state.tipo_cambio_sf = _tc_sf
            st.session_state.pct_iva = _iva_sf
            st.session_state.pauta_manual = _pauta_sf
            st.rerun()

# ── Helper: cargar y cruzar datos ─────────────────────────────────────────────
def _cargar_datos(fecha_desde, fecha_hasta, mostrar_success=False):
    """Carga órdenes TN + pagos PN + pagos MP y ejecuta el matching automático."""
    # 1. Órdenes TN
    orders = get_tn_orders(fecha_desde, fecha_hasta)
    if orders:
        orders_filtrados = []
        for o in orders:
            try:
                dt = pd.to_datetime(o.get("created_at", ""))
                if dt.tzinfo:
                    dt = dt.tz_convert(None)
                if fecha_desde <= dt.date() <= fecha_hasta:
                    orders_filtrados.append(o)
            except Exception:
                orders_filtrados.append(o)
        orders = orders_filtrados
        df_tn = procesar_orders(orders)
        st.session_state.orders_raw = orders
    else:
        df_tn = pd.DataFrame()
        st.session_state.orders_raw = []

    # 2. Pagos Pago Nube
    pagos = get_tn_pagos(fecha_desde, fecha_hasta)
    df_pagos_pn = procesar_pagos_pn(pagos) if pagos else pd.DataFrame()
    st.session_state.df_pagos = df_pagos_pn

    # 2b. Cross-reference (opcional): si /transactions devolvió datos reales,
    # se usan en lugar del estimado por tasa+provincia. Si vino vacío, se
    # mantiene la estimación que ya hizo procesar_orders.
    if not df_tn.empty and "Retención IIBB ($)" not in df_tn.columns:
        df_tn["Retención IIBB ($)"] = 0.0

    pn_match_stats = {"intentos": 0, "matched": 0, "sin_pago": 0, "estimados": 0}
    # Contar cuántas órdenes PN están usando estimación (sin match en /transactions)
    if not df_tn.empty:
        pn_match_stats["estimados"] = int((df_tn.get("Pasarela") == "PN").sum()) if "Pasarela" in df_tn.columns else 0

    if not df_tn.empty and not df_pagos_pn.empty:
        # Mapear order_id → (fee_real, retencion_real, costo_total_real)
        pagos_por_orden = {}
        for _, p in df_pagos_pn.iterrows():
            oid = str(p.get("Orden TN", "")).strip()
            if not oid:
                continue
            pagos_por_orden[oid] = {
                "fee": float(p.get("Fee ($)", 0) or 0),
                "retencion": float(p.get("Retención ($)", 0) or 0),
                "costo_total": float(p.get("Costo total ($)", 0) or 0),
                "neto": float(p.get("Neto ($)", 0) or 0),
            }

        orders_raw = st.session_state.get("orders_raw", [])
        # Construir mapa numero_orden → tn_order_id para matchear seguro
        num_to_tnid = {}
        for o in orders_raw:
            num = str(o.get("number", ""))
            tn_id = str(o.get("id", ""))
            if num and tn_id:
                num_to_tnid[num] = tn_id

        for idx, row in df_tn.iterrows():
            if row.get("Pasarela", "PN") != "PN":
                continue
            pn_match_stats["intentos"] += 1
            num_orden = str(row.get("Orden", ""))
            tn_order_id = num_to_tnid.get(num_orden, "")
            if not tn_order_id:
                continue
            if tn_order_id in pagos_por_orden:
                p_info = pagos_por_orden[tn_order_id]
                if p_info["costo_total"] > 0:
                    total = float(row.get("Total ($)", 0))
                    df_tn.at[idx, "Comision PN ($)"]    = round(p_info["costo_total"], 2)
                    df_tn.at[idx, "Retención IIBB ($)"] = round(p_info["retencion"], 2)
                    df_tn.at[idx, "Costo PN (%)"]       = round(
                        (p_info["costo_total"] / total * 100) if total > 0 else 0, 2
                    )
                    df_tn.at[idx, "Neto cobrado ($)"]   = round(total - p_info["costo_total"], 2)
                    pn_match_stats["matched"] += 1
            else:
                pn_match_stats["sin_pago"] += 1

    st.session_state.pn_match_stats = pn_match_stats

    # 2c. Cargar set de órdenes marcadas como efectivo (de Google Sheets) y aplicarlas
    # ANTES del matching MP para que no se peguen a un pago coincidente por accidente.
    ordenes_efectivo_raw = gs_read("OrdenesEfectivo") or {}
    ordenes_efectivo_set = set()
    if isinstance(ordenes_efectivo_raw, dict):
        for k, v in ordenes_efectivo_raw.items():
            if v:  # truthy = marcada
                ordenes_efectivo_set.add(str(k))
    st.session_state.ordenes_efectivo = ordenes_efectivo_set

    if not df_tn.empty and ordenes_efectivo_set:
        for idx, row in df_tn.iterrows():
            num_orden = str(row.get("Orden", ""))
            if num_orden in ordenes_efectivo_set:
                total = float(row.get("Total ($)", 0))
                df_tn.at[idx, "Pasarela"]          = "Efectivo"
                df_tn.at[idx, "Comision PN ($)"]   = 0.0
                df_tn.at[idx, "Costo PN (%)"]      = 0.0
                df_tn.at[idx, "Neto cobrado ($)"]  = total
                df_tn.at[idx, "ID MP"]             = ""
                # Recalcular margen sin descontar comisión
                costo_prods = float(row.get("Costo Productos ($)", 0))
                costo_envio = float(row.get("Envio costo ($)", 0))
                df_tn.at[idx, "Margen ($)"] = round(total - costo_prods - costo_envio, 2)
                df_tn.at[idx, "Margen (%)"] = round(
                    ((total - costo_prods - costo_envio) / total * 100) if total > 0 else 0, 2
                )

    # 3. Pagos Mercado Pago + matching automático con órdenes "a convenir"
    # match_mp_with_tn() ya saltea órdenes con Pasarela == "Efectivo"
    if MP_ACCESS_TOKEN and not df_tn.empty:
        mp_raw = get_mp_payments(str(fecha_desde), str(fecha_hasta))
        st.session_state.mp_raw = mp_raw
        if mp_raw:
            df_tn, n_matched, n_sin = match_mp_with_tn(df_tn, mp_raw)
            st.session_state.mp_match_stats = {"matched": n_matched, "sin_match": n_sin}
        else:
            st.session_state.mp_match_stats = {"matched": 0, "sin_match": 0}
    else:
        st.session_state.mp_raw = []
        st.session_state.mp_match_stats = {"matched": 0, "sin_match": 0}

    st.session_state.df_tn = df_tn
    st.session_state.ids_venta_local = set()

    if mostrar_success and orders:
        n_matched = st.session_state.mp_match_stats["matched"]
        n_convenir= int((df_tn["Pasarela"] == "Convenir").sum()) if not df_tn.empty and "Pasarela" in df_tn.columns else 0
        msg = f"✅ {len(orders)} órdenes cargadas"
        if n_matched:
            msg += f" · 🔀 {n_matched} órdenes 'a convenir' cruzadas con MP"
        if n_convenir:
            msg += f" · ⚠️ {n_convenir} sin cruzar (sin pago MP coincidente)"
        st.success(msg)
    elif mostrar_success:
        st.info("No se encontraron órdenes en el período.")

@st.cache_data(ttl=1800, show_spinner="Cargando histórico de ventas...")
def _cargar_ordenes_historico(dias_historia):
    """Órdenes sobre una ventana amplia, independiente del período del sidebar.

    Cacheado por dias_historia (TTL 30 min) para no refetchear en cada rerun.
    Liviano: solo fetch + procesar_orders (sin matching PN/MP, innecesario para velocidad).
    """
    desde = (date.today() - timedelta(days=dias_historia)).isoformat()
    hasta = date.today().isoformat()
    try:
        orders = get_tn_orders(desde, hasta)
    except Exception:
        return pd.DataFrame()
    return procesar_orders(orders) if orders else pd.DataFrame()

# ── Búsqueda ───────────────────────────────────────────────────────────────────
if buscar:
    _cargar_datos(fecha_desde, fecha_hasta, mostrar_success=True)
    _sd = st.session_state.get("stock_tn")
    if _sd is not None:
        _sm = {}
        for _, _r in _sd.iterrows():
            _sv = _r["Stock"]
            if isinstance(_sv, (int, float)):
                _sm[_r["Producto"]] = _sm.get(_r["Producto"], 0) + int(_sv)
        if _sm:
            gs_append_snapshot(_sm)

# ── Auto-carga del mes actual en primera visita ──────────────────────────────
if st.session_state.df_tn is None:
    with st.spinner("Cargando datos..."):
        _cargar_datos(fecha_desde, fecha_hasta, mostrar_success=False)

# ── Contenido principal ───────────────────────────────────────────────────────
dolar_blue = get_dolar_blue()

if st.session_state.df_tn is not None:
    # ── Helper: Extraer precios individuales por producto desde órdenes raw ──
    def _build_product_rows_from_raw(orders_raw):
        """
        Extrae precio individual de cada producto desde la API de TN.
        Cada producto en la orden tiene su propio 'price', así que no
        dividimos el total de la orden entre productos.
        """
        product_rows = []
        for o in orders_raw:
            pd_raw = o.get("payment_details", {})
            gateway = str(o.get("gateway", "")).lower()
            metodo = gateway
            cuotas = 1
            if isinstance(pd_raw, dict):
                metodo = pd_raw.get("method", gateway)
                cuotas = int(pd_raw.get("installments", 1) or 1)

            if metodo == "credit_card":
                label_medio = "Credito contado" if cuotas == 1 else f"Credito {cuotas} cuotas"
            elif metodo == "debit_card":
                label_medio = "Debito"
            elif any(x in str(metodo).lower() for x in ["transfer", "wire"]):
                label_medio = "Transferencia"
            elif "account_money" in str(metodo).lower():
                label_medio = "Dinero en cuenta"
            else:
                label_medio = str(metodo).replace("_", " ").title() if metodo else str(gateway)

            try:
                fecha = pd.to_datetime(o.get("created_at", "")).strftime("%Y-%m-%d")
            except Exception:
                fecha = ""

            order_total = float(o.get("total", 0))
            n_products = len(o.get("products", []))
            costo_envio = float(o.get("shipping_cost_owner", 0) or 0)
            shipping_customer = float(o.get("shipping_cost_customer", 0) or 0)
            tasa = tasa_pasarela(gateway, metodo, cuotas)

            # Calcular subtotal real de productos (precio de lista * qty)
            order_subtotal = 0.0
            for _p in o.get("products", []):
                try:
                    order_subtotal += float(_p.get("price", 0) or 0) * int(_p.get("quantity", 1) or 1)
                except Exception:
                    pass
            # Revenue real de productos = total pagado por el cliente - lo que pagó por envío
            # Esto absorbe descuentos por medio de pago, cupones y promos
            product_revenue = order_total - shipping_customer
            ratio_neto = (product_revenue / order_subtotal) if order_subtotal > 0 else 1.0

            for p in o.get("products", []):
                nombre = _extraer_nombre_producto(p.get("name", ""))
                precio_unit = float(p.get("price", 0) or 0)
                qty = int(p.get("quantity", 1) or 1)

                # Si el precio individual es 0, fallback a dividir total
                if precio_unit <= 0 and n_products > 0:
                    precio_unit = order_total / n_products

                # Precio neto post-descuentos prorrateado por participación en subtotal
                precio_neto_unit = round(precio_unit * ratio_neto, 2)
                descuento_unit = round(precio_unit - precio_neto_unit, 2)

                # Comisión proporcional al precio del producto respecto al total
                if order_total > 0:
                    peso_en_orden = (precio_unit * qty) / order_total
                else:
                    peso_en_orden = 1.0 / max(n_products, 1)
                comision_unit = round(order_total * tasa * peso_en_orden / qty, 2)
                envio_unit = round(costo_envio * peso_en_orden / qty, 2)

                for _ in range(qty):
                    product_rows.append({
                        "Producto": nombre,
                        "Precio ($)": round(precio_unit, 0),
                        "Precio neto ($)": round(precio_neto_unit, 0),
                        "Descuento ($)": round(descuento_unit, 0),
                        "Medio de Pago": label_medio,
                        "Cuotas": cuotas,
                        "Comisión PN ($)": round(comision_unit, 0),
                        "Tasa PN (%)": round(tasa * 100, 2),
                        "Envío ($)": round(envio_unit, 0),
                        "Fecha": fecha,
                        "Orden Total ($)": order_total,
                    })
        return product_rows

    df_tn = st.session_state.df_tn.copy()
    df_pagos = st.session_state.df_pagos.copy() if st.session_state.df_pagos is not None else pd.DataFrame()

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 1: DASHBOARD
    # ══════════════════════════════════════════════════════════════════════════
    if seccion == "📊 Dashboard":
        if df_tn.empty:
            st.info("No hay órdenes en este período.")
        else:
            # ── Helpers ──
            def _truncar(nombre, max_len=32):
                return nombre if len(nombre) <= max_len else nombre[:max_len].rstrip() + "…"

            def _fmt_compact(n):
                n = float(n)
                if abs(n) >= 1_000_000:
                    return f"${n/1_000_000:.1f}M"
                elif abs(n) >= 1_000:
                    return f"${n/1_000:.0f}K"
                return f"${n:.0f}"

            _ACCESORIOS_KW = ("estuche", "funda", "micro sd", "microsd", "lectora", "cable", "joystick cap")

            def _es_accesorio(nombre):
                n = nombre.lower()
                return any(kw in n for kw in _ACCESORIOS_KW)

            # Helper de card uniforme (mismo patrón que Salud Financiera)
            def _kpi_dash(label, value, sub="", val_color=None, accent=False):
                vc = val_color or MG_TEXT
                border = f"border-left:2px solid {MG_RED};padding-left:0.75rem;" if accent else ""
                sub_html = (
                    f'<div style="font-size:0.68rem;color:{MG_MUTED};margin-top:0.25rem;'
                    f'white-space:nowrap;overflow:hidden;text-overflow:ellipsis;">{sub}</div>'
                    if sub else
                    f'<div style="font-size:0.68rem;color:transparent;margin-top:0.25rem;">—</div>'
                )
                return (
                    f'<div style="background:{MG_SURF};border-radius:8px;padding:0.9rem 1rem;'
                    f'min-height:90px;display:flex;flex-direction:column;justify-content:flex-start;{border}">'
                    f'<div style="font-size:0.58rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                    f'text-transform:uppercase;letter-spacing:0.08em;margin-bottom:0.4rem;'
                    f'white-space:nowrap;overflow:hidden;text-overflow:ellipsis;">{label}</div>'
                    f'<div style="font-size:1.35rem;font-weight:700;color:{vc};line-height:1.1;">{value}</div>'
                    f'{sub_html}'
                    f'</div>'
                )

            # ── Métricas base — separadas por pasarela ──
            mask_pn = df_tn["Pasarela"] == "PN" if "Pasarela" in df_tn.columns else pd.Series([False]*len(df_tn))
            mask_mp = df_tn["Pasarela"].isin(["MP", "Convenir"]) if "Pasarela" in df_tn.columns else pd.Series([False]*len(df_tn))

            bruto_pn = df_tn.loc[mask_pn, "Total ($)"].sum() if mask_pn.any() else 0
            com_pn   = df_tn.loc[mask_pn, "Comision PN ($)"].sum() if mask_pn.any() else 0
            pct_pn   = (com_pn / bruto_pn * 100) if bruto_pn > 0 else 0.0

            bruto_mp = df_tn.loc[mask_mp, "Total ($)"].sum() if mask_mp.any() else 0
            com_mp   = df_tn.loc[mask_mp, "Comision PN ($)"].sum() if mask_mp.any() else 0
            pct_mp   = (com_mp / bruto_mp * 100) if bruto_mp > 0 else 0.0

            total_facturado    = df_tn["Total ($)"].sum()
            total_comision     = com_pn + com_mp
            neto_cobrado       = total_facturado - total_comision
            costo_ponderado_pn = round(total_comision / total_facturado * 100, 2) if total_facturado > 0 else 0
            ticket_prom        = total_facturado / len(df_tn) if len(df_tn) > 0 else 0

            # ── Fila 1: Resumen ventas (4 cards alineadas) ─────────────────────
            r1c1, r1c2, r1c3, r1c4 = st.columns(4)
            r1c1.markdown(
                _kpi_dash("Órdenes", str(len(df_tn)), f"ticket promedio {fmt(ticket_prom)}"),
                unsafe_allow_html=True,
            )
            r1c2.markdown(
                _kpi_dash("Facturación bruta", fmt(total_facturado), "total vendido"),
                unsafe_allow_html=True,
            )
            r1c3.markdown(
                _kpi_dash("Comisiones totales", fmt(total_comision), f"PN + MP · {costo_ponderado_pn:.2f}%", val_color=MG_RED),
                unsafe_allow_html=True,
            )
            r1c4.markdown(
                _kpi_dash("Neto cobrado", fmt(neto_cobrado), "después de comisiones"),
                unsafe_allow_html=True,
            )

            # ── Fila 2: Costos por pasarela (2 cards) ──────────────────────────
            st.markdown("")
            r2c1, r2c2 = st.columns(2)
            r2c1.markdown(
                _kpi_dash(
                    "🔵 Costo Pago Nube",
                    fmt(com_pn),
                    f"sobre {fmt(bruto_pn)} · {pct_pn:.2f}% efectivo",
                    val_color=MG_RED,
                ),
                unsafe_allow_html=True,
            )
            r2c2.markdown(
                _kpi_dash(
                    "💳 Costo Mercado Pago",
                    fmt(com_mp),
                    f"sobre {fmt(bruto_mp)} · {pct_mp:.2f}% efectivo",
                    val_color=MG_RED,
                    accent=True,
                ),
                unsafe_allow_html=True,
            )

            st.markdown("")

            # ── Toggle accesorios (afecta los tops de productos) ───────────────
            _excluir_acc = st.toggle(
                "Excluir accesorios (estuches, micros SD, etc.)",
                value=False,
                key="dash_excluir_acc",
            )

            # ── Ventas por día (BARRAS) + Top 10 por unidades — alineadas ─────
            _CHART_HEIGHT = 360
            col_a, col_b = st.columns(2)
            with col_a:
                df_dia = df_tn.groupby("Fecha")["Total ($)"].sum().reset_index()
                df_dia["Fecha"] = pd.to_datetime(df_dia["Fecha"])
                df_dia = df_dia.sort_values("Fecha")
                fig_dia = px.bar(
                    df_dia, x="Fecha", y="Total ($)",
                    title="Facturación diaria",
                    color_discrete_sequence=[MG_RED],
                )
                fig_dia.update_layout(
                    yaxis_tickformat="$,.0f",
                    xaxis_tickformat="%d %b",
                    bargap=0.25,
                    height=_CHART_HEIGHT,
                    margin=dict(t=50, b=40, l=10, r=10),
                )
                fig_dia.update_traces(
                    hovertemplate="%{x|%d/%m}<br>$%{y:,.0f}<extra></extra>",
                )
                st.plotly_chart(fig_dia, use_container_width=True)

            with col_b:
                top_prods = {}
                for _, row in df_tn.iterrows():
                    for p in str(row.get("Productos", "")).split(" / "):
                        p = p.strip()
                        if p:
                            if _excluir_acc and _es_accesorio(p):
                                continue
                            top_prods[p] = top_prods.get(p, 0) + 1
                if top_prods:
                    df_tp = pd.DataFrame(list(top_prods.items()), columns=["Producto", "Unidades"])
                    df_tp = df_tp.sort_values("Unidades", ascending=False).head(10)
                    df_tp["Label"] = df_tp["Producto"].apply(_truncar)
                    fig_tp = px.bar(
                        df_tp, x="Unidades", y="Label", orientation="h",
                        title="Top 10 productos (unidades)",
                        color="Unidades",
                        color_continuous_scale=[[0, "#26272b"], [1, "#009EE3"]],
                        text="Unidades",
                        custom_data=["Producto"],
                    )
                    fig_tp.update_layout(
                        yaxis={"categoryorder": "total ascending", "title": ""},
                        coloraxis_showscale=False,
                        height=_CHART_HEIGHT,
                        margin=dict(t=50, b=40, l=10, r=10),
                    )
                    fig_tp.update_traces(
                        textposition="outside", textfont_size=11,
                        hovertemplate="<b>%{customdata[0]}</b><br>%{x} unidades<extra></extra>",
                    )
                    st.plotly_chart(fig_tp, use_container_width=True)

            # ── Top 10 facturación (full width) ───────────────────────────────
            top_revenue = {}
            for _, row in df_tn.iterrows():
                prods_list = [p.strip() for p in str(row.get("Productos", "")).split(" / ") if p.strip()]
                n_prods = max(len(prods_list), 1)
                for p in prods_list:
                    top_revenue[p] = top_revenue.get(p, 0) + row.get("Total ($)", 0) / n_prods
            if top_revenue:
                df_rev = pd.DataFrame(list(top_revenue.items()), columns=["Producto", "Monto ($)"])
                df_rev["Monto ($)"] = df_rev["Monto ($)"].round(0)
                df_rev = df_rev.sort_values("Monto ($)", ascending=False).head(10)
                df_rev["Label"] = df_rev["Producto"].apply(_truncar)
                df_rev["Texto"] = df_rev["Monto ($)"].apply(_fmt_compact)
                fig_rev = px.bar(
                    df_rev, x="Monto ($)", y="Label", orientation="h",
                    title="Top 10 productos (facturación)",
                    color="Monto ($)",
                    color_continuous_scale=[[0, "#26272b"], [1, MG_RED]],
                    text="Texto",
                    custom_data=["Producto"],
                )
                fig_rev.update_layout(
                    yaxis={"categoryorder": "total ascending", "title": ""},
                    xaxis_tickformat="$,.0f",
                    coloraxis_showscale=False,
                    height=380,
                    margin=dict(t=50, b=40, l=10, r=10),
                )
                fig_rev.update_traces(
                    textposition="inside", textfont_size=11, textfont_color=MG_TEXT,
                    hovertemplate="<b>%{customdata[0]}</b><br>$%{x:,.0f}<extra></extra>",
                )
                st.plotly_chart(fig_rev, use_container_width=True)

            # ── Donuts separados: Pago Nube vs Mercado Pago ───────────────────
            st.markdown(
                f'<p style="font-size:0.72rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                f'letter-spacing:0.06em;text-transform:uppercase;margin-top:1rem;margin-bottom:0.4rem;">'
                f'Distribución de comisiones por pasarela</p>',
                unsafe_allow_html=True,
            )

            def _donut_pasarela(df_subset, titulo, color_main):
                """Donut con top 3 medios de pago + 'Resto' para una pasarela."""
                if df_subset.empty:
                    return None
                cm = df_subset.groupby("Medio de Pago").agg(
                    Comision=("Comision PN ($)", "sum"),
                ).reset_index().sort_values("Comision", ascending=False)
                cm = cm[cm["Comision"] > 0]
                if cm.empty:
                    return None
                top3 = cm.head(3).copy()
                resto = cm.iloc[3:]
                if not resto.empty:
                    resto_row = pd.DataFrame([{
                        "Medio de Pago": "Resto",
                        "Comision": resto["Comision"].sum(),
                    }])
                    df_donut = pd.concat([top3, resto_row], ignore_index=True)
                else:
                    df_donut = top3.copy()

                fig = px.pie(
                    df_donut, names="Medio de Pago", values="Comision",
                    title=titulo,
                    color_discrete_sequence=[color_main, "#fbbf24", "#a78bfa", MG_DIM],
                    hole=0.55,
                )
                fig.update_traces(
                    textinfo="label+percent",
                    textfont_size=11,
                    hovertemplate="<b>%{label}</b><br>$%{value:,.0f} (%{percent})<extra></extra>",
                )
                fig.update_layout(
                    showlegend=False,
                    height=360,
                    margin=dict(t=50, b=20, l=10, r=10),
                )
                return fig

            col_dpn, col_dmp = st.columns(2)
            with col_dpn:
                fig_dpn = _donut_pasarela(
                    df_tn[mask_pn] if mask_pn.any() else pd.DataFrame(),
                    f"🔵 Pago Nube — total {fmt(com_pn)} ({pct_pn:.2f}%)",
                    "#009EE3",
                )
                if fig_dpn is not None:
                    st.plotly_chart(fig_dpn, use_container_width=True)
                else:
                    st.info("No hay órdenes vía Pago Nube en el período.")
            with col_dmp:
                fig_dmp = _donut_pasarela(
                    df_tn[mask_mp] if mask_mp.any() else pd.DataFrame(),
                    f"💳 Mercado Pago — total {fmt(com_mp)} ({pct_mp:.2f}%)",
                    MG_RED,
                )
                if fig_dmp is not None:
                    st.plotly_chart(fig_dmp, use_container_width=True)
                else:
                    st.info("No hay órdenes vía Mercado Pago en el período.")

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 2: DETALLE Y AJUSTES
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🔍 Detalle y ajustes":
        st.subheader("🛍️ Detalle de órdenes — Tienda Nube")
        st.caption("🔧 v0.8.0")

        if df_tn.empty:
            st.info("No hay órdenes en este período.")
        else:
            _dolar_det = dolar_blue or 1200
            _costos_gs = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
            df_det = df_tn.copy()
            df_det["Costo Productos ($)"] = df_det.apply(
                lambda r: costo_final_row(r, _dolar_det, _costos_gs), axis=1
            )
            df_det["Neto cobrado ($)"] = df_det["Total ($)"] - df_det["Comision PN ($)"]
            df_det["Margen ($)"] = (
                df_det["Neto cobrado ($)"] - df_det["Costo Productos ($)"]
                - df_det["Envio costo ($)"].fillna(0)
            ).round(2)
            df_det["Margen (%)"] = df_det.apply(
                lambda r: round((r["Margen ($)"] / r["Total ($)"] * 100) if r["Total ($)"] > 0 else 0, 2), axis=1
            )

            # ══════════════════════════════════════════════════════════════════
            # 1) KPIs ARRIBA — 5 cards uniformes
            # ══════════════════════════════════════════════════════════════════
            _total_bruto = df_det["Total ($)"].sum()
            _total_costo = df_det["Costo Productos ($)"].sum()
            _total_comis = df_det["Comision PN ($)"].sum()
            _total_margen = df_det["Margen ($)"].sum()
            _n_ordenes = len(df_det)
            _pct_margen_tot = round(_total_margen / _total_bruto * 100, 1) if _total_bruto > 0 else 0
            _color_margen_tot = "#4ade80" if _total_margen >= 0 else MG_RED

            def _kpi_det(label, value, sub="", color=None):
                vc = color or MG_TEXT
                return (
                    f'<div style="background:{MG_SURF};border-radius:8px;padding:0.85rem 1rem;'
                    f'min-height:88px;display:flex;flex-direction:column;justify-content:flex-start;">'
                    f'<div style="font-size:0.58rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                    f'text-transform:uppercase;letter-spacing:0.08em;margin-bottom:0.4rem;">{label}</div>'
                    f'<div style="font-size:1.3rem;font-weight:700;color:{vc};line-height:1.1;">{value}</div>'
                    f'<div style="font-size:0.66rem;color:{MG_MUTED};margin-top:0.2rem;">{sub or "—"}</div>'
                    f'</div>'
                )
            _k1, _k2, _k3, _k4, _k5 = st.columns(5)
            _k1.markdown(_kpi_det("Órdenes", str(_n_ordenes), "del período"), unsafe_allow_html=True)
            _k2.markdown(_kpi_det("Facturación bruta", fmt(_total_bruto), "total vendido"), unsafe_allow_html=True)
            _k3.markdown(_kpi_det("Costo productos", fmt(_total_costo), f"−{fmt(_total_costo)}", color=MG_RED), unsafe_allow_html=True)
            _k4.markdown(_kpi_det("Costo financiero", fmt(_total_comis), f"−{fmt(_total_comis)}", color=MG_RED), unsafe_allow_html=True)
            _k5.markdown(_kpi_det("Margen neto", fmt(_total_margen), f"{_pct_margen_tot}% sobre bruto", color=_color_margen_tot), unsafe_allow_html=True)

            # ══════════════════════════════════════════════════════════════════
            # 2) TABLA PRINCIPAL — protagonista
            # ══════════════════════════════════════════════════════════════════
            st.divider()
            cols_tn = [
                "Orden", "Fecha", "Cliente", "Medio de Pago", "Cuotas", "Pasarela", "ID MP",
                "Total ($)", "Descuento ($)", "Envio costo ($)",
                "Comision PN ($)", "Costo PN (%)",
                "Neto cobrado ($)", "Costo Productos ($)", "Margen ($)", "Margen (%)",
                "Estado Envio", "Productos", "Cantidad", "Canal",
            ]
            cols_tn = [c for c in cols_tn if c in df_det.columns]

            def _color_margen_row(row):
                styles = [""] * len(row)
                if "Margen ($)" in row.index:
                    idx = list(row.index).index("Margen ($)")
                    if float(row["Margen ($)"]) >= 0:
                        styles[idx] = "background-color: #1e4620; color: white"
                    else:
                        styles[idx] = "background-color: #4a1010; color: white"
                return styles

            sin_costo_tn = int((df_det["Costo Productos ($)"] == 0).sum())
            _caption_parts = [f"💵 TC: ${_dolar_det:,.0f} ARS/USD"]
            if sin_costo_tn > 0:
                _caption_parts.append(f"⚠️ {sin_costo_tn} sin costo")
            _caption_parts.append("ℹ️ PN: tasas oficiales públicas; retenciones IIBB no se exponen vía API")
            st.caption(" · ".join(_caption_parts))

            _rename_view = {
                "Comision PN ($)": "Costo fin. ($)",
                "Costo PN (%)":    "Costo fin. %",
            }
            df_view = df_det[cols_tn].rename(columns=_rename_view)
            st.dataframe(
                df_view.style
                    .format({
                        "Total ($)": "${:,.0f}", "Descuento ($)": "${:,.0f}",
                        "Envio costo ($)": "${:,.0f}", "Costo fin. ($)": "${:,.0f}",
                        "Costo fin. %": "{:.2f}%", "Neto cobrado ($)": "${:,.0f}",
                        "Costo Productos ($)": "${:,.0f}", "Margen ($)": "${:,.0f}",
                        "Margen (%)": "{:.2f}%",
                    })
                    .apply(_color_margen_row, axis=1),
                use_container_width=True, hide_index=True, height=460,
            )

            # ══════════════════════════════════════════════════════════════════
            # 3) ACCIONES — fila compacta de 3 columnas
            # ══════════════════════════════════════════════════════════════════
            _orders_raw_dbg = st.session_state.get("orders_raw", [])
            _ac1, _ac2, _ac3 = st.columns(3)

            with _ac1:
                st.download_button(
                    "⬇️ Descargar CSV",
                    df_view.to_csv(index=False).encode("utf-8"),
                    "ordenes_tn.csv", "text/csv",
                    use_container_width=True,
                )

            with _ac2.expander("💵 Marcar efectivo", expanded=False):
                st.caption("Las marcadas no tienen costo fin. ni se cruzan con MP. Persiste en Sheets.")
                _ordenes_actuales = st.session_state.get("ordenes_efectivo", set()) or set()
                _opciones_ord = []
                _label_a_num = {}
                for _, row in df_det.iterrows():
                    n = str(row.get("Orden", ""))
                    cli = str(row.get("Cliente", ""))[:25]
                    tot = float(row.get("Total ($)", 0) or 0)
                    label = f"{n} — {cli} — ${tot:,.0f}"
                    _opciones_ord.append(label)
                    _label_a_num[label] = n
                _default_labels = [lbl for lbl, n in _label_a_num.items() if n in _ordenes_actuales]
                _seleccion_efec = st.multiselect(
                    "Órdenes en efectivo", options=_opciones_ord, default=_default_labels,
                    label_visibility="collapsed",
                )
                if st.button("💾 Guardar", use_container_width=True, key="btn_save_efec"):
                    nuevos_nums = {_label_a_num[l] for l in _seleccion_efec}
                    actual_gs = gs_read("OrdenesEfectivo") or {}
                    if not isinstance(actual_gs, dict):
                        actual_gs = {}
                    nums_periodo = {str(o) for o in df_det["Orden"].tolist()}
                    actual_gs = {k: v for k, v in actual_gs.items() if str(k) not in nums_periodo}
                    for n in nuevos_nums:
                        actual_gs[str(n)] = True
                    if gs_write("OrdenesEfectivo", actual_gs):
                        st.success(f"✅ {len(nuevos_nums)} marcadas. Click 'Actualizar datos' para aplicar.")
                    else:
                        st.error("❌ Error guardando en Sheets.")

            with _ac3.expander(f"🔧 Debug orden ({len(_orders_raw_dbg)})", expanded=False):
                if not _orders_raw_dbg:
                    st.warning("No hay órdenes en sesión. Click 'Actualizar datos'.")
                else:
                    _opciones = ["— elegí una orden —"]
                    _orden_por_label = {}
                    for o in _orders_raw_dbg:
                        n = o.get("number", "?")
                        cli = str(o.get("contact_name", ""))[:25]
                        tot = float(o.get("total", 0) or 0)
                        label = f"{n} — {cli} — ${tot:,.0f}"
                        _opciones.append(label)
                        _orden_por_label[label] = o
                    _sel = st.selectbox(
                        "Elegí orden", _opciones, index=0, key="dbg_order_sel",
                        label_visibility="collapsed",
                    )
                    if _sel and _sel in _orden_por_label:
                        o = _orden_por_label[_sel]
                        st.markdown(f"**#{o.get('number')}** · ID `{o.get('id')}` · gw `{o.get('gateway')}`")
                        # Match de costos
                        st.markdown("**🔍 Match de costos:**")
                        _costos_dbg = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
                        _tc_dbg = float(st.session_state.tipo_cambio_sf or (dolar_blue or 1400))
                        _filas_costo_dbg = []
                        for _p in o.get("products", []) or []:
                            _nombre_raw = _p.get("name", {})
                            if isinstance(_nombre_raw, dict):
                                _nombre = _nombre_raw.get("es", "") or next(iter(_nombre_raw.values()), "")
                            else:
                                _nombre = str(_nombre_raw)
                            _nombre = _extraer_nombre_producto(_nombre)
                            _qty = int(_p.get("quantity", 1) or 1)
                            _costo_tn = float(_p.get("cost") or 0)
                            _matched_key = "(no match)"
                            _nombre_norm = _normalizar(_nombre)
                            _nombre_comp = _norm_compact(_nombre)
                            _all_sources = []
                            for k, v in (_costos_dbg or {}).items():
                                if not k.startswith("_") and isinstance(v, dict):
                                    _all_sources.append(("Sheets", k, v))
                            for k, v in FOB_DEFAULTS.items():
                                _all_sources.append(("Default", k, v))
                            for src, k, v in _all_sources:
                                if _normalizar(k) == _nombre_norm:
                                    _matched_key = f"[{src}] {k} (exact)"
                                    break
                            else:
                                for src, k, v in _all_sources:
                                    if _norm_compact(k) and _norm_compact(k) == _nombre_comp:
                                        _matched_key = f"[{src}] {k} (compact)"
                                        break
                                else:
                                    _best_len = 0
                                    for src, k, v in _all_sources:
                                        if _normalizar(k) in _nombre_norm and len(_normalizar(k)) > _best_len:
                                            _matched_key = f"[{src}] {k} (key in nombre)"
                                            _best_len = len(_normalizar(k))
                                    if _best_len == 0:
                                        for src, k, v in _all_sources:
                                            if _norm_compact(k) and _nombre_comp in _norm_compact(k) and len(_nombre_comp) > _best_len:
                                                _matched_key = f"[{src}] {k} (nombre in key)"
                                                _best_len = len(_nombre_comp)
                            _fob_real, _imp_real, _total_real = _match_costo_entry(_nombre, _costos_dbg)
                            _filas_costo_dbg.append({
                                "Producto": _nombre,
                                "Qty": _qty,
                                "Match": _matched_key,
                                "FOB USD": round(_fob_real, 2),
                                "Imp USD": round(_imp_real, 2),
                                "Tot USD": round(_total_real, 2),
                                "Tot ARS": round(_total_real * _qty * _tc_dbg, 0),
                                "TN cost": _costo_tn,
                            })
                        if _filas_costo_dbg:
                            st.dataframe(pd.DataFrame(_filas_costo_dbg), use_container_width=True, hide_index=True)
                        st.markdown("**payment_details:**")
                        _pd_data = o.get("payment_details", {}) or {}
                        if not _pd_data:
                            st.warning("⚠️ Sin payment_details (efectivo o convenir manual).")
                        else:
                            st.json(_pd_data)
                        with st.expander("Ver orden completa (JSON)", expanded=False):
                            st.json(o)

            # ══════════════════════════════════════════════════════════════════
            # 4) DESGLOSES — fila compacta de 2 columnas
            # ══════════════════════════════════════════════════════════════════
            _df_pn_det = st.session_state.get("df_pagos")
            _has_pn = _df_pn_det is not None and not _df_pn_det.empty
            _df_mp_det = pd.DataFrame()
            if MP_ACCESS_TOKEN:
                _mp_raw_det = get_mp_payments(str(fecha_desde), str(fecha_hasta)) or []
                if _mp_raw_det:
                    _montos_tn_det = {round(float(t)) for t in df_det["Total ($)"]} if not df_det.empty else set()
                    _df_mp_det = procesar_mp_payments(_mp_raw_det, montos_validos_tn=_montos_tn_det)
            _has_mp = not _df_mp_det.empty

            if _has_pn or _has_mp:
                _d1, _d2 = st.columns(2)
                if _has_pn:
                    with _d1.expander(f"🔵 Desglose Pago Nube ({len(_df_pn_det)})", expanded=False):
                        st.caption("Costo de procesamiento (fee) + retenciones impositivas (IIBB, IVA).")
                        _cols_pn = [
                            "Orden TN", "Fecha", "Método", "Cuotas",
                            "Monto ($)", "Fee ($)", "Retención ($)", "Costo total ($)", "Neto ($)",
                        ]
                        _cols_pn = [c for c in _cols_pn if c in _df_pn_det.columns]
                        _fmt_pn = {c: "${:,.2f}" for c in _cols_pn if "($)" in c}
                        st.dataframe(
                            _df_pn_det[_cols_pn].style.format(_fmt_pn),
                            use_container_width=True, hide_index=True,
                        )
                if _has_mp:
                    with _d2.expander(f"💳 Desglose Mercado Pago ({len(_df_mp_det)})", expanded=False):
                        st.caption("Cargo MP, financiación, impuestos, plataforma de terceros.")
                        _cols_mp = [
                            "ID MP", "Fecha", "Tipo", "Cuotas", "Medio",
                            "Bruto ($)", "Cargo MP ($)", "Cargo financiación ($)",
                            "Impuestos ($)", "Cargo plataforma ($)",
                            "Fee total ($)", "Costo %", "Neto ($)",
                        ]
                        _cols_mp = [c for c in _cols_mp if c in _df_mp_det.columns]
                        _fmt_money = {c: "${:,.2f}" for c in _cols_mp if "($)" in c}
                        _fmt_money["Costo %"] = "{:.2f}%"
                        st.dataframe(
                            _df_mp_det[_cols_mp].style.format(_fmt_money),
                            use_container_width=True, hide_index=True,
                        )

            # Estado de costos
            st.divider()
            st.subheader("📋 Estado de costos por producto")
            todos_prods = set()
            for _, row in df_tn.iterrows():
                for p in str(row.get("Productos", "")).split(" / "):
                    p = p.strip()
                    if p:
                        todos_prods.add(p)

            filas_estado = []
            for p in sorted(todos_prods):
                fob = get_fob_usd(p, _costos_gs)
                filas_estado.append({
                    "Estado": "✅" if fob > 0 else "❌",
                    "Producto (nombre en TN)": p,
                    "FOB (USD)": fob if fob > 0 else None,
                    "Costo ARS": round(fob * _dolar_det) if fob > 0 else None,
                })
            df_estado = pd.DataFrame(filas_estado)
            sin_costo_n = len(df_estado[df_estado["Estado"] == "❌"])
            cc1, cc2 = st.columns(2)
            cc1.metric("✅ Con costo", len(df_estado) - sin_costo_n)
            cc2.metric("❌ Sin costo", sin_costo_n)
            if sin_costo_n > 0:
                st.info("💡 Cargá los costos faltantes en la solapa **💻 Costos de consolas**.")
            st.dataframe(
                df_estado.style.map(
                    lambda v: "background-color: #1a3a1a" if v == "✅" else ("background-color: #3a1a1a" if v == "❌" else ""),
                    subset=["Estado"],
                ).format({"FOB (USD)": lambda x: f"${x:.2f}" if x else "—", "Costo ARS": lambda x: f"${x:,.0f}" if x else "—"}),
                use_container_width=True, hide_index=True,
            )

        # Resumen por medio de pago
        st.divider()
        st.subheader("📊 Resumen por medio de pago")
        if not df_tn.empty:
            res = df_tn.groupby("Medio de Pago").agg(
                Cantidad=("Orden", "count"), Bruto=("Total ($)", "sum"),
                Comision=("Comision PN ($)", "sum"), Neto=("Neto cobrado ($)", "sum"),
                CostoProds=("Costo Productos ($)", "sum"), Margen=("Margen ($)", "sum"),
            ).reset_index()
            res["Costo %"] = (res["Comision"] / res["Bruto"] * 100).round(2).apply(fmt_pct)
            res["Margen %"] = (res["Margen"] / res["Bruto"] * 100).round(2).apply(fmt_pct)
            for col in ["Bruto", "Comision", "Neto", "CostoProds", "Margen"]:
                res[col] = res[col].apply(fmt)
            res.columns = ["Medio de Pago", "Cantidad", "Bruto ($)", "Comision PN ($)", "Neto ($)", "Costo Prods ($)", "Margen ($)", "Costo PN %", "Margen %"]
            st.dataframe(res, use_container_width=True, hide_index=True)

        if not df_pagos.empty:
            st.divider()
            st.subheader("💳 Detalle transacciones — Pago Nube")
            st.dataframe(
                df_pagos.style.format({"Monto ($)": "${:,.0f}", "Fee ($)": "${:,.0f}", "Neto ($)": "${:,.0f}"}),
                use_container_width=True, hide_index=True,
            )

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 3: SALUD FINANCIERA
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "💚 Salud Financiera":
        st.subheader("💚 Salud Financiera del Período")

        # ── Usar valores de session state (configurados en sidebar) ──
        dolar_default = int(dolar_blue) if dolar_blue else 1200
        tipo_cambio = st.session_state.tipo_cambio_sf or dolar_default
        pct_iva = st.session_state.pct_iva
        pauta_manual = st.session_state.pauta_manual

        st.divider()
        st.subheader("💰 Resumen financiero del período")

        if df_tn.empty:
            st.info("Buscá primero para ver los datos financieros.")
        else:
            _costos_gs_sf = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
            df_calc = df_tn.copy()
            # Usar comisiones reales de TN (no recalcular con tasas teóricas)
            df_calc["Neto cobrado ($)"] = df_calc["Total ($)"] - df_calc["Comision PN ($)"]
            df_calc["Costo Productos ($)"] = df_calc.apply(
                lambda r: costo_final_row(r, tipo_cambio, _costos_gs_sf), axis=1
            )
            df_calc["Margen ($)"] = (
                df_calc["Neto cobrado ($)"] - df_calc["Costo Productos ($)"] - df_calc["Envio costo ($)"]
            )
            df_calc["Margen (%)"] = df_calc.apply(
                lambda r: round((r["Margen ($)"] / r["Total ($)"] * 100) if r["Total ($)"] > 0 else 0, 2),
                axis=1,
            )

            facturacion_bruta = df_calc["Total ($)"].sum()
            comisiones_pn = df_calc["Comision PN ($)"].sum()
            neto_cobrado = df_calc["Neto cobrado ($)"].sum()
            costo_productos = df_calc["Costo Productos ($)"].sum()
            costo_envios = df_calc["Envio costo ($)"].sum()
            costo_iva = facturacion_bruta * (pct_iva / 100)
            margen_bruto = df_calc["Margen ($)"].sum()

            # ── Gastos fijos prorrateados ──
            gastos_fijos_saved = gs_read("GastosFijos") or {}
            total_gastos_fijos_mes = sum(
                v for k, v in gastos_fijos_saved.items()
                if isinstance(v, (int, float)) and v > 0
            )
            dias_periodo = max((fecha_hasta - fecha_desde).days + 1, 1)
            factor_prorrateo = dias_periodo / 30
            gastos_fijos_periodo = round(total_gastos_fijos_mes * factor_prorrateo)

            resultado_final = margen_bruto - costo_iva - pauta_manual - gastos_fijos_periodo

            # ── Helper: card KPI uniforme ──────────────────────────────────────
            def _kpi_html(label, value, sub="", val_color=None, accent_border=False):
                """Card de altura fija para alinear valores verticalmente."""
                vc = val_color or MG_TEXT
                border = f"border-left:2px solid {MG_RED};padding-left:0.75rem;" if accent_border else ""
                sub_html = (
                    f'<div style="font-size:0.68rem;color:{MG_MUTED};margin-top:0.25rem;'
                    f'white-space:nowrap;overflow:hidden;text-overflow:ellipsis;">{sub}</div>'
                    if sub else
                    f'<div style="font-size:0.68rem;color:transparent;margin-top:0.25rem;">—</div>'
                )
                return (
                    f'<div style="background:{MG_SURF};border-radius:8px;padding:0.9rem 1rem;'
                    f'min-height:90px;display:flex;flex-direction:column;justify-content:flex-start;{border}">'
                    f'<div style="font-size:0.58rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                    f'text-transform:uppercase;letter-spacing:0.08em;margin-bottom:0.4rem;'
                    f'white-space:nowrap;overflow:hidden;text-overflow:ellipsis;">{label}</div>'
                    f'<div style="font-size:1.35rem;font-weight:700;color:{vc};line-height:1.1;">{value}</div>'
                    f'{sub_html}'
                    f'</div>'
                )

            # ── Cargar datos MP del período (para detalle pasarela) ────────────
            _mp_raw_sf = []
            _df_mp_sf = pd.DataFrame()
            if MP_ACCESS_TOKEN:
                _mp_raw_sf = get_mp_payments(str(fecha_desde), str(fecha_hasta)) or []
                if _mp_raw_sf:
                    _montos_tn_sf = {round(float(t)) for t in df_calc["Total ($)"]} if not df_calc.empty else set()
                    _df_mp_sf = procesar_mp_payments(_mp_raw_sf, montos_validos_tn=_montos_tn_sf)

            # ── 1) RESUMEN DEL PERÍODO — 4 cards ──────────────────────────────
            comisiones_pasarela = comisiones_pn  # ya incluye PN + MP (calculado en procesar_orders)
            neto_cobrado = facturacion_bruta - comisiones_pasarela
            costo_operativo = costo_productos + costo_envios

            st.markdown(
                f'<p style="font-size:0.62rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                f'letter-spacing:0.08em;text-transform:uppercase;margin-bottom:0.5rem;">'
                f'💰 Resumen del período</p>',
                unsafe_allow_html=True,
            )
            _r1, _r2, _r3, _r4 = st.columns(4)
            _pct_total = (comisiones_pasarela / facturacion_bruta * 100) if facturacion_bruta > 0 else 0.0
            _r1.markdown(_kpi_html("Facturación bruta", fmt(facturacion_bruta), "total vendido"), unsafe_allow_html=True)
            _r2.markdown(_kpi_html("Comisiones pasarela", fmt(comisiones_pasarela), f"PN + MP · {_pct_total:.2f}%", val_color=MG_RED), unsafe_allow_html=True)
            _r3.markdown(_kpi_html("Neto cobrado", fmt(neto_cobrado), "bruto − comisiones"), unsafe_allow_html=True)
            _r4.markdown(_kpi_html("Costo operativo", fmt(costo_operativo), f"prods {fmt(costo_productos)} + envíos {fmt(costo_envios)}", val_color=MG_RED), unsafe_allow_html=True)

            # ── 2) DETALLE POR PASARELA ────────────────────────────────────────
            st.markdown(
                f'<p style="font-size:0.62rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                f'letter-spacing:0.08em;text-transform:uppercase;margin-top:1.5rem;margin-bottom:0.5rem;">'
                f'🔀 Detalle por pasarela</p>',
                unsafe_allow_html=True,
            )

            _mask_pn = df_calc["Pasarela"] == "PN"
            _mask_mp = df_calc["Pasarela"].isin(["MP", "Convenir"]) if "Pasarela" in df_calc.columns else pd.Series([False]*len(df_calc))

            _bruto_pn = df_calc.loc[_mask_pn, "Total ($)"].sum() if _mask_pn.any() else 0
            _com_pn   = df_calc.loc[_mask_pn, "Comision PN ($)"].sum() if _mask_pn.any() else 0
            _ord_pn   = int(_mask_pn.sum())
            _pct_pn   = (_com_pn / _bruto_pn * 100) if _bruto_pn > 0 else 0.0

            _bruto_mp = df_calc.loc[_mask_mp, "Total ($)"].sum() if _mask_mp.any() else 0
            _com_mp   = df_calc.loc[_mask_mp, "Comision PN ($)"].sum() if _mask_mp.any() else 0
            _ord_mp   = int(_mask_mp.sum())
            _pct_mp   = (_com_mp / _bruto_mp * 100) if _bruto_mp > 0 else 0.0

            def _pasarela_card(titulo, icon, ordenes, bruto, comision, pct, color):
                return (
                    f'<div style="background:{MG_SURF};border-radius:8px;padding:1rem 1.2rem;'
                    f'border-left:3px solid {color};">'
                    f'<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:0.8rem;">'
                    f'<div style="font-size:0.85rem;font-weight:700;color:{MG_TEXT};">{icon} {titulo}</div>'
                    f'<div style="font-size:0.7rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;">{ordenes} órdenes</div>'
                    f'</div>'
                    f'<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:0.8rem;">'
                    f'<div>'
                    f'<div style="font-size:0.55rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;text-transform:uppercase;letter-spacing:0.08em;">Bruto</div>'
                    f'<div style="font-size:1.05rem;font-weight:700;color:{MG_TEXT};">{fmt(bruto)}</div>'
                    f'</div>'
                    f'<div>'
                    f'<div style="font-size:0.55rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;text-transform:uppercase;letter-spacing:0.08em;">Comisión</div>'
                    f'<div style="font-size:1.05rem;font-weight:700;color:{MG_RED};">{fmt(comision)}</div>'
                    f'</div>'
                    f'<div>'
                    f'<div style="font-size:0.55rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;text-transform:uppercase;letter-spacing:0.08em;">Costo %</div>'
                    f'<div style="font-size:1.05rem;font-weight:700;color:{color};">{pct:.2f}%</div>'
                    f'</div>'
                    f'</div>'
                    f'</div>'
                )

            _pc1, _pc2 = st.columns(2)
            _pc1.markdown(
                _pasarela_card("Pago Nube", "🔵", _ord_pn, _bruto_pn, _com_pn, _pct_pn, "#009EE3"),
                unsafe_allow_html=True,
            )
            _pc2.markdown(
                _pasarela_card("Mercado Pago", "💳", _ord_mp, _bruto_mp, _com_mp, _pct_mp, MG_RED),
                unsafe_allow_html=True,
            )

            # Validación cruzada con MP API (si los datos no cuadran avisamos)
            if not _df_mp_sf.empty and _bruto_mp > 0:
                _mp_api_bruto = _df_mp_sf["Bruto ($)"].sum()
                _mp_api_fee   = _df_mp_sf["Fee total ($)"].sum()
                _diff_pct = abs(_mp_api_bruto - _bruto_mp) / _bruto_mp * 100 if _bruto_mp > 0 else 0
                if _diff_pct > 5:  # más de 5% de diferencia
                    st.caption(
                        f"⚠️ Diferencia entre MP (API) y TN (órdenes MP): "
                        f"API ${_mp_api_bruto:,.0f} vs TN ${_bruto_mp:,.0f} "
                        f"({_diff_pct:.1f}% diff) · Fee real API: {fmt(_mp_api_fee)}"
                    )
                else:
                    st.caption(
                        f"✅ Conciliado con API MP: Fee real {fmt(_mp_api_fee)} "
                        f"vs estimado {fmt(_com_mp)}"
                    )

            # ── Debug MP (sigue disponible) ────────────────────────────────────
            if MP_ACCESS_TOKEN:
                # Persistencia de resultados en session state
                if "mp_debug_payment" not in st.session_state:
                    st.session_state.mp_debug_payment = None
                if "mp_debug_sample" not in st.session_state:
                    st.session_state.mp_debug_sample = None

                with st.expander("🔧 Debug MP — inspeccionar estructura de pagos", expanded=False):
                    st.caption("Útil para ver qué campos vienen en la API y por qué el filtro no captura una liquidación.")

                    # OPCIÓN A: Buscar por ID
                    st.markdown("**A) Buscar pago por ID de operación**")
                    with st.form("dbg_mp_form_id"):
                        _dbg_id = st.text_input("ID operación MP", placeholder="156885323147")
                        _dbg_submit = st.form_submit_button("Buscar pago", use_container_width=True)
                        if _dbg_submit and _dbg_id:
                            try:
                                _r = requests.get(
                                    f"https://api.mercadopago.com/v1/payments/{_dbg_id.strip()}",
                                    headers={"Authorization": f"Bearer {MP_ACCESS_TOKEN}"},
                                    timeout=10,
                                )
                                if _r.status_code == 200:
                                    st.session_state.mp_debug_payment = _r.json()
                                else:
                                    st.session_state.mp_debug_payment = {
                                        "_error": f"HTTP {_r.status_code}: {_r.text[:400]}"
                                    }
                            except Exception as _e:
                                st.session_state.mp_debug_payment = {"_error": str(_e)}

                    # Renderizar resultado de búsqueda (persiste tras rerun)
                    _dbg_data = st.session_state.mp_debug_payment
                    if _dbg_data:
                        if "_error" in _dbg_data:
                            st.error(_dbg_data["_error"])
                        else:
                            st.markdown("**Campos clave del pago:**")
                            _resumen = {
                                "id": _dbg_data.get("id"),
                                "status": _dbg_data.get("status"),
                                "payment_type_id": _dbg_data.get("payment_type_id"),
                                "payment_method_id": _dbg_data.get("payment_method_id"),
                                "transaction_amount": _dbg_data.get("transaction_amount"),
                                "description": _dbg_data.get("description"),
                                "statement_descriptor": _dbg_data.get("statement_descriptor"),
                                "payer.identification": _dbg_data.get("payer", {}).get("identification"),
                                "payer.first_name": _dbg_data.get("payer", {}).get("first_name"),
                                "payer.last_name": _dbg_data.get("payer", {}).get("last_name"),
                                "payer.email": _dbg_data.get("payer", {}).get("email"),
                                "payer.entity_type": _dbg_data.get("payer", {}).get("entity_type"),
                                "additional_info.payer": _dbg_data.get("additional_info", {}).get("payer"),
                            }
                            st.json(_resumen, expanded=True)
                            st.markdown(
                                f"**¿Detectado como liquidación externa?** "
                                f"`{_es_liquidacion_externa(_dbg_data)}` · "
                                f"CUIT extraído: `{_cuit_payer(_dbg_data) or '(vacío)'}`"
                            )
                            with st.expander("Ver JSON completo", expanded=False):
                                st.json(_dbg_data)

                    st.divider()
                    # OPCIÓN B: Dump de los pagos del período actual ordenados por monto
                    st.markdown("**B) Top 5 pagos del período por monto** (para ver de dónde viene el bruto inflado)")
                    if st.button("Mostrar top 5 pagos crudos", use_container_width=True, key="dbg_mp_sample_btn"):
                        sorted_pagos = sorted(_mp_raw_sf, key=lambda x: float(x.get("transaction_amount", 0)), reverse=True)[:5]
                        st.session_state.mp_debug_sample = sorted_pagos

                    if st.session_state.mp_debug_sample:
                        for _i, _p in enumerate(st.session_state.mp_debug_sample, 1):
                            _bruto = float(_p.get("transaction_amount", 0))
                            _tipo = _p.get("payment_type_id", "")
                            _liq = _es_liquidacion_externa(_p)
                            _cuit = _cuit_payer(_p)
                            st.markdown(
                                f"**{_i}.** ID `{_p.get('id')}` · "
                                f"${_bruto:,.0f} · "
                                f"tipo: `{_tipo}` · "
                                f"CUIT payer: `{_cuit or '(vacío)'}` · "
                                f"filtrado: `{_liq}`"
                            )
                            _resumen_p = {
                                "payment_type_id": _p.get("payment_type_id"),
                                "payment_method_id": _p.get("payment_method_id"),
                                "description": _p.get("description"),
                                "statement_descriptor": _p.get("statement_descriptor"),
                                "payer": _p.get("payer"),
                                "additional_info.payer": _p.get("additional_info", {}).get("payer"),
                            }
                            with st.expander(f"Ver campos pago #{_i}", expanded=False):
                                st.json(_resumen_p)

            # ── 3) GASTOS DEL PERÍODO — 4 columnas ────────────────────────────
            st.markdown(
                f'<p style="font-size:0.62rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                f'letter-spacing:0.08em;text-transform:uppercase;margin-top:1.5rem;margin-bottom:0.5rem;">'
                f'📊 Gastos del período</p>',
                unsafe_allow_html=True,
            )
            g1, g2, g3, g4 = st.columns(4)
            g1.markdown(_kpi_html("Margen bruto", fmt(margen_bruto), "neto − costos − envíos"), unsafe_allow_html=True)
            g2.markdown(_kpi_html(f"IVA ({pct_iva:.1f}%)", fmt(costo_iva), f"−{fmt(costo_iva)}", val_color=MG_RED), unsafe_allow_html=True)
            g3.markdown(_kpi_html("Pauta publicitaria", fmt(pauta_manual), f"−{fmt(pauta_manual)}" if pauta_manual else "sin datos aún", val_color=MG_RED if pauta_manual else MG_MUTED), unsafe_allow_html=True)
            g4.markdown(_kpi_html(
                f"Gastos fijos ({dias_periodo}d)",
                fmt(gastos_fijos_periodo),
                f"${total_gastos_fijos_mes:,.0f}/mes × {factor_prorrateo:.2f}",
                val_color=MG_RED,
            ), unsafe_allow_html=True)

            # ── Resultado final ────────────────────────────────────────────────
            _res_color = "#4ade80" if resultado_final >= 0 else MG_RED
            _res_icon  = "🟢" if resultado_final >= 0 else "🔴"
            st.markdown(
                f'<div style="background:{MG_SURF};border-radius:10px;padding:1.2rem 1.5rem;'
                f'border-left:3px solid {_res_color};margin-top:1.5rem;margin-bottom:1rem;">'
                f'<div style="font-size:0.62rem;color:{MG_MUTED};font-family:\'Space Mono\',monospace;'
                f'text-transform:uppercase;letter-spacing:0.08em;margin-bottom:0.4rem;">'
                f'{_res_icon} Resultado final del período</div>'
                f'<div style="font-size:2rem;font-weight:700;color:{_res_color};">{fmt(resultado_final)}</div>'
                f'</div>',
                unsafe_allow_html=True,
            )

            # ── Gráfico de resultado día a día ──
            st.markdown("<div style='margin-top:1.5rem'></div>", unsafe_allow_html=True)
            st.subheader("📈 Resultado diario — Tendencia del período")

            # Calcular P&L por día
            df_daily = df_calc.groupby("Fecha").agg(
                Facturacion=("Total ($)", "sum"),
                Comision=("Comision PN ($)", "sum"),
                Costo_Prods=("Costo Productos ($)", "sum"),
                Costo_Envio=("Envio costo ($)", "sum"),
            ).reset_index()

            # Prorratear gastos diarios
            iva_diario = df_daily["Facturacion"] * (pct_iva / 100)
            pauta_diaria = pauta_manual / dias_periodo
            gastos_fijos_diario = gastos_fijos_periodo / dias_periodo

            df_daily["Margen_Bruto"] = (
                df_daily["Facturacion"] - df_daily["Comision"] - df_daily["Costo_Prods"] - df_daily["Costo_Envio"]
            )
            df_daily["Resultado"] = (
                df_daily["Margen_Bruto"] - iva_diario - pauta_diaria - gastos_fijos_diario
            ).round(0)
            df_daily["Resultado_Acum"] = df_daily["Resultado"].cumsum()

            # Colores de marca
            _COLOR_POS = "#4ade80"   # verde
            _COLOR_NEG = MG_RED      # rojo MG
            _COLOR_ACUM = "#009EE3"  # azul
            df_daily["Color"] = df_daily["Resultado"].apply(lambda x: _COLOR_POS if x >= 0 else _COLOR_NEG)

            # Label compacto: solo K/M para no solapar
            def _lbl(v):
                if abs(v) >= 1_000_000:
                    return f"${v/1_000_000:.1f}M"
                elif abs(v) >= 1_000:
                    return f"${v/1_000:.0f}K"
                return f"${v:.0f}"

            df_daily["Label"] = df_daily["Resultado"].apply(_lbl)

            fig_daily = go.Figure()

            # Barras de resultado diario
            fig_daily.add_trace(go.Bar(
                x=df_daily["Fecha"], y=df_daily["Resultado"],
                name="Resultado diario",
                marker_color=df_daily["Color"],
                marker_line_width=0,
                text=df_daily["Label"],
                textposition="outside",
                textfont=dict(size=10, color=MG_MUTED),
                hovertemplate="<b>%{x}</b><br>Resultado: $%{y:,.0f}<extra></extra>",
            ))

            # Línea acumulada (eje derecho)
            fig_daily.add_trace(go.Scatter(
                x=df_daily["Fecha"], y=df_daily["Resultado_Acum"],
                name="Acumulado",
                mode="lines+markers",
                line=dict(color=_COLOR_ACUM, width=2),
                marker=dict(size=5, color=_COLOR_ACUM),
                yaxis="y2",
                hovertemplate="Acumulado: $%{y:,.0f}<extra></extra>",
            ))

            fig_daily.update_layout(
                yaxis=dict(
                    title="", tickformat="$,.0f",
                    gridcolor=MG_BORDER, zeroline=True, zerolinecolor=MG_MUTED, zerolinewidth=1,
                ),
                yaxis2=dict(
                    title="", overlaying="y", side="right",
                    tickformat="$,.0f", showgrid=False,
                    tickfont=dict(color=_COLOR_ACUM, size=10),
                ),
                legend=dict(orientation="h", y=1.08, x=0, xanchor="left"),
                hovermode="x unified",
                bargap=0.3,
                height=400,
            )
            st.plotly_chart(fig_daily, use_container_width=True)

            # Indicador de tendencia
            if len(df_daily) >= 3:
                ultimos_3 = df_daily["Resultado"].tail(3).mean()
                primeros_3 = df_daily["Resultado"].head(3).mean()
                dias_positivos = (df_daily["Resultado"] >= 0).sum()
                dias_totales = len(df_daily)
                if ultimos_3 > primeros_3 * 1.1:
                    st.success(f"📈 Tendencia alcista — últimos días {_lbl(ultimos_3)}/día vs {_lbl(primeros_3)}/día al arranque. {dias_positivos}/{dias_totales} días en verde.")
                elif ultimos_3 < primeros_3 * 0.9:
                    st.warning(f"📉 Tendencia bajista — últimos días {_lbl(ultimos_3)}/día vs {_lbl(primeros_3)}/día al arranque. {dias_positivos}/{dias_totales} días en verde.")
                else:
                    st.info(f"➡️ Tendencia estable — promedio {_lbl(df_daily['Resultado'].mean())}/día. {dias_positivos}/{dias_totales} días en verde.")

            # Detalle por orden
            st.divider()
            st.subheader("📋 Detalle por orden con margen real")
            cols_fin = [
                "Orden", "Fecha", "Cliente", "Medio de Pago", "Cuotas", "Total ($)",
                "Comision PN ($)", "Neto cobrado ($)", "Costo Productos ($)",
                "Envio costo ($)", "Margen ($)", "Margen (%)",
            ]
            cols_fin = [c for c in cols_fin if c in df_calc.columns]
            st.dataframe(
                df_calc[cols_fin].style.format({
                    "Total ($)": "${:,.0f}", "Comision PN ($)": "${:,.0f}",
                    "Neto cobrado ($)": "${:,.0f}", "Costo Productos ($)": "${:,.0f}",
                    "Envio costo ($)": "${:,.0f}", "Margen ($)": "${:,.0f}", "Margen (%)": "{:.1f}%",
                }),
                use_container_width=True, hide_index=True,
            )

            # Cascada
            st.divider()
            st.subheader("📊 Cascada de resultados")
            wf = pd.DataFrame({
                "Concepto": [
                    "Facturación bruta", "Comisiones PN", "Costo productos",
                    "Costo envíos", f"IVA ({pct_iva:.1f}%)", "Pauta",
                    "Gastos fijos", "Resultado final",
                ],
                "Monto": [
                    facturacion_bruta, -comisiones_pn, -costo_productos,
                    -costo_envios, -costo_iva, -pauta_manual,
                    -gastos_fijos_periodo, resultado_final,
                ],
                "Color": [
                    "#00C49F", "#FF9900", "#FF5733", "#FF5733", "#FF5733",
                    "#FF9900", "#FF9900",
                    "#009EE3" if resultado_final >= 0 else "#FF0000",
                ],
            })
            fig_wf = px.bar(wf, x="Concepto", y="Monto", color="Concepto",
                color_discrete_sequence=wf["Color"].tolist(), title="Cascada financiera")
            fig_wf.update_layout(showlegend=False, yaxis_tickformat="$,.0f")
            fig_wf.update_traces(texttemplate="%{y:$,.0f}", textposition="outside")
            st.plotly_chart(fig_wf, use_container_width=True)

            # ── Desglose por pasarela (MP vs PN) ──
            if "Pasarela" in df_calc.columns:
                st.divider()
                st.subheader("🔀 Comisiones por pasarela")
                agg_pasarela = (
                    df_calc.groupby("Pasarela")
                    .agg(
                        Órdenes=("Orden", "count"),
                        Facturación=("Total ($)", "sum"),
                        Comisión=("Comision PN ($)", "sum"),
                        Neto=("Neto cobrado ($)", "sum"),
                    )
                    .reset_index()
                )
                agg_pasarela["Costo %"] = (
                    agg_pasarela["Comisión"] / agg_pasarela["Facturación"] * 100
                ).round(2)

                pas_c1, pas_c2 = st.columns(2)
                for _, row_p in agg_pasarela.iterrows():
                    col_p = pas_c1 if row_p["Pasarela"] == "PN" else pas_c2
                    label = "Pago Nube" if row_p["Pasarela"] == "PN" else "Mercado Pago"
                    icon = "💳" if row_p["Pasarela"] == "MP" else "🔵"
                    with col_p:
                        st.markdown(
                            f'<p style="font-size:0.68rem;color:{MG_MUTED};font-family:\'Space Mono\','
                            f'monospace;letter-spacing:0.08em;text-transform:uppercase;margin-bottom:0.3rem;">'
                            f'{icon} {label}</p>',
                            unsafe_allow_html=True,
                        )
                        pm1, pm2, pm3 = st.columns(3)
                        pm1.metric("Órdenes", int(row_p["Órdenes"]))
                        pm2.metric("Facturación", fmt(row_p["Facturación"]))
                        pm3.metric("Comisión", fmt(row_p["Comisión"]), help=f"Costo %: {row_p['Costo %']:.2f}%")

                # Si hay token de MP, mostrar fees reales
                if MP_ACCESS_TOKEN:
                    with st.spinner("Consultando fees reales de Mercado Pago..."):
                        mp_raw = get_mp_payments(str(fecha_desde), str(fecha_hasta))
                    if mp_raw:
                        _montos_tn_legacy = {round(float(t)) for t in df_calc["Total ($)"]} if not df_calc.empty else set()
                        df_mp = procesar_mp_payments(mp_raw, montos_validos_tn=_montos_tn_legacy)
                        if not df_mp.empty:
                            st.markdown("**Fees reales desde API de Mercado Pago**")
                            # Resumen por tipo
                            mp_res = (
                                df_mp.groupby("Tipo")
                                .agg(
                                    Ops=("ID MP", "count"),
                                    Bruto=("Bruto ($)", "sum"),
                                    FeeTotal=("Fee total ($)", "sum"),
                                    Neto=("Neto ($)", "sum"),
                                )
                                .reset_index()
                            )
                            mp_res["Costo %"] = (
                                mp_res["FeeTotal"] / mp_res["Bruto"] * 100
                            ).round(2)
                            mp_res_fmt = mp_res.copy()
                            for col in ["Bruto", "FeeTotal", "Neto"]:
                                mp_res_fmt[col] = mp_res[col].apply(fmt)
                            mp_res_fmt["Costo %"] = mp_res["Costo %"].apply(fmt_pct)
                            mp_res_fmt.columns = [
                                "Tipo", "Ops", "Bruto ($)", "Fee MP ($)", "Neto ($)", "Costo %",
                            ]
                            st.dataframe(mp_res_fmt, use_container_width=True, hide_index=True)
                    else:
                        st.caption("No se encontraron pagos MP en este período (o el token venció).")
                else:
                    st.caption(
                        f'💡 Agregá `MP_ACCESS_TOKEN` en secrets para ver fees reales de MP '
                        f'(comisión + costo financiero por operación).'
                    )

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 4: STOCK
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "📦 Stock":
        st.subheader("📦 Stock — Tienda Nube")

        @st.fragment
        def cargar_stock():
            if st.button("🔄 Cargar stock desde Tienda Nube", use_container_width=False):
                with st.spinner("Consultando productos y stock..."):
                    productos = get_tn_products()
                if productos:
                    stock_rows = []
                    for p in productos:
                        nombre_raw = p.get("name", {})
                        nombre = nombre_raw.get("es", "") if isinstance(nombre_raw, dict) else str(nombre_raw)
                        for v in p.get("variants", []):
                            stock = v.get("stock", None)
                            vals = v.get("values", [])
                            variante = " / ".join([
                                (val.get("es", "") or next(iter(val.values()), ""))
                                if isinstance(val, dict) else str(val)
                                for val in vals
                            ]) if vals else "—"
                            stock_rows.append({
                                "Producto": nombre,
                                "Variante": variante,
                                "SKU": v.get("sku", ""),
                                "Stock": stock if stock is not None else "Sin límite",
                                "Precio ($)": float(v.get("price", 0) or 0),
                            })
                    st.session_state.stock_tn = pd.DataFrame(stock_rows)
                    _snap_map = {}
                    for _r in stock_rows:
                        _sv = _r["Stock"]
                        if isinstance(_sv, (int, float)):
                            _snap_map[_r["Producto"]] = _snap_map.get(_r["Producto"], 0) + int(_sv)
                    gs_append_snapshot(_snap_map)
                    st.success(f"✅ {len(stock_rows)} variantes cargadas")
                else:
                    st.warning("No se pudieron cargar productos.")

            if "stock_tn" in st.session_state and st.session_state.stock_tn is not None:
                df_stock = st.session_state.stock_tn
                col_f1, col_f2 = st.columns(2)
                buscar_prod = col_f1.text_input("🔍 Buscar producto", "")
                mostrar_sin_stock = col_f2.checkbox("Solo sin stock o bajo", value=False)

                df_stock_f = df_stock.copy()
                if buscar_prod:
                    df_stock_f = df_stock_f[df_stock_f["Producto"].str.contains(buscar_prod, case=False, na=False)]
                if mostrar_sin_stock:
                    df_stock_f = df_stock_f[df_stock_f["Stock"].apply(lambda x: isinstance(x, (int, float)) and x <= 3)]

                st.dataframe(
                    df_stock_f.style.format({"Precio ($)": "${:,.0f}"}),
                    use_container_width=True, hide_index=True,
                )

                st.divider()
                st.subheader("⚠️ Alertas de stock bajo")
                umbral = st.slider("Umbral (unidades)", 1, 20, 5)
                alertas = df_stock[df_stock["Stock"].apply(lambda x: isinstance(x, (int, float)) and x <= umbral)]
                if alertas.empty:
                    st.success(f"✅ Todos con más de {umbral} unidades.")
                else:
                    st.warning(f"⚠️ {len(alertas)} variante(s) con stock ≤ {umbral}")
                    st.dataframe(alertas.style.format({"Precio ($)": "${:,.0f}"}),
                        use_container_width=True, hide_index=True)

                st.divider()
                stock_num = df_stock[df_stock["Stock"].apply(lambda x: isinstance(x, (int, float)))]
                r1, r2, r3 = st.columns(3)
                r1.metric("Total productos", df_stock["Producto"].nunique())
                r2.metric("Total variantes", len(df_stock))
                r3.metric("Total unidades", int(stock_num["Stock"].sum()))
                st.download_button("⬇️ Descargar stock",
                    df_stock.to_csv(index=False).encode("utf-8"), "stock_marketgamer.csv", "text/csv")
            else:
                st.info("👆 Hacé clic en 'Cargar stock' para ver el inventario.")

        cargar_stock()

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 5: VELOCIDAD DE VENTAS
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🔥 Velocidad de ventas":
        from velocidad_restock import calcular_velocidad_restock
        st.subheader("🔥 Velocidad de ventas y planificación de restock")
        st.caption("📌 Esta solapa usa su propia ventana histórica, independiente del período del panel lateral.")

        with st.expander("⚙️ Configuración de alertas", expanded=False):
            cg, ch = st.columns(2)
            p_hist = cg.slider("📚 Historia a analizar (días)", 90, 1095, 365)
            p_vent = ch.slider("🕐 Ventana reciente (días)", 14, 180, 90)
            ca, cb, cc = st.columns(3)
            p_lead = ca.slider("📦 Lead time (días)", 1, 45, 7)
            p_colchon = cb.slider("🛟 Colchón de seguridad (días)", 0, 30, 7)
            p_cob = cc.slider("🎯 Cobertura objetivo (días)", 7, 90, 30)
            ce, cf = st.columns(2)
            p_minu = ce.slider("Mín. unidades p/ confianza", 1, 20, 5)
            p_mind = cf.slider("Mín. días distintos p/ confianza", 1, 10, 3)

        df_full = _cargar_ordenes_historico(p_hist)
        if df_full is None or df_full.empty:
            st.info("No hay órdenes en el rango histórico seleccionado.")
        else:
            stock_map = {}
            precio_map = {}
            stock_df = st.session_state.get("stock_tn")
            if stock_df is not None:
                for _, srow in stock_df.iterrows():
                    pn = srow["Producto"]
                    sv = srow["Stock"]
                    if isinstance(sv, (int, float)):
                        stock_map[pn] = stock_map.get(pn, 0) + int(sv)
                    pr = srow.get("Precio ($)", 0)
                    if pr and pn not in precio_map:
                        precio_map[pn] = float(pr)

            if not stock_map:
                st.warning(
                    "📦 No cargaste el stock todavía. Andá a la solapa **Stock** y tocá "
                    "**'Cargar stock desde Tienda Nube'**. Sin stock no se puede calcular el "
                    "restock: todos los productos figuran como 'Sin límite'."
                )

            historial = gs_read("HistorialStock") or {}

            params = {
                "lead_time": p_lead, "colchon": p_colchon, "cobertura": p_cob,
                "ventana_reciente": p_vent, "min_unidades_conf": p_minu,
                "min_dias_conf": p_mind,
            }

            df_vel = calcular_velocidad_restock(
                df_full, stock_map, historial, precio_map, params,
                date.today().isoformat(),
            )

            if df_vel.empty:
                st.info("No hay datos de productos.")
            else:
                df_ok = df_vel[df_vel["Confianza"] == "ok"]
                df_baja = df_vel[df_vel["Confianza"] == "baja"]
                criticos = df_ok[df_ok["_necesita_restock"]]

                v1, v2, v3, v4 = st.columns(4)
                v1.metric("Productos", len(df_vel))
                v2.metric("🔴 A reponer", len(criticos))
                v3.metric("Ventana reciente", f"{p_vent} días")
                v4.metric("Snapshots", len(historial))

                if not criticos.empty:
                    st.divider()
                    st.error(f"⚠️ {len(criticos)} producto(s) necesitan restock")
                    for _, crow in criticos.iterrows():
                        st.warning(
                            f"🔴 **{crow['Producto']}** — vel. reciente "
                            f"{crow['Vel. reciente']} u/día | stock {crow['Stock actual']} "
                            f"| ROP {crow['ROP']} | pedir **{crow['Restock sugerido']} u** "
                            f"| riesgo ${crow['Facturación en riesgo']:,.0f}"
                        )

                st.divider()
                df_chart = (df_ok[["Producto", "Vel. reciente"]]
                            .sort_values("Vel. reciente", ascending=True).tail(15))
                if not df_chart.empty:
                    fig_vel = px.bar(
                        df_chart, x="Vel. reciente", y="Producto", orientation="h",
                        title="Velocidad reciente (u/día)", color="Vel. reciente",
                        color_continuous_scale=["#00C49F", "#FFD700", "#FF5733"],
                        text="Vel. reciente",
                    )
                    fig_vel.update_layout(showlegend=False, coloraxis_showscale=False,
                                          yaxis={"categoryorder": "total ascending"})
                    fig_vel.update_traces(texttemplate="%{text:.2f}", textposition="outside")
                    st.plotly_chart(fig_vel, use_container_width=True)

                cols_show = [
                    "Producto", "Unidades", "Vel. histórica", "Vel. reciente",
                    "Stock actual", "ROP", "Días restantes", "Restock sugerido",
                    "Facturación en riesgo",
                ]
                st.dataframe(
                    df_ok[cols_show].style.format({
                        "Vel. histórica": "{:.3f}", "Vel. reciente": "{:.3f}",
                        "ROP": "{:.1f}", "Facturación en riesgo": "${:,.0f}",
                    }),
                    use_container_width=True, hide_index=True,
                )

                if not df_baja.empty:
                    with st.expander(f"🔸 {len(df_baja)} productos de baja confianza (pocas ventas)", expanded=False):
                        st.dataframe(
                            df_baja[["Producto", "Unidades", "Vel. reciente", "Stock actual"]],
                            use_container_width=True, hide_index=True,
                        )

                st.divider()
                st.subheader("📈 Evolución de ventas")
                prod_sel = st.selectbox("Producto", df_vel["Producto"].tolist())
                if prod_sel:
                    df_evo = df_full[df_full["Productos"].str.contains(prod_sel, na=False, regex=False)]
                    if not df_evo.empty:
                        df_evo_g = df_evo.groupby("Fecha").agg(
                            Unidades=("Cantidad", "sum"), Revenue=("Total ($)", "sum"),
                        ).reset_index()
                        df_evo_g["Acumulado"] = df_evo_g["Unidades"].cumsum()
                        fig_evo = px.bar(df_evo_g, x="Fecha", y="Unidades",
                            title=f"Ventas diarias — {prod_sel}", color_discrete_sequence=["#009EE3"])
                        fig_evo.add_scatter(x=df_evo_g["Fecha"], y=df_evo_g["Acumulado"],
                            mode="lines+markers", name="Acumulado",
                            line=dict(color="#FFD700", width=2), yaxis="y2")
                        fig_evo.update_layout(
                            yaxis2=dict(overlaying="y", side="right", showgrid=False),
                            legend=dict(orientation="h"),
                        )
                        st.plotly_chart(fig_evo, use_container_width=True)

                st.download_button("⬇️ Descargar análisis",
                    df_vel[cols_show].to_csv(index=False).encode("utf-8"),
                    "restock_analysis.csv", "text/csv")

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 6: GASTOS FIJOS
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🏗️ Gastos fijos":
        st.subheader("🏗️ Gastos fijos mensuales")
        st.caption("Los datos quedan guardados en Google Sheets.")

        if "gastos_fijos" not in st.session_state:
            saved = gs_read("GastosFijos")
            st.session_state.gastos_fijos = saved if saved else {
                "Bruno": 0, "Coco": 0, "Agencia": 0, "Local": 0,
                "Sueldo Agus": 0, "Sueldo Stella": 0, "Contador": 0,
                "Sueldo Facu": 0, "Otros": 0,
            }

        gastos = st.session_state.gastos_fijos.copy()

        st.markdown("**Editá, agregá o eliminá gastos fijos (ARS):**")
        nuevos_gastos = {}
        gastos_a_eliminar = []
        items = list(gastos.items())

        for i, (k, v) in enumerate(items):
            col_name, col_val, col_del = st.columns([3, 2, 0.5])
            with col_name:
                st.markdown(f"**{k}**")
            with col_val:
                nuevos_gastos[k] = st.number_input(
                    f"monto_{k}", value=int(v), step=50_000,
                    key=f"gf_{k}", label_visibility="collapsed",
                )
            with col_del:
                if st.button("🗑️", key=f"del_{k}", help=f"Eliminar {k}"):
                    gastos_a_eliminar.append(k)

        if gastos_a_eliminar:
            for k in gastos_a_eliminar:
                nuevos_gastos.pop(k, None)
                st.session_state.gastos_fijos.pop(k, None)
            gs_write("GastosFijos", nuevos_gastos)
            st.session_state.gastos_fijos = nuevos_gastos
            st.rerun()

        st.divider()
        with st.expander("➕ Agregar gasto"):
            ng1, ng2, ng3 = st.columns([2, 2, 1])
            nuevo_nombre = ng1.text_input("Nombre")
            nuevo_monto = ng2.number_input("Monto mensual (ARS)", value=0, step=50_000, key="ng_monto")
            with ng3:
                st.write("")
                st.write("")
                if st.button("Agregar", type="primary"):
                    if nuevo_nombre and nuevo_nombre not in nuevos_gastos:
                        nuevos_gastos[nuevo_nombre] = nuevo_monto
                        st.session_state.gastos_fijos = nuevos_gastos
                        gs_write("GastosFijos", nuevos_gastos)
                        st.rerun()

        if st.button("💾 Guardar gastos fijos", use_container_width=True, type="primary"):
            st.session_state.gastos_fijos = nuevos_gastos
            ok = gs_write("GastosFijos", nuevos_gastos)
            st.success("✅ Guardado en Google Sheets" if ok else "⚠️ Solo en sesión")

        st.divider()
        total_gastos = sum(v for v in nuevos_gastos.values() if isinstance(v, (int, float)))
        st.metric("💰 Total gastos fijos mensuales", fmt(total_gastos))

        df_gastos = pd.DataFrame(
            [(k, v) for k, v in nuevos_gastos.items() if isinstance(v, (int, float)) and v > 0],
            columns=["Concepto", "Monto (ARS)"],
        ).sort_values("Monto (ARS)", ascending=False)

        if not df_gastos.empty:
            fig_g = px.pie(df_gastos, names="Concepto", values="Monto (ARS)", hole=0.4,
                title="Distribución de gastos fijos", color_discrete_sequence=COLORES)
            st.plotly_chart(fig_g, use_container_width=True)

        if st.session_state.df_tn is not None and not st.session_state.df_tn.empty:
            st.divider()
            dias_p = max((fecha_hasta - fecha_desde).days + 1, 1)
            factor = dias_p / 30
            st.metric(f"📐 Gastos para {dias_p} días ({factor:.2f}x mes)", fmt(round(total_gastos * factor)))

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 7: COSTOS DE CONSOLAS
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "💻 Costos de consolas":
        st.subheader("💻 Costos de consolas")
        st.caption("FOB + importación. Los productos de TN se agregan sin sobreescribir datos existentes.")

        if "costos_consolas" not in st.session_state:
            saved = gs_read("CostosConsolas")
            st.session_state.costos_consolas = saved if saved else FOB_DEFAULTS.copy()

        tc_consolas = int(dolar_blue) if dolar_blue else 1200

        @st.fragment
        def cargar_productos_tn():
            if st.button("🔄 Actualizar productos desde Tienda Nube", key="btn_load_consolas"):
                with st.spinner("Cargando..."):
                    productos_tn = get_tn_products()
                if productos_tn:
                    prods_map = {}
                    prods_urls = {}
                    multi_variant_skipped = 0
                    for p in productos_tn:
                        nombre_raw = p.get("name", {})
                        nombre = nombre_raw.get("es", "") if isinstance(nombre_raw, dict) else str(nombre_raw)
                        variants_list = p.get("variants", []) or []
                        # Productos con multiples variantes: no agregar la fila padre
                        # — el usuario gestiona cada variante con nombre propio en la sheet.
                        if len(variants_list) > 1:
                            multi_variant_skipped += 1
                            continue
                        peso_kg = None
                        for v in variants_list:
                            w = v.get("weight")
                            if w:
                                try:
                                    peso_kg = float(w)
                                    break
                                except Exception:
                                    pass
                        if not peso_kg:
                            try:
                                peso_kg = float(p.get("weight")) if p.get("weight") else None
                            except Exception:
                                peso_kg = None
                        prods_map[nombre] = peso_kg
                        handle_raw = p.get("handle") or {}
                        if isinstance(handle_raw, dict):
                            handle = handle_raw.get("es") or handle_raw.get("pt") or handle_raw.get("en") or ""
                        else:
                            handle = str(handle_raw or "")
                        url = ""
                        if handle:
                            url = f"https://www.marketgamer.com.ar/productos/{handle}/"
                        else:
                            permalink = p.get("permalink") or p.get("canonical_url") or ""
                            if permalink:
                                url = str(permalink)
                            elif p.get("id"):
                                url = f"https://marketgamer.mitiendanube.com/admin/v2/products/{p.get('id')}"
                        if url:
                            prods_urls[nombre] = url
                    st.session_state.productos_tn_map = prods_map
                    st.session_state.productos_tn_urls = prods_urls

                    costos_actual = st.session_state.costos_consolas.copy()
                    nuevos = 0
                    for nombre, peso in prods_map.items():
                        if nombre not in costos_actual:
                            fob_def = FOB_DEFAULTS.get(nombre, {}).get("fob_usd", 0.0)
                            peso_def = peso or FOB_DEFAULTS.get(nombre, {}).get("peso_kg", 0.0)
                            costos_actual[nombre] = {
                                "fob_usd": fob_def,
                                "peso_kg": peso_def,
                            }
                            nuevos += 1
                        else:
                            existing = costos_actual[nombre]
                            if isinstance(existing, dict) and peso:
                                if not existing.get("peso_kg"):
                                    existing["peso_kg"] = peso
                    st.session_state.costos_consolas = costos_actual
                    st.session_state._costos_needs_refresh = True
                    msg = f"✅ {len(prods_map)} productos cargados ({nuevos} nuevos agregados, existentes conservados)"
                    if multi_variant_skipped:
                        msg += f" · {multi_variant_skipped} con múltiples variantes omitidos (cargá cada variante a mano)"
                    st.success(msg)
                    try:
                        st.rerun(scope="app")
                    except TypeError:
                        st.rerun()
                else:
                    st.warning("No se pudieron cargar productos.")

        cargar_productos_tn()

        costos = st.session_state.costos_consolas.copy()
        productos_map = st.session_state.get("productos_tn_map", {})

        with st.expander(f"⚙️ Costos de referencia — importación ${costos.get('_costo_kg_usd', 65.0):.2f} USD/kg · dólar ${tc_consolas:,.0f} ARS", expanded=False):
            col_imp1, col_imp2 = st.columns(2)
            costo_kg_usd = col_imp1.number_input(
                "Costo importación (USD/kg)",
                value=float(costos.get("_costo_kg_usd", 65.0)), step=0.5, key="ckg",
            )
            col_imp2.metric("Dólar blue", f"${tc_consolas:,.0f} ARS")

        # ── Construir DF editable UNA sola vez en session_state ──
        def _build_costos_df():
            _costos = st.session_state.costos_consolas.copy()
            _prods_map = st.session_state.get("productos_tn_map", {})
            _prods_urls = st.session_state.get("productos_tn_urls", {})
            _all = set(_prods_map.keys())
            for k in _costos:
                if not k.startswith("_"):
                    _all.add(k)
            # Construir filas y deduplicar por nombre compacto
            # Si dos productos normalizan igual, queda el que tiene FOB > 0
            _seen_compact = {}  # norm_compact → (prod, peso, fob)
            for prod in sorted(_all):
                pd_ = _costos.get(prod, {})
                fob_s = float(pd_.get("fob_usd", 0.0) or 0.0) if isinstance(pd_, dict) else 0.0
                peso_s = float(pd_.get("peso_kg", 0.0) or 0.0) if isinstance(pd_, dict) else 0.0
                peso_tn = _prods_map.get(prod)
                peso_def = FOB_DEFAULTS.get(prod, {}).get("peso_kg", 0.0)
                peso_f = float(peso_tn or peso_s or peso_def)
                fob_def = FOB_DEFAULTS.get(prod, {}).get("fob_usd", 0.0)
                fob_f = fob_s if fob_s > 0 else float(fob_def)

                url_f = _prods_urls.get(prod, "")
                row_data = {"Producto": prod, "Peso (kg)": peso_f, "FOB (USD)": fob_f, "URL": url_f}
                nc = _norm_compact(prod)
                if nc in _seen_compact:
                    # Ya existe: quedarse con el que tenga mejor FOB
                    _prev = _seen_compact[nc]
                    if fob_f > _prev["FOB (USD)"]:
                        _seen_compact[nc] = row_data
                    elif fob_f == _prev["FOB (USD)"] and peso_f > _prev["Peso (kg)"]:
                        _seen_compact[nc] = row_data
                    elif url_f and not _prev.get("URL"):
                        _seen_compact[nc] = row_data
                else:
                    _seen_compact[nc] = row_data

            rows = sorted(_seen_compact.values(), key=lambda r: r["Producto"])
            return pd.DataFrame(rows) if rows else pd.DataFrame(columns=["Producto", "Peso (kg)", "FOB (USD)", "URL"])

        if "costos_df_editor" not in st.session_state:
            st.session_state.costos_df_editor = _build_costos_df()

        # Botón para refrescar desde datos guardados (tras cargar TN)
        if st.session_state.get("_costos_needs_refresh"):
            st.session_state.costos_df_editor = _build_costos_df()
            del st.session_state["_costos_needs_refresh"]

        if st.session_state.costos_df_editor.empty:
            st.info("Cargá productos desde TN con el botón de arriba.")
        else:
            # ── Controles de vista ──
            n_sin_precio = (st.session_state.costos_df_editor["FOB (USD)"] == 0).sum()
            ctrl1, ctrl2 = st.columns([3, 1])
            with ctrl1:
                busqueda = st.text_input(
                    "Buscar producto",
                    placeholder="ej: anbernic rg ds",
                    key="costos_busqueda",
                    label_visibility="collapsed",
                )
            with ctrl2:
                solo_sin_precio = st.toggle(
                    f"Sin precio ({n_sin_precio})",
                    value=False,
                    key="costos_filtro_sin_precio",
                )

            # ── Preparar DF para el editor ──
            df_edit_base = st.session_state.costos_df_editor.copy().sort_values("Producto").reset_index(drop=True)
            if busqueda.strip():
                mask = df_edit_base["Producto"].str.contains(busqueda.strip(), case=False, na=False)
                df_edit_base = df_edit_base[mask].copy()
            if solo_sin_precio:
                df_edit_base = df_edit_base[df_edit_base["FOB (USD)"] == 0].copy()

            # Agregar columnas calculadas como vista previa (read-only en el editor)
            df_edit_base["Import (USD)"] = (df_edit_base["Peso (kg)"] * costo_kg_usd).round(2)
            df_edit_base["Total (USD)"] = (df_edit_base["FOB (USD)"] + df_edit_base["Import (USD)"]).round(2)
            df_edit_base["Total (ARS)"] = (df_edit_base["Total (USD)"] * tc_consolas).round(0)

            n_sin_precio = (st.session_state.costos_df_editor["FOB (USD)"] == 0).sum()
            if n_sin_precio > 0:
                st.caption(f"⚠️ {n_sin_precio} producto(s) sin FOB cargado — usá el filtro para encontrarlos.")

            edited_df = st.data_editor(
                df_edit_base,
                column_config={
                    "Producto": st.column_config.TextColumn("Producto", width="large"),
                    "Peso (kg)": st.column_config.NumberColumn("Peso (kg)", min_value=0.0, step=0.01, format="%.3f"),
                    "FOB (USD)": st.column_config.NumberColumn("FOB (USD)", min_value=0.0, step=0.5, format="$%.2f"),
                    "Import (USD)": st.column_config.NumberColumn("Import (USD)", disabled=True, format="$%.2f"),
                    "Total (USD)": st.column_config.NumberColumn("Total (USD)", disabled=True, format="$%.2f"),
                    "Total (ARS)": st.column_config.NumberColumn("Total (ARS)", disabled=True, format="$%.0f"),
                    "URL": st.column_config.LinkColumn("Ver en TN", display_text="🔗 Abrir", disabled=True, help="Abre el producto en la tienda. Sin link = no encontrado en TN (posiblemente eliminado)."),
                },
                hide_index=True,
                use_container_width=True,
                num_rows="dynamic",
                height=600,
                key="costos_editor",
            )

            # NO sobreescribir costos_df_editor acá — el widget gestiona sus edits
            # internamente via key="costos_editor". Sobreescribirlo resetea el editor.
            if edited_df is not None:
                edited_df = edited_df.fillna({"Producto": "", "Peso (kg)": 0.0, "FOB (USD)": 0.0})
                edited_df["Import (USD)"] = (edited_df["Peso (kg)"] * costo_kg_usd).round(2)
                edited_df["Total (USD)"] = (edited_df["FOB (USD)"] + edited_df["Import (USD)"]).round(2)
                edited_df["Total (ARS)"] = (edited_df["Total (USD)"] * tc_consolas).round(0)

            if st.button("💾 Guardar costos de consolas", use_container_width=True, type="primary"):
                nuevos_costos = {k: v for k, v in st.session_state.costos_consolas.items()}
                nuevos_costos["_costo_kg_usd"] = costo_kg_usd

                # Productos visibles en el editor (puede ser un subconjunto si hay filtro)
                _prods_visibles = set(df_edit_base["Producto"].tolist())
                _editor_prods = set()

                for _, row in edited_df.iterrows():
                    name = str(row.get("Producto", "")).strip()
                    if not name:
                        continue
                    _editor_prods.add(name)
                    nuevos_costos[name] = {
                        "fob_usd": float(row["FOB (USD)"]),
                        "peso_kg": float(row["Peso (kg)"]),
                        "costo_import_usd": float(row["Import (USD)"]),
                        "costo_total_usd": float(row["Total (USD)"]),
                    }

                # Solo eliminar productos que estaban visibles y el usuario borró explícitamente
                # Los productos ocultos por búsqueda/filtro NO se tocan
                for p in _prods_visibles - _editor_prods:
                    nuevos_costos.pop(p, None)

                st.session_state.costos_consolas = nuevos_costos
                st.session_state._costos_needs_refresh = True
                ok = gs_write("CostosConsolas", nuevos_costos)
                st.success("✅ Guardado en Google Sheets" if ok else "⚠️ Solo en sesión")
                st.rerun()

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 8: MARGEN TEÓRICO POR CONSOLA
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "📐 Margen teórico":
        st.subheader("📐 Margen teórico por consola")

        _costos_gs_mt = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
        _tc_mt = int(dolar_blue) if dolar_blue else 1200

        # ── Configuración de costos adicionales (colapsado por default) ──
        with st.expander("⚙️ Costos adicionales", expanded=False):
            cc1, cc2, cc3 = st.columns(3)
            with cc1:
                iva_mt = cc1.number_input("IVA (%)", value=10.5, step=0.5, key="iva_mt")
            with cc2:
                packaging_ars = cc2.number_input("Packaging por unidad ($)", value=2500, step=500, key="pkg_mt")
            with cc3:
                cc3.metric("Dólar blue", f"${_tc_mt:,.0f}")

        # ── Cargar catálogo completo de TN (todos los productos publicados) ──
        @st.cache_data(ttl=300, show_spinner=False)
        def _get_catalogo_tn():
            prods = get_tn_products()
            catalogo = {}
            for p in prods:
                nombre_raw = p.get("name", {})
                nombre = nombre_raw.get("es", "") if isinstance(nombre_raw, dict) else str(nombre_raw)
                nombre = _extraer_nombre_producto(nombre)
                variantes = p.get("variants", [])
                if variantes:
                    # Precio de la primera variante activa
                    precio = float(variantes[0].get("price", 0) or 0)
                else:
                    precio = 0.0
                if nombre and precio > 0:
                    catalogo[nombre] = precio
            return catalogo

        catalogo_tn = _get_catalogo_tn()

        # ── Combinar productos de órdenes + catálogo TN ──
        orders_raw = st.session_state.orders_raw
        product_rows = _build_product_rows_from_raw(orders_raw)
        df_prod_raw = pd.DataFrame(product_rows) if product_rows else pd.DataFrame()

        # Precios de órdenes (promedio real vendido)
        if not df_prod_raw.empty:
            precios_ordenes = df_prod_raw.groupby("Producto").agg(
                Precio_prom=("Precio ($)", "mean"),
                Unidades=("Precio ($)", "count"),
            ).reset_index()
            precios_dict = dict(zip(precios_ordenes["Producto"], precios_ordenes["Precio_prom"]))
            unidades_dict = dict(zip(precios_ordenes["Producto"], precios_ordenes["Unidades"]))
        else:
            precios_dict = {}
            unidades_dict = {}

        # Combinar: empezamos con catálogo TN, completamos con datos de órdenes
        todos_productos = {}
        for nombre, precio_tn in catalogo_tn.items():
            todos_productos[nombre] = {
                "precio": precios_dict.get(nombre, precio_tn),
                "unidades": unidades_dict.get(nombre, 0),
                "fuente": "vendido" if nombre in precios_dict else "catálogo",
            }
        # Agregar los de órdenes que no están en catálogo
        for nombre, precio in precios_dict.items():
            if nombre not in todos_productos:
                todos_productos[nombre] = {
                    "precio": precio,
                    "unidades": unidades_dict.get(nombre, 0),
                    "fuente": "vendido",
                }

        if not todos_productos:
            st.info("No hay productos disponibles. Verificá que TN esté conectado.")
        else:
            # ── Inferir marca de cada producto ──
            def _inferir_marca(nombre):
                nc = _norm_compact(nombre)
                for key, marca in BRAND_CATALOG.items():
                    if key in nc:
                        return marca
                return "Otra"

            # ── Métricas del período ──
            if not df_tn.empty:
                envio_prom = df_tn["Envio costo ($)"].mean()
                total_fact = df_tn["Total ($)"].sum()
                total_com = df_tn["Comision PN ($)"].sum()
                tasa_ponderada = total_com / total_fact if total_fact > 0 else 0.0415
            else:
                envio_prom = 5000.0
                tasa_ponderada = 0.0415

            m1, m2, m3, m4 = st.columns(4)
            m1.metric("Tasa PN ponderada", f"{tasa_ponderada*100:.2f}%")
            m2.metric("Envío promedio", fmt(envio_prom))
            m3.metric("Packaging/u", fmt(packaging_ars))
            m4.metric("IVA", f"{iva_mt}%")

            # ── Construir tabla completa ──
            rows_mt = []
            for prod, datos in todos_productos.items():
                precio_prom = datos["precio"]
                unidades = datos["unidades"]
                fuente = datos["fuente"]
                marca = _inferir_marca(prod)

                costo_total_usd = get_costo_total_usd(prod, _costos_gs_mt)
                costo_total_ars = costo_total_usd * _tc_mt
                costo_iva = round(precio_prom * (iva_mt / 100), 0)
                costo_full = round(costo_total_ars + packaging_ars + costo_iva, 0)
                comision_est = round(precio_prom * tasa_ponderada, 0)
                margen_teorico = round(precio_prom - costo_full - comision_est - envio_prom, 0)
                margen_pct = round(margen_teorico / precio_prom * 100, 1) if precio_prom > 0 else 0

                rows_mt.append({
                    "Marca": marca,
                    "Producto": prod,
                    "Precio ($)": round(precio_prom, 0),
                    "Costo prod ($)": round(costo_total_ars, 0),
                    "Costo full ($)": costo_full,
                    "Comisión PN ($)": comision_est,
                    "Envío ($)": round(envio_prom, 0),
                    "Margen ($)": margen_teorico,
                    "Margen (%)": margen_pct,
                    "Uds vendidas": unidades,
                    "_fuente": fuente,
                })

            df_mt = pd.DataFrame(rows_mt).sort_values("Margen (%)", ascending=False)

            # ── Filtro por marca (pills clickeables) ──
            marcas_disponibles = sorted(df_mt["Marca"].unique().tolist())
            marcas_sel = st.pills(
                "Marca",
                options=marcas_disponibles,
                default=marcas_disponibles,
                selection_mode="multi",
                key="filtro_marcas_mt",
            )
            if marcas_sel:
                df_mt = df_mt[df_mt["Marca"].isin(marcas_sel)]

            st.caption(
                f"Mostrando **{len(df_mt)}** productos "
                f"({df_mt[df_mt['_fuente']=='vendido'].shape[0]} vendidos en el período, "
                f"{df_mt[df_mt['_fuente']=='catálogo'].shape[0]} solo en catálogo TN). "
                "Precio de catálogo TN para productos no vendidos en el período."
            )

            # ── Tabla principal ──
            df_display = df_mt.drop(columns=["_fuente"]).copy()

            def _color_row_margen(row):
                v = row["Margen (%)"]
                if v >= 20:
                    color = "background-color: #1e4620; color: #4ade80"
                elif v >= 10:
                    color = "background-color: #3a3000; color: #fbbf24"
                elif v >= 0:
                    color = "background-color: #2a1a00; color: #fb923c"
                else:
                    color = "background-color: #3a0f0f; color: #f87171"
                return [color if c == "Margen (%)" else "" for c in row.index]

            st.dataframe(
                df_display.style
                    .format({
                        "Precio ($)": "${:,.0f}",
                        "Costo prod ($)": "${:,.0f}",
                        "Costo full ($)": "${:,.0f}",
                        "Comisión PN ($)": "${:,.0f}",
                        "Envío ($)": "${:,.0f}",
                        "Margen ($)": "${:,.0f}",
                        "Margen (%)": "{:.1f}%",
                    })
                    .apply(_color_row_margen, axis=1),
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Marca": st.column_config.TextColumn("Marca", width="small"),
                    "Producto": st.column_config.TextColumn("Producto", width="large"),
                    "Uds vendidas": st.column_config.NumberColumn("Uds", width="small"),
                    "Margen (%)": st.column_config.TextColumn("Margen %", width="small"),
                },
            )

            # ══════════════════════════════════════════════════════════════
            # TABLA DE RENTABILIDAD POR CUOTAS
            # ══════════════════════════════════════════════════════════════
            st.subheader("💳 Rentabilidad por cantidad de cuotas")
            st.caption("Cómo cambia el margen según la forma de pago del cliente.")

            cuotas_config = {
                "Débito/Trans": 0.0415,
                "1 cuota": 0.0415,
                "3 cuotas": 0.1420,
                "6 cuotas": 0.2370,
                "12 cuotas": 0.4324,
            }

            with st.expander("⚙️ Ajustar tasas por cuota", expanded=False):
                tc_cols = st.columns(len(cuotas_config))
                cuotas_ajustadas = {}
                for i, (label, default) in enumerate(cuotas_config.items()):
                    with tc_cols[i]:
                        cuotas_ajustadas[label] = st.number_input(
                            f"{label} (%)", value=round(default * 100, 2),
                            step=0.1, key=f"cuota_mt_{label}",
                        ) / 100

            rows_cuotas = []
            for prod, datos in todos_productos.items():
                precio = datos["precio"]
                if precio <= 0:
                    continue

                costo_total_usd = get_costo_total_usd(prod, _costos_gs_mt)
                costo_total_ars = costo_total_usd * _tc_mt
                costo_iva = round(precio * (iva_mt / 100), 0)
                costo_full = round(costo_total_ars + packaging_ars + costo_iva, 0)

                row_data = {
                    "Producto": prod,
                    "Precio ($)": round(precio, 0),
                    "Costo full ($)": costo_full,
                }

                for label, tasa in cuotas_ajustadas.items():
                    comision = round(precio * tasa, 0)
                    margen = round(precio - costo_full - comision - envio_prom, 0)
                    margen_pct = round(margen / precio * 100, 1) if precio > 0 else 0
                    row_data[f"M {label} ($)"] = margen
                    row_data[f"M {label} (%)"] = margen_pct

                rows_cuotas.append(row_data)

            df_cuotas = pd.DataFrame(rows_cuotas).sort_values("Precio ($)", ascending=False)

            # Aplicar filtro de marca a cuotas también
            if marcas_sel:
                df_cuotas_filt = df_cuotas[df_cuotas["Producto"].apply(
                    lambda p: _inferir_marca(p) in marcas_sel
                )]
            else:
                df_cuotas_filt = df_cuotas

            fmt_cuotas = {
                "Precio ($)": "${:,.0f}",
                "Costo full ($)": "${:,.0f}",
            }
            pct_cols = []
            for label in cuotas_ajustadas.keys():
                fmt_cuotas[f"M {label} ($)"] = "${:,.0f}"
                fmt_cuotas[f"M {label} (%)"] = "{:.1f}%"
                pct_cols.append(f"M {label} (%)")

            def _color_margen_cuotas(v):
                if isinstance(v, (int, float)):
                    if v >= 20:
                        return "background-color: #1e4620; color: #4ade80"
                    elif v >= 10:
                        return "background-color: #3a3000; color: #fbbf24"
                    elif v >= 0:
                        return "background-color: #2a1a00; color: #fb923c"
                    else:
                        return "background-color: #3a0f0f; color: #f87171"
                return ""

            st.dataframe(
                df_cuotas_filt.style
                    .format(fmt_cuotas)
                    .map(_color_margen_cuotas, subset=pct_cols),
                use_container_width=True, hide_index=True,
            )

            # Gráfico comparativo para un producto seleccionado
            prod_sel_cuotas = st.selectbox(
                "Ver detalle por cuotas de:",
                df_cuotas_filt["Producto"].tolist(),
                key="sel_cuotas_mt",
            )
            if prod_sel_cuotas:
                row_sel = df_cuotas_filt[df_cuotas_filt["Producto"] == prod_sel_cuotas].iloc[0]
                chart_data = []
                for label in cuotas_ajustadas.keys():
                    chart_data.append({
                        "Cuotas": label,
                        "Margen ($)": row_sel[f"M {label} ($)"],
                        "Margen (%)": row_sel[f"M {label} (%)"],
                        "Tasa PN (%)": round(cuotas_ajustadas[label] * 100, 2),
                    })
                df_chart_cuotas = pd.DataFrame(chart_data)

                fig_cuotas = go.Figure()
                colors_cuotas = [
                    "#00C49F" if m >= 0 else "#FF5733"
                    for m in df_chart_cuotas["Margen ($)"]
                ]
                fig_cuotas.add_trace(go.Bar(
                    x=df_chart_cuotas["Cuotas"],
                    y=df_chart_cuotas["Margen ($)"],
                    name="Margen ($)",
                    marker_color=colors_cuotas,
                    text=df_chart_cuotas["Margen ($)"].apply(lambda x: f"${x:,.0f}"),
                    textposition="outside",
                ))
                fig_cuotas.add_trace(go.Scatter(
                    x=df_chart_cuotas["Cuotas"],
                    y=df_chart_cuotas["Margen (%)"],
                    name="Margen (%)",
                    mode="lines+markers+text",
                    line=dict(color="#009EE3", width=3),
                    marker=dict(size=10),
                    text=df_chart_cuotas["Margen (%)"].apply(lambda x: f"{x:.1f}%"),
                    textposition="top center",
                    yaxis="y2",
                ))
                fig_cuotas.update_layout(
                    title=f"Margen por cuotas — {prod_sel_cuotas} (Precio: {fmt(row_sel['Precio ($)'])})",
                    yaxis=dict(title="Margen ($)", tickformat="$,.0f"),
                    yaxis2=dict(title="Margen (%)", overlaying="y", side="right", ticksuffix="%", showgrid=False),
                    legend=dict(orientation="h", y=-0.15),
                    hovermode="x unified",
                )
                st.plotly_chart(fig_cuotas, use_container_width=True)

            col_dl1, col_dl2 = st.columns(2)
            with col_dl1:
                st.download_button("⬇️ Descargar margen teórico",
                    df_mt.to_csv(index=False).encode("utf-8"), "margen_teorico.csv", "text/csv")
            with col_dl2:
                st.download_button("⬇️ Descargar rentabilidad por cuotas",
                    df_cuotas.to_csv(index=False).encode("utf-8"), "rentabilidad_cuotas.csv", "text/csv")

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 9: MARGEN REAL POR CONSOLA
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "📈 Margen real":
        st.subheader("📈 Margen real por consola y medio de pago")
        st.caption(
            "Margen real calculado orden por orden usando el **precio individual** de cada producto. "
            "Muestra cómo el medio de pago impacta en la rentabilidad."
        )

        if df_tn.empty:
            st.info("Buscá primero para ver los datos.")
        else:
            _costos_gs_mr = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
            _tc_mr = int(dolar_blue) if dolar_blue else 1200
            orders_raw_mr = st.session_state.orders_raw

            # Config costos adicionales
            with st.expander("⚙️ Costos adicionales", expanded=False):
                cr1, cr2 = st.columns(2)
                iva_mr = cr1.number_input("IVA (%)", value=10.5, step=0.5, key="iva_mr")
                packaging_mr = cr2.number_input("Packaging/u ($)", value=2500, step=500, key="pkg_mr")

            # Extraer precios individuales desde raw orders
            product_rows_mr = _build_product_rows_from_raw(orders_raw_mr)

            if not product_rows_mr:
                st.info("No hay datos de productos.")
            else:
                # Calcular margen real por cada línea de producto
                rows_real = []
                for pr in product_rows_mr:
                    prod = pr["Producto"]
                    precio_lista = pr["Precio ($)"]
                    precio_neto = pr.get("Precio neto ($)", precio_lista)
                    descuento = pr.get("Descuento ($)", 0)
                    comision = pr["Comisión PN ($)"]
                    envio = pr["Envío ($)"]

                    _fob_usd, _import_usd, costo_total_usd = _match_costo_entry(prod, _costos_gs_mr)
                    costo_total_ars = costo_total_usd * _tc_mr
                    # IVA se calcula sobre el precio neto (lo que efectivamente se factura)
                    costo_iva = round(precio_neto * (iva_mr / 100), 0)
                    costo_full = round(costo_total_ars + packaging_mr + costo_iva, 0)

                    neto = precio_neto - comision
                    margen = round(neto - costo_full - envio, 0)

                    rows_real.append({
                        "Producto": prod,
                        "Medio de Pago": pr["Medio de Pago"],
                        "Cuotas": pr["Cuotas"],
                        "Precio lista ($)": round(precio_lista, 0),
                        "Descuento ($)": round(descuento, 0),
                        "Precio neto ($)": round(precio_neto, 0),
                        "Comisión PN ($)": round(comision, 0),
                        "Costo PN (%)": round(comision / precio_neto * 100, 2) if precio_neto > 0 else 0,
                        "FOB (USD)": round(_fob_usd, 2),
                        "Import (USD)": round(_import_usd, 2),
                        "← Costo prod ($)": round(costo_total_ars, 0),
                        "Packaging ($)": packaging_mr,
                        f"IVA ({iva_mr}%)": costo_iva,
                        "Envío ($)": round(envio, 0),
                        "Margen ($)": margen,
                        "Margen (%)": round(margen / precio_neto * 100, 1) if precio_neto > 0 else 0,
                    })

                df_real = pd.DataFrame(rows_real)

                if df_real.empty:
                    st.info("No hay datos suficientes.")
                else:
                    # ── Vista 1: Margen promedio por producto ──
                    st.markdown("### Margen promedio real por consola")
                    col_iva_mr = f"IVA ({iva_mr}%)"
                    df_avg = df_real.groupby("Producto").agg(
                        Ventas=("Precio neto ($)", "count"),
                        Precio_lista_prom=("Precio lista ($)", "mean"),
                        Descuento_prom=("Descuento ($)", "mean"),
                        Precio_prom=("Precio neto ($)", "mean"),
                        Comision_prom=("Comisión PN ($)", "mean"),
                        FOB_prom=("FOB (USD)", "mean"),
                        Import_prom=("Import (USD)", "mean"),
                        Costo_prom=("← Costo prod ($)", "mean"),
                        Packaging_prom=("Packaging ($)", "mean"),
                        IVA_prom=(col_iva_mr, "mean"),
                        Envio_prom=("Envío ($)", "mean"),
                        Margen_prom=("Margen ($)", "mean"),
                        Margen_pct_prom=("Margen (%)", "mean"),
                    ).reset_index()
                    df_avg.columns = [
                        "Producto", "Ventas", "Precio lista prom ($)", "Descuento prom ($)",
                        "Precio neto prom ($)", "Comisión prom ($)",
                        "FOB (USD)", "Import (USD)", "← Costo prod ($)",
                        "Packaging ($)", col_iva_mr, "Envío ($)",
                        "Margen prom ($)", "Margen (%)",
                    ]
                    df_avg = df_avg.sort_values("Margen (%)", ascending=False)

                    # Gráfico
                    df_avg_chart = df_avg.sort_values("Margen (%)", ascending=True).tail(20)
                    fig_mr = px.bar(
                        df_avg_chart, x="Margen (%)", y="Producto", orientation="h",
                        title="Margen real promedio por consola (%)",
                        color="Margen (%)",
                        color_continuous_scale=["#FF5733", "#FFD700", "#00C49F"],
                        text="Margen (%)",
                    )
                    fig_mr.update_layout(
                        yaxis={"categoryorder": "total ascending"},
                        coloraxis_showscale=False,
                    )
                    fig_mr.update_traces(texttemplate="%{text:.1f}%", textposition="outside")
                    st.plotly_chart(fig_mr, use_container_width=True)

                    # ── Margen promedio ponderado del total ──
                    _total_revenue = df_avg["Precio neto prom ($)"].mul(df_avg["Ventas"]).sum()
                    _total_margen = df_avg["Margen prom ($)"].mul(df_avg["Ventas"]).sum()
                    _total_ventas = df_avg["Ventas"].sum()
                    _margen_pct_pond = (_total_margen / _total_revenue * 100) if _total_revenue > 0 else 0
                    _margen_prom_unit = _total_margen / _total_ventas if _total_ventas > 0 else 0
                    _color_mg = "#00C49F" if _margen_pct_pond >= 20 else "#FFD700" if _margen_pct_pond >= 10 else "#FF5733"

                    _mc1, _mc2, _mc3, _mc4 = st.columns(4)
                    _mc1.metric("📦 Unidades vendidas", f"{int(_total_ventas)}")
                    _mc2.metric("💰 Revenue total", f"${_total_revenue:,.0f}")
                    _mc3.metric("📊 Margen total", f"${_total_margen:,.0f}")
                    _mc4.metric("📈 Margen promedio pond.", f"{_margen_pct_pond:.1f}%")
                    st.caption(f"Margen promedio por unidad: **${_margen_prom_unit:,.0f}** · Comisiones PN son reales (calculadas con la tasa según medio de pago y cuotas de cada orden)")

                    st.dataframe(
                        df_avg.style.format({
                            "Precio lista prom ($)": "${:,.0f}",
                            "Descuento prom ($)": "${:,.0f}",
                            "Precio neto prom ($)": "${:,.0f}",
                            "Comisión prom ($)": "${:,.0f}",
                            "FOB (USD)": "${:,.2f}",
                            "Import (USD)": "${:,.2f}",
                            "← Costo prod ($)": "${:,.0f}",
                            "Packaging ($)": "${:,.0f}",
                            col_iva_mr: "${:,.0f}",
                            "Envío ($)": "${:,.0f}",
                            "Margen prom ($)": "${:,.0f}",
                            "Margen (%)": "{:.1f}%",
                        }).map(lambda v: (
                            "background-color: #1e4620; color: #00C49F" if isinstance(v, (int, float)) and v >= 20
                            else "background-color: #5a4a1a; color: #ffd700" if isinstance(v, (int, float)) and v >= 10
                            else "background-color: #4a1010; color: #ff6b6b" if isinstance(v, (int, float))
                            else ""
                        ), subset=["Margen (%)"]),
                        use_container_width=True, hide_index=True,
                    )

                    # ── Vista 2: Impacto del medio de pago por consola ──
                    st.divider()
                    st.markdown("### Impacto del medio de pago en el margen")

                    prod_sel_mr = st.selectbox(
                        "Seleccioná un producto",
                        df_avg["Producto"].tolist(),
                        key="sel_margen_real",
                    )

                    if prod_sel_mr:
                        df_prod = df_real[df_real["Producto"] == prod_sel_mr]
                        df_by_medio = df_prod.groupby("Medio de Pago").agg(
                            Ventas=("Precio neto ($)", "count"),
                            Precio_lista_prom=("Precio lista ($)", "mean"),
                            Descuento_prom=("Descuento ($)", "mean"),
                            Precio_prom=("Precio neto ($)", "mean"),
                            Comision_prom=("Comisión PN ($)", "mean"),
                            Costo_PN_pct=("Costo PN (%)", "mean"),
                            Margen_prom=("Margen ($)", "mean"),
                            Margen_pct=("Margen (%)", "mean"),
                        ).reset_index().sort_values("Margen_pct", ascending=False)
                        df_by_medio.columns = [
                            "Medio de Pago", "Ventas", "Precio lista ($)", "Descuento ($)",
                            "Precio neto ($)", "Comisión prom ($)", "Costo PN (%)",
                            "Margen prom ($)", "Margen (%)",
                        ]

                        col_mr1, col_mr2 = st.columns(2)
                        with col_mr1:
                            fig_medio = px.bar(
                                df_by_medio, x="Medio de Pago", y="Margen (%)",
                                title=f"Margen real por medio de pago — {prod_sel_mr}",
                                color="Margen (%)",
                                color_continuous_scale=["#FF5733", "#FFD700", "#00C49F"],
                                text="Margen (%)",
                            )
                            fig_medio.update_traces(texttemplate="%{text:.1f}%", textposition="outside")
                            fig_medio.update_layout(coloraxis_showscale=False)
                            st.plotly_chart(fig_medio, use_container_width=True)

                        with col_mr2:
                            fig_dist = px.pie(
                                df_by_medio, names="Medio de Pago", values="Ventas",
                                title=f"Distribución de ventas — {prod_sel_mr}",
                                color_discrete_sequence=COLORES, hole=0.35,
                            )
                            fig_dist.update_traces(textinfo="label+value+percent")
                            st.plotly_chart(fig_dist, use_container_width=True)

                        st.dataframe(
                            df_by_medio.style.format({
                                "Precio lista ($)": "${:,.0f}",
                                "Descuento ($)": "${:,.0f}",
                                "Precio neto ($)": "${:,.0f}",
                                "Comisión prom ($)": "${:,.0f}",
                                "Costo PN (%)": "{:.2f}%",
                                "Margen prom ($)": "${:,.0f}",
                                "Margen (%)": "{:.1f}%",
                            }),
                            use_container_width=True, hide_index=True,
                        )

                        # Diferencia entre mejor y peor medio
                        if len(df_by_medio) >= 2:
                            mejor = df_by_medio.iloc[0]
                            peor = df_by_medio.iloc[-1]
                            diff = round(mejor["Margen (%)"] - peor["Margen (%)"], 1)
                            st.info(
                                f"💡 Diferencia de margen entre **{mejor['Medio de Pago']}** "
                                f"({mejor['Margen (%)']:.1f}%) y **{peor['Medio de Pago']}** "
                                f"({peor['Margen (%)']:.1f}%): **{diff} puntos porcentuales**"
                            )

                    st.divider()
                    st.download_button("⬇️ Descargar margen real detallado",
                        df_real.to_csv(index=False).encode("utf-8"), "margen_real_detallado.csv", "text/csv")

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 10: PROVEEDORES
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🏭 Proveedores":
        st.markdown("""
        <style>
        /* ── Market Gamer · Proveedores ───────────────────────────────── */
        .prov-page-header{display:flex;align-items:baseline;gap:14px;padding-bottom:20px;border-bottom:1px solid rgba(255,255,255,0.06);margin-bottom:8px}
        .prov-page-title{font-size:22px;font-weight:700;color:#e2e8f0;letter-spacing:-.3px;margin:0}
        .prov-stat-pill{font-size:11px;font-weight:600;color:#6b7280;background:rgba(255,255,255,0.05);border:1px solid rgba(255,255,255,0.08);padding:3px 10px;border-radius:20px}
        .prov-section{font-size:10px;font-weight:700;color:#4b5563;letter-spacing:1.2px;text-transform:uppercase;padding:20px 0 10px;border-bottom:1px solid rgba(255,255,255,0.05);margin-bottom:14px}
        .sup-contact{font-size:12px;color:#6b7280;display:flex;gap:16px;margin-bottom:14px;flex-wrap:wrap}
        .sup-contact span{display:flex;align-items:center;gap:5px}
        .cat-tbl{width:100%;border-collapse:collapse;margin-bottom:16px}
        .cat-tbl th{font-size:10px;font-weight:700;color:#4b5563;text-transform:uppercase;letter-spacing:.7px;padding:6px 10px;border-bottom:1px solid rgba(255,255,255,0.07);text-align:left}
        .cat-tbl th.r{text-align:right}
        .cat-tbl td{font-size:13px;color:#94a3b8;padding:8px 10px;border-bottom:1px solid rgba(255,255,255,0.03)}
        .cat-tbl td.pn{color:#e2e8f0;font-weight:500}
        .cat-tbl td.pr{text-align:right;font-family:ui-monospace,monospace;font-size:14px;color:#cbd5e1;font-weight:600}
        .cat-tbl tr:hover td{background:rgba(255,255,255,0.02)}
        .brand-tag{display:inline-block;font-size:10px;font-weight:600;padding:2px 7px;border-radius:4px;background:rgba(255,255,255,0.07);color:#94a3b8}
        .brand-anbernic{background:rgba(96,165,250,0.12);color:#60a5fa}
        .brand-powkiddy{background:rgba(251,191,36,0.12);color:#fbbf24}
        .brand-trimui{background:rgba(74,222,128,0.12);color:#4ade80}
        .brand-miyoo{background:rgba(192,132,252,0.12);color:#c084fc}
        .brand-retroid{background:rgba(251,146,60,0.12);color:#fb923c}
        .comp-banner{background:rgba(74,222,128,0.05);border:1px solid rgba(74,222,128,0.15);border-radius:8px;padding:14px 18px;margin:8px 0 16px;display:flex;align-items:center;gap:12px}
        .comp-banner-title{font-size:15px;font-weight:700;color:#e2e8f0}
        .comp-banner-sub{font-size:12px;color:#6b7280;margin-top:2px}
        </style>
        """, unsafe_allow_html=True)

        # ── Cargar datos guardados ──
        if "proveedores_data" not in st.session_state:
            saved_prov = gs_read("Proveedores")
            if saved_prov and "suppliers" in saved_prov:
                st.session_state.proveedores_data = saved_prov
            else:
                st.session_state.proveedores_data = {
                    "suppliers": {
                        "Anne - Qbuy Technology": {
                            "contacto": "Anne Lee | WhatsApp: +86 181 7167 5976 | Wechat: qbuy18",
                            "web": "qbuytech.en.alibaba.com",
                            "updated_at": "",
                            "productos": {
                                "R36S": {"precio_usd": 23.2, "storage": "64G", "marca": "R36 Series", "screen": "3.5in 640x480", "cpu": "RK3326", "battery": "3200mAh"},
                                "R36H": {"precio_usd": 27.5, "storage": "64G", "marca": "R36 Series", "screen": "3.5in 640x480", "cpu": "RK3326", "battery": "3000mAh"},
                                "R36S PLUS": {"precio_usd": 28.4, "storage": "64G", "marca": "R36 Series", "screen": "3.5in 640x480", "cpu": "RK3326", "battery": "3200mAh"},
                                "RG35XX PRO": {"precio_usd": 45.7, "storage": "64G", "marca": "Anbernic", "screen": "3.5in 640x480", "cpu": "H700", "battery": "3300mAh"},
                                "RG406V": {"precio_usd": 161.7, "storage": "128G", "marca": "Anbernic", "screen": "4in 960x720", "cpu": "T820", "battery": "5500mAh"},
                                "RG406H": {"precio_usd": 158.3, "storage": "128G", "marca": "Anbernic", "screen": "4in 960x720", "cpu": "T820", "battery": "5000mAh"},
                                "RG556": {"precio_usd": 182.6, "storage": "128G", "marca": "Anbernic", "screen": "5.48in AMOLED", "cpu": "T820", "battery": "—"},
                                "RG477M": {"precio_usd": 288.7, "storage": "256G", "marca": "Anbernic", "screen": "4.7in 1280x960", "cpu": "T820", "battery": "5300mAh"},
                            },
                        }
                    }
                }

        prov_data = st.session_state.proveedores_data
        suppliers = prov_data.get("suppliers", {})

        _n_sups = len(suppliers)
        _n_prods = sum(len(s.get("productos", {})) for s in suppliers.values())
        st.markdown(
            f'<div class="prov-page-header">'
            f'<span class="prov-page-title">Proveedores</span>'
            f'<span class="prov-stat-pill">{_n_sups} activos</span>'
            f'<span class="prov-stat-pill">{_n_prods} productos</span>'
            f'</div>',
            unsafe_allow_html=True,
        )

        # ════════════════════════════════════════════════════════════════════
        # SECCIÓN 1: CATÁLOGOS
        # ════════════════════════════════════════════════════════════════════
        st.markdown('<div class="prov-section">Catálogos</div>', unsafe_allow_html=True)

        def _parse_uploaded_to_products(uploaded_file) -> dict:
            """Parsea un archivo subido (CSV/Excel/PDF/PPTX) y devuelve dict de productos."""
            ext = uploaded_file.name.lower()
            raw_rows = []

            if ext.endswith(".pptx"):
                raw_rows = parse_pptx_catalog(uploaded_file.read())

            elif ext.endswith(".pdf"):
                try:
                    import pdfplumber
                    with pdfplumber.open(uploaded_file) as pdf:
                        all_tables = []
                        for page in pdf.pages:
                            all_tables.extend(page.extract_tables())

                    # Auto-detectar formato: Anbernic (EXW SZ) vs Anne/Qbuy (Price(USD))
                    def _detect_pdf_format(tables):
                        for table in tables:
                            for row in table:
                                if not row:
                                    continue
                                cells = [str(c).strip().lower() if c else "" for c in row]
                                row_text = " ".join(cells)
                                if "exw" in row_text:
                                    return "anbernic"
                                if "price(usd)" in row_text or "speicification" in row_text or "specification" in row_text:
                                    return "anne"
                        return "anne"  # default

                    pdf_format = _detect_pdf_format(all_tables)

                    if pdf_format == "anbernic":
                        # Formato Anbernic: header solo en página 1, resto son filas de datos directas
                        # ITEM | Specification | Pictures | OEM (MOQ) | MOQ (pcs) | EXW SZ | Quotation Mode | Package
                        import re as _re
                        item_col, price_col = 0, 5  # defaults confirmados por el header de pag 1
                        # Buscar header para confirmar índices de columnas
                        for table in all_tables:
                            for row in table:
                                if not row:
                                    continue
                                cells_lower = [str(c).strip().lower() if c else "" for c in row]
                                if any("exw" in c for c in cells_lower):
                                    for i, c in enumerate(cells_lower):
                                        if c == "item":
                                            item_col = i
                                        if "exw" in c:
                                            price_col = i
                                    break
                            else:
                                continue
                            break

                        SKIP_VALUES = {"item", ""}
                        # Regex para variantes con label: "8+128GB: US$201.5" o "US$29.00"
                        VARIANT_RE = _re.compile(r'^(.+?):\s*(?:US)?\$([\d,]+\.?\d*)$')
                        PLAIN_RE = _re.compile(r'(?:US)?\$([\d,]+\.?\d*)')
                        for table in all_tables:
                            for row in table:
                                if not row:
                                    continue
                                cells = [str(c).strip() if c else "" for c in row]
                                model = cells[item_col] if item_col < len(cells) else ""
                                price_raw = cells[price_col] if price_col < len(cells) else ""
                                if not model or not price_raw:
                                    continue
                                model_clean = model.split("\n")[0].strip()
                                if model_clean.lower() in SKIP_VALUES:
                                    continue
                                # Separar por línea: cada línea puede ser una variante
                                for line in price_raw.split("\n"):
                                    line = line.strip()
                                    if not line:
                                        continue
                                    vm = VARIANT_RE.match(line)
                                    if vm:
                                        # "8+128GB: US$201.5" → producto "RG 477V 8+128GB"
                                        label = vm.group(1).strip()
                                        price_str = vm.group(2).replace(",", "")
                                        product_name = f"{model_clean} {label}"
                                    else:
                                        pm = PLAIN_RE.search(line)
                                        if not pm:
                                            continue
                                        price_str = pm.group(1).replace(",", "")
                                        product_name = model_clean
                                    try:
                                        val = float(price_str)
                                        if 1 < val < 5000:
                                            raw_rows.append({
                                                "Producto": product_name,
                                                "FOB (USD)": val,
                                                "Marca": "Anbernic", "Pantalla": "—", "CPU": "—", "Storage": "—",
                                            })
                                    except ValueError:
                                        pass

                    else:
                        # Formato Anne (Qbuy): cada producto ocupa 2 filas en la tabla:
                        #   Fila header: [model_name, "Speicification", "Color", "Price(USD)", ""]
                        #   Fila data:   [""/None,    specs_text,       colores, precio,       storage]
                        current_model = None
                        SPEC_KEYWORDS = {"speicification", "specification", "price(usd)", "color"}
                        for table in all_tables:
                            for row in table:
                                if not row:
                                    continue
                                cells = [str(c).strip() if c else "" for c in row]
                                col0 = cells[0] if cells else ""
                                col1 = cells[1].lower() if len(cells) > 1 else ""
                                col3 = cells[3] if len(cells) > 3 else ""
                                col4 = cells[4] if len(cells) > 4 else ""

                                if col0 and any(kw in col1 for kw in SPEC_KEYWORDS):
                                    current_model = col0.split("\n")[0].strip()
                                    continue

                                if current_model and not col0:
                                    try:
                                        val = float(col3.replace("$", "").replace(",", "").strip())
                                        if 1 < val < 5000:
                                            storage = col4.strip() if col4.strip() else "—"
                                            # Agregar sufijo de variante SOLO cuando corresponde:
                                            # 1. El spec lista múltiples configs: "8g+128g, 12g+256g" → 2 pares
                                            # 2. Modelos conocidos multi-variante con spec en formato separado
                                            import re as _re2
                                            # Modelos Anbernic que tienen 2 variantes de RAM+storage
                                            KNOWN_MULTI_VARIANT = {"RG477M", "RG 477M"}
                                            spec_lower = col1.lower()
                                            model_key = _re2.sub(r'\s+', '', current_model).upper()
                                            pairs = _re2.findall(r'(\d+)\s*gb?\s*\+\s*(\d+)\s*gb?', spec_lower)
                                            if len(pairs) > 1:
                                                # Múltiples variantes explícitas en el spec → tomar la mayor
                                                ram, stor = pairs[-1]
                                                product_name = f"{current_model} {ram}+{stor}GB"
                                            elif current_model in KNOWN_MULTI_VARIANT or model_key in {_re2.sub(r'\s+', '', m).upper() for m in KNOWN_MULTI_VARIANT}:
                                                # Modelo multi-variante con spec en formato "12GB...256G"
                                                ram_m = _re2.search(r'(\d+)\s*gb\b', spec_lower)
                                                stor_m = _re2.search(r'\b(\d{2,3})\s*g\b', spec_lower)
                                                if ram_m and stor_m:
                                                    product_name = f"{current_model} {ram_m.group(1)}+{stor_m.group(1)}GB"
                                                else:
                                                    product_name = current_model
                                            else:
                                                product_name = current_model
                                            raw_rows.append({
                                                "Producto": product_name,
                                                "FOB (USD)": val,
                                                "Marca": "—", "Pantalla": "—", "CPU": "—", "Storage": storage,
                                            })
                                            current_model = None
                                    except ValueError:
                                        pass

                except ImportError:
                    st.warning("⚠️ `pdfplumber` no instalado. Usá CSV/Excel en su lugar.")
                    return {}

            elif ext.endswith((".csv", ".tsv")):
                sep = "\t" if ext.endswith(".tsv") else ","
                df_up = pd.read_csv(uploaded_file, sep=sep)
                cols = df_up.columns.tolist()
                col_nombre = cols[0] if cols else None
                col_precio = cols[1] if len(cols) > 1 else None
                if col_nombre and col_precio:
                    for _, r in df_up.iterrows():
                        try:
                            precio = float(str(r[col_precio]).replace("$", "").replace(",", ""))
                        except Exception:
                            continue
                        if str(r[col_nombre]).strip() and precio > 0:
                            raw_rows.append({
                                "Producto": str(r[col_nombre]).strip(),
                                "FOB (USD)": precio,
                                "Marca": str(r.get("marca", r.get("Marca", "—"))).strip(),
                                "Pantalla": "—", "CPU": "—", "Storage": "—",
                            })

            elif ext.endswith((".xlsx", ".xls")):
                df_up = pd.read_excel(uploaded_file)
                cols = df_up.columns.tolist()
                col_nombre = cols[0] if cols else None
                col_precio = cols[1] if len(cols) > 1 else None
                if col_nombre and col_precio:
                    for _, r in df_up.iterrows():
                        try:
                            precio = float(str(r[col_precio]).replace("$", "").replace(",", ""))
                        except Exception:
                            continue
                        if str(r[col_nombre]).strip() and precio > 0:
                            raw_rows.append({
                                "Producto": str(r[col_nombre]).strip(),
                                "FOB (USD)": precio,
                                "Marca": str(r.get("marca", r.get("Marca", "—"))).strip(),
                                "Pantalla": "—", "CPU": "—", "Storage": "—",
                            })

            # Convertir raw_rows a dict de productos (deduplicar por nombre)
            products = {}
            for row in raw_rows:
                nombre = row.get("Producto", "").strip()
                precio = float(row.get("FOB (USD)", 0))
                if nombre and precio > 0 and nombre not in products:
                    products[nombre] = {
                        "precio_usd": precio,
                        "marca": row.get("Marca", "—"),
                        "screen": row.get("Pantalla", "—"),
                        "cpu": row.get("CPU", "—"),
                        "storage": row.get("Storage", "—"),
                        "battery": "—",
                    }
            return products

        # Mostrar cada proveedor con su catálogo y opción de reemplazo
        for sup_name in list(suppliers.keys()):
            sup_data = suppliers[sup_name]
            updated = sup_data.get("updated_at", "") or "Sin fecha"
            prod_count = len(sup_data.get("productos", {}))

            with st.expander(f"**{sup_name}** — {prod_count} productos · {updated}"):
                contacto = sup_data.get("contacto", "—")
                web = sup_data.get("web", "—")
                st.markdown(
                    f'<div class="sup-contact">'
                    f'<span>📞 {contacto}</span>'
                    f'<span>🌐 {web}</span>'
                    f'</div>',
                    unsafe_allow_html=True,
                )

                if sup_data.get("productos"):
                    _BRAND_CSS = {
                        "Anbernic": "brand-anbernic", "Powkiddy": "brand-powkiddy",
                        "Trimui": "brand-trimui", "Miyoo": "brand-miyoo", "Retroid": "brand-retroid",
                    }
                    productos_sorted = sorted(
                        sup_data["productos"].items(),
                        key=lambda x: float(x[1].get("precio_usd", 0)),
                    )
                    rows_html = ""
                    for _pname, _pinfo in productos_sorted:
                        _precio = float(_pinfo.get("precio_usd", 0))
                        _marca = _pinfo.get("marca", "—")
                        if _marca == "—":
                            _marca = _get_brand(_pname)
                        _storage = _pinfo.get("storage", "—")
                        _bcss = _BRAND_CSS.get(_marca, "")
                        _marca_html = f'<span class="brand-tag {_bcss}">{_marca}</span>' if _marca != "—" else '<span style="color:#374151">—</span>'
                        rows_html += (
                            f"<tr>"
                            f'<td class="pn">{_pname}</td>'
                            f"<td>{_marca_html}</td>"
                            f"<td>{_storage}</td>"
                            f'<td class="pr">${_precio:.1f}</td>'
                            f"</tr>"
                        )
                    st.markdown(
                        f'<table class="cat-tbl">'
                        f"<thead><tr>"
                        f"<th>Producto</th><th>Marca</th><th>Storage</th><th class='r'>FOB (USD)</th>"
                        f"</tr></thead>"
                        f"<tbody>{rows_html}</tbody>"
                        f"</table>",
                        unsafe_allow_html=True,
                    )

                st.markdown("**Actualizar catálogo** — reemplaza la lista anterior completa")
                uploaded = st.file_uploader(
                    f"Nueva lista de precios ({sup_name})",
                    type=["csv", "xlsx", "xls", "pdf", "pptx"],
                    key=f"upload_replace_{sup_name}",
                )

                if uploaded:
                    try:
                        new_products = _parse_uploaded_to_products(uploaded)
                    except Exception as e:
                        st.error(f"Error procesando archivo: {e}")
                        new_products = {}

                    if new_products:
                        st.success(f"✅ {len(new_products)} productos leídos del archivo")
                        old_products = sup_data.get("productos", {})
                        diff = compute_catalog_diff(old_products, new_products)

                        if diff:
                            st.markdown("**Cambios respecto al catálogo actual:**")
                            emoji_map = {"cambiado": "🟡", "nuevo": "🟢", "eliminado": "🔴"}
                            df_diff = pd.DataFrame(diff)
                            df_diff.insert(0, "", df_diff["tipo"].map(emoji_map))
                            df_diff = df_diff.drop(columns=["tipo"])
                            fmt_diff = {}
                            if "Antes (USD)" in df_diff.columns:
                                fmt_diff["Antes (USD)"] = lambda x: f"${x:.1f}" if x is not None and not pd.isna(x) else "—"
                            if "Después (USD)" in df_diff.columns:
                                fmt_diff["Después (USD)"] = lambda x: f"${x:.1f}" if x is not None and not pd.isna(x) else "—"
                            st.dataframe(df_diff, use_container_width=True, hide_index=True)
                        else:
                            st.info("Sin cambios de precio respecto al catálogo anterior.")

                        if st.button(
                            f"✅ Confirmar reemplazo — {sup_name}",
                            key=f"confirm_replace_{sup_name}",
                            type="primary",
                        ):
                            suppliers[sup_name]["productos"] = new_products
                            suppliers[sup_name]["updated_at"] = str(_date.today())
                            prov_data["suppliers"] = suppliers
                            st.session_state.proveedores_data = prov_data
                            gs_write("Proveedores", prov_data)
                            st.success(f"✅ Catálogo de {sup_name} actualizado ({len(new_products)} productos)")
                            st.rerun()
                    else:
                        st.warning("No se encontraron productos válidos en el archivo.")

                col_del1, col_del2 = st.columns([4, 1])
                with col_del2:
                    if st.button("🗑️ Eliminar proveedor", key=f"del_{sup_name}"):
                        del suppliers[sup_name]
                        prov_data["suppliers"] = suppliers
                        st.session_state.proveedores_data = prov_data
                        gs_write("Proveedores", prov_data)
                        st.rerun()

        # Agregar nuevo proveedor
        with st.expander("➕ Agregar nuevo proveedor"):
            with st.form("form_nuevo_proveedor"):
                np1, np2 = st.columns(2)
                nuevo_prov_nombre = np1.text_input("Nombre del proveedor")
                nuevo_prov_contacto = np2.text_input("Contacto (WhatsApp, email, etc)")
                nuevo_prov_web = st.text_input("Web / Alibaba", "")
                submit_prov = st.form_submit_button("➕ Crear proveedor", use_container_width=True)
                if submit_prov and nuevo_prov_nombre:
                    if nuevo_prov_nombre not in suppliers:
                        suppliers[nuevo_prov_nombre] = {
                            "contacto": nuevo_prov_contacto,
                            "web": nuevo_prov_web,
                            "updated_at": "",
                            "productos": {},
                        }
                        prov_data["suppliers"] = suppliers
                        st.session_state.proveedores_data = prov_data
                        gs_write("Proveedores", prov_data)
                        st.success(f"✅ Proveedor '{nuevo_prov_nombre}' creado")
                        st.rerun()
                    else:
                        st.warning("Ya existe un proveedor con ese nombre.")

        # ════════════════════════════════════════════════════════════════════
        # SECCIÓN 2: COMPARADOR
        # ════════════════════════════════════════════════════════════════════
        st.markdown('<div class="prov-section">Comparador de precios</div>', unsafe_allow_html=True)
        st.markdown(
            '<div class="comp-banner">'
            '<div>'
            '<div class="comp-banner-title">⚖️ Comparador de precios</div>'
            '<div class="comp-banner-sub">Todos los catálogos en una tabla. Verde = más barato entre proveedores.</div>'
            '</div>'
            '</div>',
            unsafe_allow_html=True,
        )

        @st.fragment
        def _render_comparador(suppliers_snap):
            if len(suppliers_snap) < 1:
                st.info("Cargá al menos un proveedor para ver la comparación.")
                return

            _col_search, _col_stock = st.columns([3, 1])
            buscar = _col_search.text_input(
                "🔍 Buscar producto", key="search_comparador",
                placeholder="ej: RG35, Trimui, 477…"
            )
            solo_stock = _col_stock.checkbox(
                "🏪 Solo mis productos", key="comp_solo_stock",
                help="Muestra solo los productos que ya tenés en TiendaNube (con o sin stock)"
            )
            df_fuzzy = fuzzy_group_products(suppliers_snap)

            if df_fuzzy.empty:
                st.info("No hay productos cargados en los catálogos.")
                return

            # Buscar solo cuando no hay filtro TN (si hay filtro TN, se aplica al final)
            if buscar and not solo_stock:
                df_fuzzy = df_fuzzy[
                    df_fuzzy["Producto"].str.contains(buscar, case=False, na=False)
                ]

            # Filtro "Solo mis productos" — construye DESDE TN + datos CRUDOS del proveedor
            # NO usa fuzzy_group_products (que fusiona productos distintos como RG40XXH/V)
            if solo_stock:
                with st.spinner("Consultando TiendaNube…"):
                    _stock_dict = get_stock_for_planner()

                if _stock_dict:
                    _EXCLUIR   = ("estuche", "funda", "case", "protector", "cable", "adaptador")
                    _MARCAS_TN = ("anbernic", "powkiddy", "miyoo", "trimui", "retroid")

                    # Lista de consolas en TN (sin accesorios)
                    _tn_consolas = sorted([
                        tn_name for tn_name in _stock_dict
                        if any(m in tn_name.lower() for m in _MARCAS_TN)
                        and not any(tn_name.lower().startswith(e) for e in _EXCLUIR)
                    ])

                    # ── Leer datos CRUDOS de proveedores (sin fuzzy grouping) ──
                    _sup_names = list(suppliers_snap.keys())
                    # norm → {sup1: price, sup2: price, ...}
                    _raw_lookup = {}
                    for _sname, _sdata in suppliers_snap.items():
                        for _pname, _pinfo in _sdata.get("productos", {}).items():
                            try:
                                _price = float(_pinfo.get("precio_usd", 0))
                            except (TypeError, ValueError):
                                continue
                            if _price <= 0:
                                continue
                            _cl = re.sub(r'\s*\d+\+\d+\s*gb.*', '', _pname, flags=re.IGNORECASE).strip()
                            _norm = re.sub(r'[^a-z0-9]', '', _cl.lower())
                            if _norm not in _raw_lookup:
                                _raw_lookup[_norm] = {}
                            _raw_lookup[_norm][_sname] = _price

                    def _cat_variants(cat):
                        variants = [cat]
                        if cat.startswith("rg") and len(cat) > 4:
                            variants.append(cat[2:])
                        return variants

                    def _find_prices(tn_norm):
                        """Busca precios de TODOS los proveedores para un producto TN.
                        Devuelve dict {supplier: price} o {}."""
                        best_prices, best_len = {}, 0
                        for _norm, _prices in _raw_lookup.items():
                            for _v in _cat_variants(_norm):
                                if _v and _v in tn_norm and len(_v) > best_len:
                                    best_len = len(_v)
                                    best_prices = _prices
                        return best_prices

                    # Construir una fila por cada consola de TN
                    _new_rows = []
                    for _tn_name in _tn_consolas:
                        _tn_norm = re.sub(r'[^a-z0-9]', '', _tn_name.lower())
                        _prices = _find_prices(_tn_norm)
                        _new_rows.append({
                            "Producto": _tn_name,
                            **{s: _prices.get(s) for s in _sup_names}
                        })

                    df_fuzzy = pd.DataFrame(_new_rows).reset_index(drop=True)

                    if buscar:
                        df_fuzzy = df_fuzzy[
                            df_fuzzy["Producto"].str.contains(buscar, case=False, na=False)
                        ]
                else:
                    st.caption("⚠️ No se pudo conectar con Tienda Nube.")

            if df_fuzzy.empty:
                st.info("No se encontraron productos con ese filtro.")
                return

            sup_cols = [c for c in df_fuzzy.columns if c != "Producto"]
            if len(sup_cols) == 1:
                st.caption("Solo hay 1 proveedor cargado — cargá más para comparar precios.")

            # Tabla HTML custom: precios grandes, verde para el más barato
            CSS = """
            <style>
            .comp-wrap{overflow-x:auto;margin-top:8px;}
            .comp-tbl{width:100%;border-collapse:collapse;font-family:inherit;}
            .comp-tbl thead tr{border-bottom:2px solid #2d3748;}
            .comp-tbl th{padding:10px 18px;text-align:left;font-size:11px;font-weight:700;
                color:#6b7280;text-transform:uppercase;letter-spacing:.8px;}
            .comp-tbl th.pc{text-align:right;}
            .comp-tbl td{padding:11px 18px;border-bottom:1px solid #1a2030;vertical-align:middle;}
            .comp-tbl tr:hover td{background:rgba(255,255,255,0.025);}
            .pname{font-size:15px;font-weight:600;color:#e2e8f0;}
            .pval{display:block;text-align:right;font-size:20px;font-weight:700;color:#9ca3af;}
            .pval.cheap{color:#4ade80;font-size:22px;}
            .pval.none{font-size:13px;color:#374151;}
            td.cheap-cell{background:rgba(20,50,20,.55);border-radius:6px;}
            </style>
            """
            _BRAND_CSS_MAP = {
                "Anbernic": "brand-anbernic", "Powkiddy": "brand-powkiddy",
                "Trimui": "brand-trimui", "Miyoo": "brand-miyoo", "Retroid": "brand-retroid",
            }
            th_prod = '<th>Producto</th>'
            th_sups = "".join(f'<th class="pc">{c}</th>' for c in sup_cols)
            rows_html = []
            for _, row in df_fuzzy.iterrows():
                prices = {c: row[c] for c in sup_cols if pd.notna(row[c]) and row[c] > 0}
                cheapest = min(prices, key=prices.get) if len(prices) > 0 else None
                brand = _get_brand(row["Producto"])
                bcss = _BRAND_CSS_MAP.get(brand, "")
                pname = row["Producto"]
                # Prefijar marca si no está ya en el nombre (evita "Trimui Trimui Brick")
                if brand != "—" and not pname.lower().startswith(brand.lower()):
                    brand_span = f'<span class="brand-tag {bcss}" style="margin-right:7px">{brand}</span>'
                    display_name = f'{brand_span}{pname}'
                else:
                    display_name = pname
                tds = f'<td><span class="pname">{display_name}</span></td>'
                for c in sup_cols:
                    v = row[c]
                    if pd.notna(v) and v > 0:
                        if c == cheapest and len(prices) > 1:
                            tds += f'<td class="cheap-cell"><span class="pval cheap">${v:.1f}</span></td>'
                        else:
                            tds += f'<td><span class="pval">${v:.1f}</span></td>'
                    else:
                        tds += '<td><span class="pval none">—</span></td>'
                rows_html.append(f"<tr>{tds}</tr>")

            tabla = (
                f'{CSS}<div class="comp-wrap"><table class="comp-tbl">'
                f"<thead><tr>{th_prod}{th_sups}</tr></thead>"
                f'<tbody>{"".join(rows_html)}</tbody>'
                f"</table></div>"
            )
            st.markdown(tabla, unsafe_allow_html=True)

            # Comparación vs CostosConsolas
            st.divider()
            st.subheader("📊 vs costos cargados en el dashboard")
            _costos_gs_prov = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}
            if _costos_gs_prov:
                rows_comp = []
                for _, prod_row in df_fuzzy.iterrows():
                    prod_name = prod_row["Producto"]
                    fob_dash = get_fob_usd(prod_name, _costos_gs_prov)
                    prices_prov = {c: prod_row[c] for c in sup_cols if pd.notna(prod_row[c])}
                    if not prices_prov or fob_dash <= 0:
                        continue
                    best_sup = min(prices_prov, key=prices_prov.get)
                    fob_prov = prices_prov[best_sup]
                    diff_val = round(fob_prov - fob_dash, 2)
                    rows_comp.append({
                        "Producto": prod_name,
                        "FOB Dashboard (USD)": fob_dash,
                        "Mejor proveedor": f"{best_sup} ${fob_prov:.1f}",
                        "Diferencia (USD)": diff_val,
                        "Status": "✅ Más barato" if diff_val < -0.5 else ("⚠️ Más caro" if diff_val > 0.5 else "≈ Similar"),
                    })
                if rows_comp:
                    df_vs = pd.DataFrame(rows_comp).sort_values("Diferencia (USD)")
                    st.dataframe(
                        df_vs.style.format({"FOB Dashboard (USD)": "${:.2f}", "Diferencia (USD)": "${:+.2f}"}),
                        use_container_width=True, hide_index=True,
                    )
                else:
                    st.info("No hay coincidencias entre catálogos y costos del dashboard.")
            else:
                st.info("Cargá costos en 💻 Costos de consolas para comparar.")

        _render_comparador(suppliers)

        # ════════════════════════════════════════════════════════════════════
        # SECCIÓN 3: PLANIFICADOR DE COMPRA
        # ════════════════════════════════════════════════════════════════════
        st.divider()
        st.subheader("🛒 Planificador de compra")
        st.caption("Cruzá tu stock actual con los precios de los catálogos. Seleccioná qué reponer y generá el resumen del pedido por proveedor.")

        if not suppliers:
            st.info("Cargá al menos un proveedor para usar el planificador.")
        else:
            if st.button("🔄 Cargar stock desde Tienda Nube", key="btn_load_stock_planner"):
                with st.spinner("Consultando stock..."):
                    st.session_state.stock_planner = get_stock_for_planner()
                st.success(f"✅ {len(st.session_state.stock_planner)} productos cargados")

            if "stock_planner" not in st.session_state:
                st.info("Hacé clic en 'Cargar stock' para ver los niveles actuales.")
            else:
                from difflib import SequenceMatcher as _SM

                stock_dict = st.session_state.stock_planner

                def _best_stock_match(prod_name: str, stock_dict: dict) -> int | None:
                    """Busca el stock en TN usando fuzzy match. None = sin límite / no encontrado."""
                    best_ratio = 0.0
                    best_val = None
                    for tn_name, tn_stock in stock_dict.items():
                        ratio = _SM(None, _normalizar(prod_name), _normalizar(tn_name)).ratio()
                        if ratio > best_ratio:
                            best_ratio = ratio
                            best_val = tn_stock
                    return best_val if best_ratio >= 0.70 else "?"

                def _best_supplier_price(prod_name: str, suppliers: dict) -> tuple[str, float]:
                    """Devuelve (proveedor, precio) del proveedor más barato para ese producto."""
                    candidates = []
                    for sup_name, sup_data in suppliers.items():
                        best_ratio_sup = 0.0
                        best_price_sup = 0.0
                        for p_name, p_info in sup_data.get("productos", {}).items():
                            ratio = _SM(None, _normalizar(prod_name), _normalizar(p_name)).ratio()
                            if ratio >= 0.80 and ratio > best_ratio_sup:
                                best_ratio_sup = ratio
                                best_price_sup = float(p_info.get("precio_usd", 0))
                        if best_ratio_sup >= 0.80:
                            candidates.append((sup_name, best_price_sup))
                    if not candidates:
                        return ("—", 0.0)
                    return min(candidates, key=lambda x: x[1])

                # Construir tabla del planificador (un producto por fila, sin duplicados)
                seen_products: set = set()
                plan_rows = []
                for sup_name, sup_data in suppliers.items():
                    for prod_name in sup_data.get("productos", {}).keys():
                        canonical = _normalizar(prod_name)
                        if canonical in seen_products:
                            continue
                        seen_products.add(canonical)

                        stock_val = _best_stock_match(prod_name, stock_dict)
                        best_sup, best_price = _best_supplier_price(prod_name, suppliers)

                        if isinstance(stock_val, int):
                            urgency = "🔴 Crítico" if stock_val <= 2 else ("🟡 Bajo" if stock_val <= 5 else "🟢 OK")
                        elif stock_val is None:
                            urgency = "♾️ Sin límite"
                        else:
                            urgency = "❓ Sin dato"

                        if best_sup == "—":
                            continue
                        plan_rows.append({
                            "Producto": prod_name,
                            "Stock": stock_val if stock_val not in (None, "?") else ("Sin límite" if stock_val is None else "—"),
                            "Urgencia": urgency,
                            "Mejor precio (USD)": best_price,
                            "Proveedor sugerido": best_sup,
                            "Comprar": False,
                            "Cantidad": 1,
                        })

                if not plan_rows:
                    st.info("No hay productos en los catálogos cargados.")
                else:
                    # Ordenar por urgencia (Crítico primero)
                    urgency_order = {"🔴 Crítico": 0, "🟡 Bajo": 1, "🟢 OK": 2, "❓ Sin dato": 3, "♾️ Sin límite": 4}
                    plan_rows.sort(key=lambda r: urgency_order.get(r["Urgencia"], 5))

                    df_plan = pd.DataFrame(plan_rows)
                    edited_plan = st.data_editor(
                        df_plan,
                        column_config={
                            "Comprar": st.column_config.CheckboxColumn("✓ Comprar", default=False),
                            "Cantidad": st.column_config.NumberColumn("Cant.", min_value=1, step=1, default=1),
                            "Mejor precio (USD)": st.column_config.NumberColumn("Precio (USD)", format="$%.1f", disabled=True),
                            "Producto": st.column_config.TextColumn("Producto", disabled=True),
                            "Stock": st.column_config.TextColumn("Stock", disabled=True),
                            "Urgencia": st.column_config.TextColumn("Urgencia", disabled=True),
                            "Proveedor sugerido": st.column_config.TextColumn("Proveedor", disabled=True),
                        },
                        hide_index=True,
                        use_container_width=True,
                        key="plan_table",
                    )

                    selected = edited_plan[edited_plan["Comprar"] == True]
                    if not selected.empty:
                        st.divider()
                        st.subheader("📋 Resumen del pedido")

                        by_supplier: dict = {}
                        for _, row in selected.iterrows():
                            sup = str(row["Proveedor sugerido"])
                            if sup not in by_supplier:
                                by_supplier[sup] = []
                            by_supplier[sup].append({
                                "producto": str(row["Producto"]),
                                "cantidad": int(row["Cantidad"]),
                                "precio": float(row["Mejor precio (USD)"]),
                            })

                        order_text = ""
                        grand_total = 0.0
                        for sup, items in by_supplier.items():
                            subtotal = sum(i["cantidad"] * i["precio"] for i in items)
                            grand_total += subtotal
                            st.markdown(f"**📦 Pedido a {sup} — Total: ${subtotal:.0f} USD**")
                            order_text += f"Pedido a {sup} — Total: ${subtotal:.0f} USD\n"
                            for item in items:
                                line_total = item["cantidad"] * item["precio"]
                                st.markdown(f"  · {item['cantidad']}× {item['producto']} × ${item['precio']:.0f} = ${line_total:.0f}")
                                order_text += f"  · {item['cantidad']}x {item['producto']} x ${item['precio']:.0f}\n"
                            order_text += "\n"

                        st.markdown(f"**💰 Total general: ${grand_total:.0f} USD**")
                        order_text += f"Total general: ${grand_total:.0f} USD"

                        st.divider()
                        st.markdown("**Copiá esto para WhatsApp:**")
                        st.code(order_text, language=None)

        # ── Guardado global ──
        st.divider()
        if st.button("💾 Guardar todos los datos de proveedores", use_container_width=True):
            ok = gs_write("Proveedores", st.session_state.proveedores_data)
            st.success("✅ Guardado en Google Sheets" if ok else "⚠️ Solo en sesión")

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 11: ESTADÍSTICAS DE PAGO
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "💳 Estadísticas de pago":
        st.subheader("💳 Estadísticas de pago")

        if df_tn.empty:
            st.warning("Sin datos para el período seleccionado.")
        else:
            # ── Preparar columnas de fecha ──
            df_pago = df_tn.copy()
            df_pago["Fecha"] = pd.to_datetime(df_pago["Fecha"], errors="coerce")
            df_pago["Semana"] = df_pago["Fecha"].dt.to_period("W").apply(lambda r: str(r.start_time.date()))
            df_pago["Dia"] = df_pago["Fecha"].dt.date

            # ── KPIs rápidos ──
            total_fact = df_pago["Total ($)"].sum()
            total_ordenes = len(df_pago)
            ticket_prom = total_fact / total_ordenes if total_ordenes else 0
            costo_pn_prom = df_pago["Costo PN (%)"].mean() if "Costo PN (%)" in df_pago.columns else 0

            kp1, kp2, kp3, kp4 = st.columns(4)
            kp1.metric("Facturación total", fmt(total_fact))
            kp2.metric("Órdenes", total_ordenes)
            kp3.metric("Ticket promedio", fmt(ticket_prom))
            kp4.metric("Costo PN promedio", f"{costo_pn_prom:.2f}%")

            st.divider()

            # ── Sección 1: Por medio de pago ──
            st.markdown("### 💳 Por medio de pago")

            agg_medio = (
                df_pago.groupby("Medio de Pago")
                .agg(
                    Ordenes=("Orden", "count"),
                    Facturación=("Total ($)", "sum"),
                    Neto=("Neto cobrado ($)", "sum"),
                    Comisión=("Comision PN ($)", "sum"),
                )
                .reset_index()
            )
            agg_medio["Ticket prom"] = agg_medio["Facturación"] / agg_medio["Ordenes"]
            agg_medio["Costo PN %"] = (agg_medio["Comisión"] / agg_medio["Facturación"] * 100).round(2)
            agg_medio["Part. fact %"] = (agg_medio["Facturación"] / agg_medio["Facturación"].sum() * 100).round(1)
            agg_medio = agg_medio.sort_values("Facturación", ascending=False)

            col_chart1, col_chart2 = st.columns(2)

            with col_chart1:
                fig_fact_medio = px.bar(
                    agg_medio,
                    x="Medio de Pago",
                    y="Facturación",
                    color="Medio de Pago",
                    color_discrete_sequence=COLORES,
                    title="Facturación por medio de pago",
                    text=agg_medio["Facturación"].apply(lambda v: fmt(v)),
                )
                fig_fact_medio.update_traces(textposition="outside")
                fig_fact_medio.update_layout(showlegend=False, height=350, margin=dict(t=40, b=0))
                st.plotly_chart(fig_fact_medio, use_container_width=True)

            with col_chart2:
                fig_ord_medio = px.bar(
                    agg_medio,
                    x="Medio de Pago",
                    y="Ordenes",
                    color="Medio de Pago",
                    color_discrete_sequence=COLORES,
                    title="Órdenes por medio de pago",
                    text="Ordenes",
                )
                fig_ord_medio.update_traces(textposition="outside")
                fig_ord_medio.update_layout(showlegend=False, height=350, margin=dict(t=40, b=0))
                st.plotly_chart(fig_ord_medio, use_container_width=True)

            # Donut participación en facturación
            col_donut1, col_donut2 = st.columns(2)
            with col_donut1:
                fig_donut_fact = px.pie(
                    agg_medio,
                    names="Medio de Pago",
                    values="Facturación",
                    hole=0.5,
                    color_discrete_sequence=COLORES,
                    title="Participación en facturación",
                )
                fig_donut_fact.update_traces(textposition="outside", textinfo="label+percent")
                fig_donut_fact.update_layout(showlegend=False, height=320, margin=dict(t=40, b=0))
                st.plotly_chart(fig_donut_fact, use_container_width=True)

            with col_donut2:
                fig_costo_medio = px.bar(
                    agg_medio.sort_values("Costo PN %", ascending=True),
                    x="Costo PN %",
                    y="Medio de Pago",
                    orientation="h",
                    color="Costo PN %",
                    color_continuous_scale=["#00C49F", "#FFD700", "#FF5733"],
                    title="Costo Pago Nube % por método",
                    text=agg_medio.sort_values("Costo PN %", ascending=True)["Costo PN %"].apply(lambda v: f"{v:.2f}%"),
                )
                fig_costo_medio.update_traces(textposition="outside")
                fig_costo_medio.update_layout(showlegend=False, height=320, margin=dict(t=40, b=0), coloraxis_showscale=False)
                st.plotly_chart(fig_costo_medio, use_container_width=True)

            # Tabla resumen
            tabla_medio = agg_medio[["Medio de Pago", "Ordenes", "Facturación", "Ticket prom", "Costo PN %", "Part. fact %"]].copy()
            tabla_medio_fmt = tabla_medio.copy()
            tabla_medio_fmt["Facturación"] = tabla_medio["Facturación"].apply(fmt)
            tabla_medio_fmt["Ticket prom"] = tabla_medio["Ticket prom"].apply(fmt)
            tabla_medio_fmt["Costo PN %"] = tabla_medio["Costo PN %"].apply(lambda v: f"{v:.2f}%")
            tabla_medio_fmt["Part. fact %"] = tabla_medio["Part. fact %"].apply(lambda v: f"{v:.1f}%")
            st.dataframe(tabla_medio_fmt, use_container_width=True, hide_index=True)

            st.divider()

            # ── Sección 2: Por cuotas ──
            st.markdown("### 🔢 Por cuotas")

            df_pago["Cuotas_label"] = df_pago["Cuotas"].apply(
                lambda c: f"{int(c)} cuota{'s' if int(c) > 1 else ''}"
            )
            agg_cuotas = (
                df_pago.groupby(["Cuotas", "Cuotas_label"])
                .agg(
                    Ordenes=("Orden", "count"),
                    Facturación=("Total ($)", "sum"),
                    Comisión=("Comision PN ($)", "sum"),
                )
                .reset_index()
                .sort_values("Cuotas")
            )
            agg_cuotas["Ticket prom"] = agg_cuotas["Facturación"] / agg_cuotas["Ordenes"]
            agg_cuotas["Costo PN %"] = (agg_cuotas["Comisión"] / agg_cuotas["Facturación"] * 100).round(2)
            agg_cuotas["Part. fact %"] = (agg_cuotas["Facturación"] / agg_cuotas["Facturación"].sum() * 100).round(1)

            col_cuotas1, col_cuotas2 = st.columns(2)

            with col_cuotas1:
                fig_fact_cuotas = px.bar(
                    agg_cuotas,
                    x="Cuotas_label",
                    y="Facturación",
                    color="Cuotas_label",
                    color_discrete_sequence=COLORES,
                    title="Facturación por cuotas",
                    text=agg_cuotas["Facturación"].apply(lambda v: fmt(v)),
                )
                fig_fact_cuotas.update_traces(textposition="outside")
                fig_fact_cuotas.update_layout(showlegend=False, height=350, margin=dict(t=40, b=0), xaxis_title="")
                st.plotly_chart(fig_fact_cuotas, use_container_width=True)

            with col_cuotas2:
                fig_ord_cuotas = px.bar(
                    agg_cuotas,
                    x="Cuotas_label",
                    y="Ordenes",
                    color="Cuotas_label",
                    color_discrete_sequence=COLORES,
                    title="Órdenes por cuotas",
                    text="Ordenes",
                )
                fig_ord_cuotas.update_traces(textposition="outside")
                fig_ord_cuotas.update_layout(showlegend=False, height=350, margin=dict(t=40, b=0), xaxis_title="")
                st.plotly_chart(fig_ord_cuotas, use_container_width=True)

            # Tabla cuotas
            tabla_cuotas = agg_cuotas[["Cuotas_label", "Ordenes", "Facturación", "Ticket prom", "Costo PN %", "Part. fact %"]].copy()
            tabla_cuotas_fmt = tabla_cuotas.copy()
            tabla_cuotas_fmt.columns = ["Cuotas", "Órdenes", "Facturación", "Ticket prom", "Costo PN %", "Part. fact %"]
            tabla_cuotas_fmt["Facturación"] = tabla_cuotas["Facturación"].apply(fmt)
            tabla_cuotas_fmt["Ticket prom"] = tabla_cuotas["Ticket prom"].apply(fmt)
            tabla_cuotas_fmt["Costo PN %"] = tabla_cuotas["Costo PN %"].apply(lambda v: f"{v:.2f}%")
            tabla_cuotas_fmt["Part. fact %"] = tabla_cuotas["Part. fact %"].apply(lambda v: f"{v:.1f}%")
            st.dataframe(tabla_cuotas_fmt, use_container_width=True, hide_index=True)

            st.divider()

            # ── Sección 3: Tendencia por método ──
            st.markdown("### 📈 Tendencia de facturación por método")

            # Agrupar por día y método
            tend_medio = (
                df_pago.groupby(["Dia", "Medio de Pago"])["Total ($)"]
                .sum()
                .reset_index()
            )
            tend_medio["Dia"] = pd.to_datetime(tend_medio["Dia"])
            tend_medio = tend_medio.sort_values("Dia")

            # Solo mostrar si hay suficientes días
            dias_unicos = tend_medio["Dia"].nunique()
            if dias_unicos >= 3:
                fig_tend = px.line(
                    tend_medio,
                    x="Dia",
                    y="Total ($)",
                    color="Medio de Pago",
                    color_discrete_sequence=COLORES,
                    title="Facturación diaria por medio de pago",
                    markers=True,
                )
                fig_tend.update_layout(height=380, margin=dict(t=40, b=0), xaxis_title="", yaxis_title="Facturación ($)")
                st.plotly_chart(fig_tend, use_container_width=True)
            else:
                st.info("📅 Se necesitan al menos 3 días de datos para mostrar la tendencia.")

            # ── Sección 4: Mapa de calor método × cuotas ──
            st.markdown("### 🗺️ Facturación: método × cuotas")

            pivot = df_pago.pivot_table(
                index="Medio de Pago",
                columns="Cuotas",
                values="Total ($)",
                aggfunc="sum",
                fill_value=0,
            )
            pivot.columns = [f"{int(c)}c" for c in pivot.columns]
            pivot = pivot.reset_index()

            fig_heat = px.imshow(
                pivot.set_index("Medio de Pago"),
                color_continuous_scale=["#0d1b2a", "#009EE3", "#00C49F"],
                text_auto=".3s",
                aspect="auto",
                title="Facturación ($) por método y cuotas",
            )
            fig_heat.update_layout(height=300, margin=dict(t=40, b=0), coloraxis_showscale=False)
            st.plotly_chart(fig_heat, use_container_width=True)

    # ══════════════════════════════════════════════════════════════════════════
    # TAB: WEB / ANALYTICS (Google Analytics 4)
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🌐 Web / Analytics":
        st.subheader("🌐 Web / Analytics")
        st.caption("Comportamiento del tráfico de la web. Fuente: Google Analytics 4.")

        if not st.secrets.get("GA4_PROPERTY_ID"):
            st.warning(
                "⚠️ Falta configurar `GA4_PROPERTY_ID` en `secrets.toml`. "
                "Cargá el ID de propiedad GA4 (numérico, no el measurement ID `G-XXXX`) "
                "para activar esta solapa."
            )
        else:
            # ── Selector de período comparativo ──
            ga4_periodos = {
                "28d vs 28d previos": "28d",
                "7d vs 7d previos": "7d",
                "Mes actual vs mes anterior (MoM)": "MoM",
                "Últimos 28d vs mismo período año anterior (YoY)": "YoY",
            }
            _periodo_label = st.radio(
                "Período comparativo",
                list(ga4_periodos.keys()),
                index=0,
                horizontal=True,
                help=(
                    "7d: detección rápida pero ruidoso. "
                    "28d: la 'verdad' de la tendencia. "
                    "MoM: alineado con cómo opera el resto del dashboard. "
                    "YoY: estacionalidad (HotSale, Navidad, Día del Padre)."
                ),
            )
            ga4_periodo_key = ga4_periodos[_periodo_label]

            with st.spinner("Consultando Google Analytics..."):
                ga4 = get_ga4_metrics(ga4_periodo_key)

            if not ga4:
                st.warning(
                    "⚠️ No se pudieron obtener datos de GA4. Revisá que el service account "
                    "tenga permiso de Lector en la propiedad y que la Data API esté habilitada."
                )
            else:
                def _delta_pct(act, prev):
                    if prev <= 0:
                        return None
                    return (act - prev) / prev * 100

                def _delta_str(d):
                    if d is None:
                        return "—"
                    signo = "+" if d >= 0 else ""
                    return f"{signo}{d:.1f}%"

                k = ga4["kpis"]
                sesiones_act, sesiones_prev = k["sesiones"]
                usuarios_act, usuarios_prev = k["usuarios"]
                interaccion_act, interaccion_prev = k["interaccion"]
                checkouts_act, checkouts_prev = k["checkouts"]
                duracion_act, duracion_prev = k["duracion"]
                bounce_act, bounce_prev = k["bounce"]

                # Tasa de conversión a checkout = checkouts ÷ sesiones
                conv_act = (checkouts_act / sesiones_act * 100) if sesiones_act > 0 else 0.0
                conv_prev = (checkouts_prev / sesiones_prev * 100) if sesiones_prev > 0 else 0.0

                def _fmt_seg(seg):
                    seg = float(seg or 0)
                    if seg >= 60:
                        m = int(seg // 60)
                        s = int(seg % 60)
                        return f"{m}m {s:02d}s"
                    return f"{seg:.0f}s"

                def _fmt_pp(act, prev):
                    """Delta en puntos porcentuales para métricas ya en %."""
                    if prev is None:
                        return "—"
                    d = act - prev
                    return f"{'+' if d >= 0 else ''}{d:.1f} pp"

                def _calidad_y_score(items, vistas_key="vistas", tiempo_key="tiempo_prom_seg"):
                    """Calcula score (Bayesian del tiempo) y cuadrante de calidad.

                    Idea: el promedio bayesiano arrastra a los items con pocas vistas
                    hacia el promedio global, evitando falsos positivos (1 vista de 5min)
                    y falsos negativos (volumen alto con engagement bajo).
                    Modifica items in-place agregando 'score_seg' y 'calidad'.
                    """
                    if not items:
                        return
                    import statistics as _stats
                    vistas_validas = [it[vistas_key] for it in items if it[vistas_key] > 0]
                    if not vistas_validas:
                        for it in items:
                            it["score_seg"] = 0.0
                            it["calidad"] = "⏸ Sin data"
                        return
                    # Tiempo promedio global ponderado por vistas
                    total_eng = sum(it[tiempo_key] * it[vistas_key] for it in items)
                    total_v = sum(vistas_validas)
                    tiempo_global = (total_eng / total_v) if total_v > 0 else 0.0
                    mediana_vistas = _stats.median(vistas_validas)
                    # Peso del prior bayesiano = mediana de vistas (productos por debajo
                    # de la mediana son arrastrados al promedio global; los que están
                    # encima conservan su tiempo real)
                    C = max(mediana_vistas, 1)
                    # Calcular score para todos
                    scores = []
                    for it in items:
                        n = it[vistas_key]
                        t = it[tiempo_key]
                        bayes = (n * t + C * tiempo_global) / (n + C) if (n + C) > 0 else 0.0
                        it["score_seg"] = bayes
                        scores.append(bayes)
                    mediana_score = _stats.median(scores) if scores else 0.0
                    # Threshold mínimo para no clasificar ruido
                    min_vistas = max(3, mediana_vistas / 4)
                    for it in items:
                        n = it[vistas_key]
                        b = it["score_seg"]
                        if n < min_vistas:
                            it["calidad"] = "⏸ Sin data"
                        elif n >= mediana_vistas and b >= mediana_score:
                            it["calidad"] = "🟢 Ganador"
                        elif n < mediana_vistas and b >= mediana_score:
                            it["calidad"] = "🔵 Subdetectado"
                        elif n >= mediana_vistas and b < mediana_score:
                            it["calidad"] = "🟡 Desperdiciado"
                        else:
                            it["calidad"] = "⚫ Bajo perfil"

                # ── Calcular score y calidad para productos y páginas ──────────
                # productos_full es el set grande y representativo para clasificar productos
                _calidad_y_score(ga4.get("productos_full") or [])
                # Map de lookup por path: para enriquecer top_productos (5)
                _quality_by_path = {
                    p["path"]: (p.get("calidad", "⏸ Sin data"), p.get("score_seg", 0.0))
                    for p in (ga4.get("productos_full") or [])
                }
                # Aplicar mismo lookup a top_productos (que vienen con key 'pagina')
                for p in (ga4.get("top_productos") or []):
                    cal, sc = _quality_by_path.get(p["pagina"], ("⏸ Sin data", 0.0))
                    p["calidad"] = cal
                    p["score_seg"] = sc
                # top_paginas (30 items mixtos): clasificar contra sí mismo
                _calidad_y_score(ga4.get("top_paginas") or [], tiempo_key="tiempo_prom_seg")

                # ── Fila 1: KPIs de tráfico ─────────────────────────────────────
                c1, c2, c3, c4 = st.columns(4)
                with c1:
                    st.metric(
                        "Sesiones",
                        f"{int(sesiones_act):,}".replace(",", "."),
                        _delta_str(_delta_pct(sesiones_act, sesiones_prev)),
                        help=(
                            "Cantidad de visitas al sitio. Una persona puede generar varias "
                            "sesiones si vuelve después de 30 min de inactividad. Es la métrica "
                            "base de volumen: si cae, llega menos gente; si sube, está "
                            "funcionando algún canal de tráfico."
                        ),
                    )
                with c2:
                    st.metric(
                        "Usuarios activos",
                        f"{int(usuarios_act):,}".replace(",", "."),
                        _delta_str(_delta_pct(usuarios_act, usuarios_prev)),
                        help=(
                            "Personas únicas que entraron. Comparado con sesiones te dice si "
                            "viene gente nueva (números parecidos) o si la misma gente vuelve "
                            "varias veces (sesiones >> usuarios)."
                        ),
                    )
                with c3:
                    st.metric(
                        "Tasa de interacción",
                        f"{interaccion_act * 100:.1f}%",
                        _fmt_pp(interaccion_act * 100, interaccion_prev * 100),
                        help=(
                            "% de sesiones donde la persona se quedó +10s, vio +1 página o "
                            "disparó algún evento. Es el inverso del bounce. Arriba de 50% "
                            "está OK; abajo de 40% indica tráfico de baja calidad."
                        ),
                    )
                with c4:
                    st.metric(
                        "Checkouts iniciados",
                        f"{int(checkouts_act):,}".replace(",", "."),
                        _delta_str(_delta_pct(checkouts_act, checkouts_prev)),
                        help=(
                            "Personas que llegaron al paso de checkout (/checkout/v3/start). "
                            "No mide ventas cerradas, mide INTENCIÓN de comprar. La caída de "
                            "ventas con este número estable = problema en pago/envío."
                        ),
                    )

                # ── Fila 2: KPIs de comportamiento / conversión ──────────────────
                c5, c6, c7, c8 = st.columns(4)
                with c5:
                    st.metric(
                        "Conv. a checkout",
                        f"{conv_act:.2f}%",
                        _fmt_pp(conv_act, conv_prev),
                        help=(
                            "Checkouts iniciados ÷ sesiones. El número más importante del "
                            "funnel: dice qué % del tráfico llega a querer comprar. "
                            "Benchmark e-commerce: 2-3% es OK, 4-5% es muy bueno. Caída acá "
                            "antes que en ventas = problema upstream (pauta mala, UX rota)."
                        ),
                    )
                with c6:
                    st.metric(
                        "Tiempo prom. sesión",
                        _fmt_seg(duracion_act),
                        _delta_str(_delta_pct(duracion_act, duracion_prev)),
                        help=(
                            "Cuánto tiempo total pasa la persona en el sitio por visita. "
                            "Suma todas las páginas que vio. >2 min indica interés real; "
                            "<30s suele ser tráfico de pauta mal segmentada."
                        ),
                    )
                with c7:
                    st.metric(
                        "Bounce rate",
                        f"{bounce_act * 100:.1f}%",
                        _fmt_pp(bounce_act * 100, bounce_prev * 100),
                        delta_color="inverse",
                        help=(
                            "% de sesiones que rebotan sin interactuar (1 página, <10s, "
                            "ningún evento). Menos es mejor. Abajo de 50% está OK; arriba "
                            "de 70% señal de tráfico que cae por error o landing no convence."
                        ),
                    )
                with c8:
                    ppu = (sesiones_act / max(usuarios_act, 1))
                    st.metric(
                        "Sesiones / usuario",
                        f"{ppu:.2f}",
                        help=(
                            "Cuántas veces vuelve la misma persona en el período. >1.5 "
                            "indica buena recurrencia (te están considerando varias veces); "
                            "≈1.0 es tráfico de descubrimiento puro (no vuelven)."
                        ),
                    )

                st.divider()

                # ── Tendencia diaria de sesiones + checkouts ────────────────────
                st.subheader("Tendencia diaria")
                st.caption(
                    "Sesiones y checkouts iniciados día a día. Detecta el día exacto "
                    "donde se rompió o despegó algo (lo que el promedio agregado esconde)."
                )
                tendencia = ga4.get("tendencia_diaria") or []
                if not tendencia:
                    st.info("Sin datos de tendencia.")
                else:
                    df_tend = pd.DataFrame(tendencia)
                    df_tend["fecha_dt"] = pd.to_datetime(df_tend["fecha"], format="%Y%m%d", errors="coerce")
                    df_tend = df_tend.sort_values("fecha_dt")
                    import plotly.graph_objects as _go
                    fig_tend = _go.Figure()
                    fig_tend.add_trace(_go.Scatter(
                        x=df_tend["fecha_dt"], y=df_tend["sesiones"],
                        name="Sesiones", mode="lines+markers",
                        line=dict(color=MG_RED, width=2),
                        marker=dict(size=5),
                    ))
                    fig_tend.add_trace(_go.Scatter(
                        x=df_tend["fecha_dt"], y=df_tend["checkouts"],
                        name="Checkouts iniciados", mode="lines+markers",
                        line=dict(color="#4ade80", width=2),
                        marker=dict(size=5),
                        yaxis="y2",
                    ))
                    fig_tend.update_layout(
                        height=320,
                        margin=dict(l=0, r=0, t=20, b=0),
                        paper_bgcolor=MG_BG, plot_bgcolor=MG_BG,
                        font=dict(color=MG_TEXT, family="Hanken Grotesk"),
                        yaxis=dict(title="Sesiones", gridcolor=MG_BORDER, zerolinecolor=MG_BORDER),
                        yaxis2=dict(
                            title="Checkouts", overlaying="y", side="right",
                            gridcolor=MG_BORDER, zerolinecolor=MG_BORDER, showgrid=False,
                        ),
                        xaxis=dict(gridcolor=MG_BORDER, zerolinecolor=MG_BORDER),
                        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
                        hovermode="x unified",
                    )
                    st.plotly_chart(fig_tend, use_container_width=True)

                st.divider()

                # ── Glosario de canales (texto explicativo previo a las tablas) ─
                with st.expander("ℹ️ Qué significa cada canal de tráfico"):
                    st.markdown("""
- **Direct** — entraron escribiendo la URL directo, desde bookmark, o desde apps de mensajería sin link rastreable (WhatsApp, IG DM). Suelen ser clientes que ya te conocen.
- **Organic Search** — vinieron desde Google, Bing, etc. SIN haber hecho click en un anuncio. SEO puro.
- **Paid Search** — Google Ads, Bing Ads (anuncios en buscadores).
- **Organic Social** — posts orgánicos de IG, Facebook, TikTok (no pauta).
- **Paid Social** — Meta Ads, TikTok Ads, anuncios de pauta en redes. Acá cae tu inversión de pauta.
- **Email** — clicks desde campañas de mail marketing.
- **Referral** — link desde otro sitio web (blog, foro, listing).
- **Affiliates** — programas de afiliados.
- **Display** — banners en la red de Google Display u otras.
- **Unassigned** — GA4 no pudo clasificar la fuente. Suele ser tracking incompleto o links sin UTMs.
                    """.strip())

                # ── Sesiones por canal + Dispositivo (lado a lado) ──────────────
                col_canal, col_device = st.columns([2, 1])
                with col_canal:
                    st.subheader("Sesiones por canal")
                    st.caption("De dónde viene el tráfico, con % sobre el total.")
                    canales = ga4["canales"]
                    if not canales:
                        st.info("Sin datos de canales en el período.")
                    else:
                        total_canales = sum(c["sesiones_act"] for c in canales) or 1
                        df_canales = pd.DataFrame([
                            {
                                "Canal": c["canal"] or "(sin dato)",
                                "Sesiones": int(c["sesiones_act"]),
                                "%": f"{c['sesiones_act'] / total_canales * 100:.1f}%",
                                "Variación": _delta_str(_delta_pct(c["sesiones_act"], c["sesiones_prev"])),
                            }
                            for c in canales
                        ])
                        st.dataframe(df_canales, hide_index=True, use_container_width=True)

                with col_device:
                    st.subheader("Dispositivo")
                    st.caption("Mobile vs desktop vs tablet.")
                    dispositivos = ga4["dispositivos"]
                    if not dispositivos:
                        st.info("Sin datos.")
                    else:
                        total_dev = sum(d["sesiones"] for d in dispositivos) or 1
                        df_dev = pd.DataFrame([
                            {
                                "Device": d["device"].title(),
                                "Sesiones": int(d["sesiones"]),
                                "%": f"{d['sesiones'] / total_dev * 100:.1f}%",
                            }
                            for d in dispositivos
                        ])
                        st.dataframe(df_dev, hide_index=True, use_container_width=True)

                # ── Sistema operativo + Resolución (para iterar UI) ─────────────
                col_os, col_res = st.columns(2)
                with col_os:
                    st.subheader("Sistema operativo")
                    st.caption("Android / iOS / Windows / Mac — útil para compatibilidad.")
                    sistemas = ga4["sistemas_op"]
                    if not sistemas:
                        st.info("Sin datos.")
                    else:
                        total_os = sum(s["sesiones"] for s in sistemas) or 1
                        df_os = pd.DataFrame([
                            {
                                "OS": s["os"],
                                "Sesiones": int(s["sesiones"]),
                                "%": f"{s['sesiones'] / total_os * 100:.1f}%",
                            }
                            for s in sistemas[:8]
                        ])
                        st.dataframe(df_os, hide_index=True, use_container_width=True)

                with col_res:
                    st.subheader("Resolución de pantalla")
                    st.caption("Top resoluciones — guía los breakpoints CSS.")
                    resoluciones = ga4["resoluciones"]
                    if not resoluciones:
                        st.info("Sin datos.")
                    else:
                        total_res = sum(r["sesiones"] for r in resoluciones) or 1
                        df_res = pd.DataFrame([
                            {
                                "Resolución": r["resolucion"],
                                "Sesiones": int(r["sesiones"]),
                                "%": f"{r['sesiones'] / total_res * 100:.1f}%",
                            }
                            for r in resoluciones
                        ])
                        st.dataframe(df_res, hide_index=True, use_container_width=True)

                st.divider()

                # ── Checkouts por canal (atribución) ─────────────────────────────
                st.subheader("Checkouts iniciados por canal")
                st.caption(
                    "Qué canal trae a la gente que LLEGA al checkout. "
                    "Comparalo con 'Sesiones por canal': si Paid Social tiene 60% del tráfico "
                    "pero 20% de los checkouts, estás pagando tráfico que no convierte."
                )
                checkouts_canal = ga4.get("checkouts_por_canal") or []
                if not checkouts_canal:
                    st.info("Sin checkouts en el período.")
                else:
                    total_checkouts_canal = sum(c["checkouts"] for c in checkouts_canal) or 1
                    # Mapear sesiones por canal para poder calcular tasa de conversión por canal
                    sesiones_por_canal = {c["canal"]: c["sesiones_act"] for c in ga4["canales"]}
                    df_cc = pd.DataFrame([
                        {
                            "Canal": c["canal"],
                            "Checkouts": int(c["checkouts"]),
                            "% del total": f"{c['checkouts'] / total_checkouts_canal * 100:.1f}%",
                            "Tasa conv.": (
                                f"{c['checkouts'] / sesiones_por_canal[c['canal']] * 100:.2f}%"
                                if sesiones_por_canal.get(c["canal"], 0) > 0 else "—"
                            ),
                        }
                        for c in checkouts_canal
                    ])
                    st.dataframe(df_cc, hide_index=True, use_container_width=True)
                    st.caption(
                        "💡 La columna 'Tasa conv.' = checkouts ÷ sesiones de ese canal. "
                        "El canal con tasa más alta es el que mejor convierte (no el que más tráfico trae)."
                    )

                st.divider()

                # ── Top categorías ──────────────────────────────────────────────
                st.subheader("Top categorías (secciones del sitio)")
                st.caption(
                    "Agrupado por primer segmento del path. La columna 'URL ejemplo' "
                    "te muestra cuál es la página más visitada de cada bucket — útil para "
                    "entender qué cae adentro."
                )
                top_cats = ga4["top_categorias"]
                if not top_cats:
                    st.info("Sin datos.")
                else:
                    total_cats = sum(c["vistas"] for c in top_cats) or 1
                    df_cats = pd.DataFrame([
                        {
                            "Categoría": c["categoria"],
                            "Vistas": int(c["vistas"]),
                            "%": f"{c['vistas'] / total_cats * 100:.1f}%",
                            "Engagement prom.": _fmt_seg(c["tiempo_prom_seg"]),
                            "URL ejemplo": c.get("url_top", ""),
                        }
                        for c in top_cats
                    ])
                    st.dataframe(df_cats, hide_index=True, use_container_width=True)

                st.divider()

                # ── Top productos + Top páginas (lado a lado) ───────────────────
                col_prod, col_pag = st.columns(2)
                with col_prod:
                    st.subheader("Top productos vistos")
                    st.caption(
                        "Páginas /productos/* más miradas, con clasificación de calidad. "
                        "La calidad se calcula contra TODO el catálogo de productos vistos, "
                        "no solo el top 5."
                    )
                    top_prod = ga4["top_productos"]
                    if not top_prod:
                        st.info("Sin vistas de productos en el período.")
                    else:
                        df_prod = pd.DataFrame([
                            {
                                "Calidad": p.get("calidad", "—"),
                                "Producto": p["pagina"].replace("/productos/", "").rstrip("/"),
                                "Vistas": int(p["vistas"]),
                                "Engagement prom.": _fmt_seg(p["tiempo_prom_seg"]),
                                "Score": _fmt_seg(p.get("score_seg", 0.0)),
                            }
                            for p in top_prod
                        ])
                        st.dataframe(df_prod, hide_index=True, use_container_width=True)

                with col_pag:
                    st.subheader("Top páginas")
                    st.caption("Top 30 URLs más vistas (cualquier tipo). Scrolleable.")
                    top = ga4["top_paginas"]
                    if not top:
                        st.info("Sin datos.")
                    else:
                        df_top = pd.DataFrame([
                            {
                                "Calidad": p.get("calidad", "—"),
                                "Página": p["pagina"],
                                "Vistas": int(p["vistas"]),
                                "Engagement prom.": _fmt_seg(p["tiempo_prom_seg"]),
                                "Score": _fmt_seg(p.get("score_seg", 0.0)),
                            }
                            for p in top
                        ])
                        st.dataframe(df_top, hide_index=True, use_container_width=True, height=300)

                # Leyenda de calidad (común a las tablas de arriba y a la tabla larga)
                with st.expander("ℹ️ Qué significa cada nivel de Calidad / Score"):
                    st.markdown("""
**Score** = promedio bayesiano del engagement (tiempo prom. corregido por sample size).
Los items con pocas vistas son arrastrados al promedio global del sitio para evitar
ruido. Resultado: ítems con 1 vista y 5min NO se rankean como ganadores, y los de
mucho tráfico con engagement bajo NO zafan por volumen.

**Calidad — cuadrantes** (cruz de vistas vs score, mediana del set como umbral):

- 🟢 **Ganador** — mucho tráfico + buen engagement. Hacé MÁS de esto: más stock, pauta enfocada, destacar en home.
- 🔵 **Subdetectado** — poco tráfico + buen engagement. Producto que cuando lo encuentran les gusta. Falta visibilidad: subilo en home, pautalo, mejorá su SEO.
- 🟡 **Desperdiciado** — mucho tráfico + bajo engagement. La gente llega pero no engancha. Revisá foto, descripción, precio, comparativas. Si pautás este producto, estás regando dinero.
- ⚫ **Bajo perfil** — poco de ambas. No hay señal clara, no prioritario.
- ⏸ **Sin data** — vistas debajo del umbral mínimo. La muestra no alcanza para juzgar.
                    """.strip())

                st.divider()

                # ── Marcas (drill-down de /categoria/<marca>) ───────────────────
                st.subheader("Marcas")
                st.caption(
                    "Tráfico a cada landing de marca (`/categoria/<marca>`). "
                    "Elegí una marca para ver qué productos miran dentro."
                )
                marcas_lista = ga4.get("marcas") or []
                if not marcas_lista:
                    st.info("Sin vistas a páginas de marca en el período.")
                    _marca_seleccionada = None
                else:
                    total_marcas = sum(m["vistas"] for m in marcas_lista) or 1
                    df_marcas = pd.DataFrame([
                        {
                            "Marca": m["marca"].capitalize(),
                            "Vistas": int(m["vistas"]),
                            "% del total": f"{m['vistas'] / total_marcas * 100:.1f}%",
                            "Engagement prom.": _fmt_seg(m["tiempo_prom_seg"]),
                            "Engagement total": _fmt_seg(m["engagement_seg"]),
                        }
                        for m in marcas_lista
                    ])
                    st.dataframe(df_marcas, hide_index=True, use_container_width=True)

                    _opciones_marca = ["(Ver todas)"] + [m["marca"].capitalize() for m in marcas_lista]
                    _marca_sel = st.selectbox(
                        "Drill-down: productos de la marca",
                        _opciones_marca,
                        index=0,
                        key="ga4_marca_drill",
                    )
                    _marca_seleccionada = (
                        None if _marca_sel == "(Ver todas)" else _marca_sel.lower()
                    )

                    if _marca_seleccionada:
                        # Filtrar productos cuyo slug empieza con la marca
                        productos_marca = [
                            p for p in (ga4.get("productos_full") or [])
                            if p["path"].replace("/productos/", "").lower().startswith(_marca_seleccionada)
                        ]
                        productos_marca.sort(key=lambda x: x["vistas"], reverse=True)
                        if not productos_marca:
                            st.info(
                                f"No detectamos productos cuyo slug empiece con `{_marca_seleccionada}`. "
                                "Si los productos de esta marca no llevan el nombre de marca en el slug, "
                                "usá el buscador de la tabla de abajo."
                            )
                        else:
                            total_vistas_marca = sum(p["vistas"] for p in productos_marca)
                            engagement_total_marca = sum(p["engagement_seg"] for p in productos_marca)
                            tiempo_prom_marca = (
                                engagement_total_marca / total_vistas_marca
                                if total_vistas_marca > 0 else 0.0
                            )
                            st.caption(
                                f"**{len(productos_marca)} productos** de **{_marca_sel}** · "
                                f"Vistas totales: **{int(total_vistas_marca):,}** · "
                                f"Engagement prom.: **{_fmt_seg(tiempo_prom_marca)}**".replace(",", ".")
                            )
                            df_prod_marca = pd.DataFrame([
                                {
                                    "Producto": p["path"].replace("/productos/", "").rstrip("/"),
                                    "Título": p.get("titulo", ""),
                                    "Vistas": int(p["vistas"]),
                                    "% de la marca": (
                                        f"{p['vistas'] / total_vistas_marca * 100:.1f}%"
                                        if total_vistas_marca > 0 else "—"
                                    ),
                                    "Engagement prom.": _fmt_seg(p["tiempo_prom_seg"]),
                                }
                                for p in productos_marca
                            ])
                            st.dataframe(df_prod_marca, hide_index=True, use_container_width=True, height=350)

                st.divider()

                # ── Tabla larga de productos buscable, con calidad y filtros ───
                st.subheader("Explorar productos vistos")
                st.caption(
                    "Tabla completa de productos con tráfico. La columna Calidad clasifica "
                    "cada producto en uno de 5 cuadrantes accionables (ver leyenda más arriba)."
                )
                productos_full = ga4.get("productos_full") or []
                if not productos_full:
                    st.info("Sin productos vistos en el período.")
                else:
                    # Calidad ya fue computada arriba en _quality_by_path
                    col_f1, col_f2 = st.columns([2, 1])
                    with col_f1:
                        _filtro = st.text_input(
                            "Buscar producto (slug o título)",
                            placeholder="ej: anbernic, rg40, switch...",
                            key="ga4_prod_filtro",
                        ).strip().lower()
                    with col_f2:
                        _filtro_calidad = st.multiselect(
                            "Filtrar por calidad",
                            ["🟢 Ganador", "🔵 Subdetectado", "🟡 Desperdiciado", "⚫ Bajo perfil", "⏸ Sin data"],
                            default=[],
                            key="ga4_prod_calidad_filtro",
                            help="Vacío = todos. Tip: filtrá '🟡 Desperdiciado' para ver dónde corregir fichas, '🔵 Subdetectado' para detectar joyas ocultas.",
                        )

                    rows = []
                    for p in productos_full:
                        slug = p["path"].replace("/productos/", "").rstrip("/")
                        titulo = p.get("titulo", "")
                        if _filtro and _filtro not in slug.lower() and _filtro not in titulo.lower():
                            continue
                        if _filtro_calidad and p.get("calidad") not in _filtro_calidad:
                            continue
                        rows.append({
                            "Calidad": p.get("calidad", "—"),
                            "Producto": slug,
                            "Título": titulo,
                            "Vistas": int(p["vistas"]),
                            "Engagement prom.": _fmt_seg(p["tiempo_prom_seg"]),
                            "Score": _fmt_seg(p.get("score_seg", 0.0)),
                            "Engagement total": _fmt_seg(p["engagement_seg"]),
                        })
                    if not rows:
                        st.info("Ningún producto matchea los filtros.")
                    else:
                        df_pf = pd.DataFrame(rows)
                        # Resumen del set actual (con filtros aplicados)
                        from collections import Counter as _Counter
                        _conteo = _Counter(r["Calidad"] for r in rows)
                        _resumen = " · ".join(
                            f"{c} {_conteo.get(c, 0)}" for c in
                            ["🟢 Ganador", "🔵 Subdetectado", "🟡 Desperdiciado", "⚫ Bajo perfil", "⏸ Sin data"]
                            if _conteo.get(c, 0) > 0
                        )
                        st.caption(
                            f"**{len(rows)} producto{'s' if len(rows) != 1 else ''}** · {_resumen}"
                        )
                        st.dataframe(df_pf, hide_index=True, use_container_width=True, height=420)

                st.divider()

                # ── Búsquedas internas en el sitio ──────────────────────────────
                st.subheader("Búsquedas internas")
                st.caption(
                    "Qué buscan los usuarios DENTRO del sitio (buscador interno). "
                    "Oro puro: si buscan algo seguido que no tenés bien ranqueado, "
                    "te falta SEO interno o directamente el producto."
                )
                busquedas = ga4.get("busquedas_internas") or []
                urls_query = ga4.get("urls_con_query") or []
                if busquedas:
                    total_busq = sum(b["vistas"] for b in busquedas) or 1
                    df_busq = pd.DataFrame([
                        {
                            "Término buscado": b["termino"],
                            "Búsquedas": int(b["vistas"]),
                            "% del total": f"{b['vistas'] / total_busq * 100:.1f}%",
                        }
                        for b in busquedas
                    ])
                    st.dataframe(df_busq, hide_index=True, use_container_width=True, height=320)
                else:
                    st.info(
                        "No detectamos búsquedas internas con los formatos comunes "
                        "(`?q=`, `?_q=`, `?query=`, `?s=`). Abajo te muestro los top URLs "
                        "con query string que sí están llegando a GA4 — si ves algún patrón "
                        "que parezca búsqueda, decime el formato y lo agrego."
                    )

                # Fallback de debug: top URLs con query string (siempre visible si hay data)
                if urls_query:
                    with st.expander(
                        "🔬 Debug: top URLs con query string en el sitio "
                        f"({len(urls_query)} URLs detectadas)"
                    ):
                        st.caption(
                            "Te muestro las primeras URLs con `?` para identificar cómo trackea "
                            "TN tu buscador, paginación, filtros, etc."
                        )
                        df_urlq = pd.DataFrame([
                            {"URL": u["url"], "Vistas": int(u["vistas"])}
                            for u in urls_query
                        ])
                        st.dataframe(df_urlq, hide_index=True, use_container_width=True, height=300)

                st.divider()

                # ── Landing pages ───────────────────────────────────────────────
                st.subheader("Landing pages")
                st.caption("Por dónde entran a la web (primera página de la sesión).")
                landing = ga4["landing_pages"]
                if not landing:
                    st.info("Sin datos.")
                else:
                    df_land = pd.DataFrame([
                        {"Página de entrada": l["pagina"], "Sesiones": int(l["sesiones"])}
                        for l in landing
                    ])
                    st.dataframe(df_land, hide_index=True, use_container_width=True)

                # ── Eventos GA4 (investigación de add_to_cart / view_item) ──
                with st.expander("🔬 Eventos GA4 disponibles (investigación)"):
                    st.caption(
                        "Top 20 eventos que está disparando el sitio en GA4. "
                        "Si ves `add_to_cart`, `view_item` o `begin_checkout` con volumen real, "
                        "Tiendanube está enviando los eventos de ecommerce y podemos atribuir "
                        "interés real por producto (no solo vistas). Si solo ves `page_view`, "
                        "`session_start`, `user_engagement`, falta configurar el datalayer de TN."
                    )
                    eventos = ga4.get("eventos") or []
                    if not eventos:
                        st.info("Sin eventos en el período.")
                    else:
                        df_ev = pd.DataFrame([
                            {"Evento": e["evento"], "Count": int(e["count"])}
                            for e in eventos
                        ])
                        st.dataframe(df_ev, hide_index=True, use_container_width=True)

    # ══════════════════════════════════════════════════════════════════════════
    # TAB 12: ANALISTA IA
    # ══════════════════════════════════════════════════════════════════════════
    elif seccion == "🤖 Analista IA":
        if "analyst_messages" not in st.session_state:
            st.session_state.analyst_messages = []

        st.subheader("🤖 Analista Financiero IA")
        st.caption(
            "Preguntale lo que quieras sobre el negocio. "
            "Usa los datos reales del período seleccionado."
        )

        if not ANTHROPIC_KEY:
            st.error(
                "⚠️ Falta la API key de Anthropic.\n\n"
                "Agregá `ANTHROPIC_KEY = \"sk-ant-...\"` en `secrets.toml`."
            )
            st.info("💡 Obtené tu key en [console.anthropic.com](https://console.anthropic.com/)")
        else:
            # ── Construir contexto de datos ──
            def _build_analyst_context():
                lines = []
                _dias_p = max((fecha_hasta - fecha_desde).days + 1, 1)
                _tc = int(dolar_blue) if dolar_blue else 1200
                _gastos_gs = gs_read("GastosFijos") or {}
                _costos_gs = st.session_state.get("costos_consolas") or gs_read("CostosConsolas") or {}

                lines.append(f"=== MARKET GAMER — ANÁLISIS FINANCIERO ===")
                lines.append(f"Período: {fecha_desde.strftime('%d/%m/%Y')} → {fecha_hasta.strftime('%d/%m/%Y')} ({_dias_p} días)")
                lines.append(f"Tipo de cambio: ${_tc:,} ARS/USD (dólar blue)")
                lines.append("")

                if df_tn.empty:
                    lines.append("No hay órdenes cargadas en este período.")
                    return "\n".join(lines)

                _total_ordenes = len(df_tn)
                _fact_bruta = df_tn["Total ($)"].sum()
                _total_comision = df_tn["Comision PN ($)"].sum()
                _neto = df_tn["Neto cobrado ($)"].sum()
                _envios = df_tn["Envio costo ($)"].sum()
                _ticket_prom = _fact_bruta / _total_ordenes if _total_ordenes > 0 else 0

                lines.append("--- RESUMEN GENERAL ---")
                lines.append(f"Órdenes: {_total_ordenes}")
                lines.append(f"Facturación bruta: ${_fact_bruta:,.0f}")
                if _fact_bruta > 0:
                    lines.append(f"Comisiones PN: ${_total_comision:,.0f} ({_total_comision/_fact_bruta*100:.2f}%)")
                lines.append(f"Neto cobrado: ${_neto:,.0f}")
                lines.append(f"Costo envíos dueño: ${_envios:,.0f}")
                lines.append(f"Ticket promedio: ${_ticket_prom:,.0f}")
                lines.append(f"Órdenes/día: {_total_ordenes / _dias_p:.1f}")
                lines.append(f"Facturación/día: ${_fact_bruta / _dias_p:,.0f}")
                lines.append("")

                # Ventas por producto
                lines.append("--- VENTAS POR PRODUCTO ---")
                _prod_stats = {}
                for _, _row in df_tn.iterrows():
                    _prods = [p.strip() for p in str(_row.get("Productos", "")).split(" / ") if p.strip()]
                    _n = max(len(_prods), 1)
                    for _p in _prods:
                        if _p not in _prod_stats:
                            _prod_stats[_p] = {"uds": 0, "rev": 0.0, "ords": 0}
                        _prod_stats[_p]["uds"] += int(_row.get("Cantidad", 1) or 1)
                        _prod_stats[_p]["rev"] += _row.get("Total ($)", 0) / _n
                        _prod_stats[_p]["ords"] += 1

                _orders_raw = st.session_state.orders_raw
                _prod_rows = _build_product_rows_from_raw(_orders_raw) if _orders_raw else []
                _margen_map = {}
                if _prod_rows:
                    _df_pr = pd.DataFrame(_prod_rows)
                    for _pn, _grp in _df_pr.groupby("Producto"):
                        _precio_p = _grp["Precio ($)"].mean()
                        _costo_usd = get_costo_total_usd(_pn, _costos_gs)
                        _costo_ars = _costo_usd * _tc
                        _costo_full = _costo_ars + 2500 + (_precio_p * 0.105)
                        _com_p = _grp["Comisión PN ($)"].mean()
                        _env_p = _grp["Envío ($)"].mean()
                        _margen_p = _precio_p - _costo_full - _com_p - _env_p
                        _margen_pct = (_margen_p / _precio_p * 100) if _precio_p > 0 else 0
                        _margen_map[_pn] = {
                            "precio": _precio_p, "costo_usd": _costo_usd,
                            "costo_full": _costo_full, "margen": _margen_p, "margen_pct": _margen_pct,
                        }

                lines.append(f"{'Producto':<45} {'Uds':>5} {'Revenue':>12} {'Margen%':>8} {'FOB$':>7}")
                lines.append("-" * 82)
                for _pn, _ps in sorted(_prod_stats.items(), key=lambda x: x[1]["rev"], reverse=True):
                    _mi = _margen_map.get(_pn, {})
                    _mpct = f"{_mi.get('margen_pct', 0):.1f}%" if _mi else "s/d"
                    _fob = f"${_mi.get('costo_usd', 0):.0f}" if _mi.get('costo_usd') else "s/d"
                    lines.append(f"{_pn:<45} {_ps['uds']:>5} ${_ps['rev']:>10,.0f} {_mpct:>8} {_fob:>7}")
                lines.append("")

                # Evolución diaria
                lines.append("--- EVOLUCIÓN DIARIA ---")
                _df_dia = df_tn.groupby("Fecha").agg(
                    Ordenes=("Orden", "count"), Facturacion=("Total ($)", "sum"), Unidades=("Cantidad", "sum"),
                ).reset_index().sort_values("Fecha")
                for _, _r in _df_dia.iterrows():
                    lines.append(f"{_r['Fecha']:<12} {int(_r['Ordenes']):>5} órd  ${_r['Facturacion']:>12,.0f}  {int(_r['Unidades']):>4} uds")
                if len(_df_dia) >= 6:
                    _1h = _df_dia.head(len(_df_dia) // 2)["Facturacion"].mean()
                    _2h = _df_dia.tail(len(_df_dia) // 2)["Facturacion"].mean()
                    if _2h > _1h * 1.1:
                        lines.append(f"📈 TENDENCIA ALCISTA (+{(_2h/_1h - 1)*100:.0f}%)")
                    elif _2h < _1h * 0.9:
                        lines.append(f"📉 TENDENCIA BAJISTA ({(_2h/_1h - 1)*100:.0f}%)")
                    else:
                        lines.append("➡️ TENDENCIA ESTABLE")
                lines.append("")

                # Medios de pago
                lines.append("--- MEDIOS DE PAGO ---")
                _med = df_tn.groupby("Medio de Pago").agg(
                    Ordenes=("Orden", "count"), Fact=("Total ($)", "sum"), Com=("Comision PN ($)", "sum"),
                ).reset_index().sort_values("Fact", ascending=False)
                for _, _r in _med.iterrows():
                    _cp = (_r["Com"] / _r["Fact"] * 100) if _r["Fact"] > 0 else 0
                    lines.append(f"{_r['Medio de Pago']:<25} {int(_r['Ordenes']):>5} órd  ${_r['Fact']:>12,.0f}  com {_cp:.2f}%")
                lines.append("")

                # Cuotas
                lines.append("--- DISTRIBUCIÓN POR CUOTAS ---")
                _cuotas = df_tn.groupby("Cuotas").agg(Ordenes=("Orden", "count"), Fact=("Total ($)", "sum")).reset_index()
                for _, _r in _cuotas.iterrows():
                    lines.append(f"  {int(_r['Cuotas'])} cuota(s): {int(_r['Ordenes'])} órdenes — ${_r['Fact']:,.0f}")
                lines.append("")

                # Gastos fijos
                lines.append("--- GASTOS FIJOS MENSUALES ---")
                _total_gf = 0
                if _gastos_gs:
                    for _k, _v in _gastos_gs.items():
                        if isinstance(_v, (int, float)) and _v > 0:
                            lines.append(f"  {_k}: ${_v:,.0f}")
                            _total_gf += _v
                    lines.append(f"  TOTAL mensual: ${_total_gf:,.0f}")
                    lines.append(f"  Prorrateado ({_dias_p}d): ${round(_total_gf * _dias_p / 30):,.0f}")
                lines.append("")

                # Resultado financiero
                lines.append("--- RESULTADO FINANCIERO ---")
                _costo_prods = 0
                for _pn, _ps in _prod_stats.items():
                    _costo_prods += get_costo_total_usd(_pn, _costos_gs) * _tc * _ps["uds"]
                _iva = _fact_bruta * 0.105
                _gf_per = round(_total_gf * _dias_p / 30)
                _margen_b = _neto - _costo_prods - _envios
                _resultado = _margen_b - _iva - _gf_per

                lines.append(f"  Facturación bruta:  ${_fact_bruta:>12,.0f}")
                lines.append(f"  - Comisiones PN:    ${_total_comision:>12,.0f}")
                lines.append(f"  = Neto cobrado:     ${_neto:>12,.0f}")
                lines.append(f"  - Costo productos:  ${_costo_prods:>12,.0f}")
                lines.append(f"  - Costo envíos:     ${_envios:>12,.0f}")
                lines.append(f"  = Margen bruto:     ${_margen_b:>12,.0f}")
                lines.append(f"  - IVA (10.5%):      ${_iva:>12,.0f}")
                lines.append(f"  - Gastos fijos:     ${_gf_per:>12,.0f}")
                lines.append(f"  = RESULTADO:        ${_resultado:>12,.0f} {'✅' if _resultado >= 0 else '🔴'}")
                if _fact_bruta > 0:
                    lines.append(f"  Margen neto/bruto:  {_resultado/_fact_bruta*100:.1f}%")
                lines.append("")

                # Stock
                _stock_df = st.session_state.get("stock_tn")
                if _stock_df is not None and not _stock_df.empty:
                    lines.append("--- STOCK ACTUAL ---")
                    _stock_sum = _stock_df.groupby("Producto").agg(
                        Stock=("Stock", lambda x: sum(v for v in x if isinstance(v, (int, float)))),
                    ).reset_index().sort_values("Stock")
                    for _, _r in _stock_sum.iterrows():
                        _pn = _r["Producto"]
                        _st = int(_r["Stock"])
                        _vel = _prod_stats.get(_pn, {}).get("uds", 0) / _dias_p
                        _dr = int(_st / _vel) if _vel > 0 else 999
                        _alert = " ⚠️" if _dr <= 14 else ""
                        lines.append(f"{_pn:<45} stock:{_st:>4}  vel:{_vel:.2f}/día  {_dr}d{_alert}")
                    lines.append("")

                # Costos cargados
                if _costos_gs:
                    lines.append("--- COSTOS DE CONSOLAS ---")
                    _ckg = float(_costos_gs.get("_costo_kg_usd", 65.0) or 65.0)
                    for _k, _v in sorted(_costos_gs.items()):
                        if _k.startswith("_") or not isinstance(_v, dict):
                            continue
                        _fob = float(_v.get("fob_usd", 0) or 0)
                        _tot = float(_v.get("costo_total_usd", 0) or 0)
                        if _tot == 0 and _fob > 0:
                            _tot = _fob + float(_v.get("peso_kg", 0) or 0) * _ckg
                        if _fob > 0:
                            lines.append(f"{_k:<40} FOB ${_fob:.0f}  Total ${_tot:.0f}  ARS ${_tot*_tc:,.0f}")
                    lines.append("")

                # Top clientes
                lines.append("--- TOP 10 CLIENTES ---")
                _top_cli = df_tn.groupby("Cliente").agg(
                    Ords=("Orden", "count"), Fact=("Total ($)", "sum"),
                ).reset_index().sort_values("Fact", ascending=False).head(10)
                for _, _r in _top_cli.iterrows():
                    lines.append(f"  {_r['Cliente']}: {int(_r['Ords'])} ord — ${_r['Fact']:,.0f}")
                lines.append("")

                # Días de la semana
                lines.append("--- VENTAS POR DÍA DE SEMANA ---")
                try:
                    _dfw = df_tn.copy()
                    _dfw["DOW"] = pd.to_datetime(_dfw["Fecha"]).dt.day_name()
                    _labels = {"Monday": "Lunes", "Tuesday": "Martes", "Wednesday": "Miércoles",
                               "Thursday": "Jueves", "Friday": "Viernes", "Saturday": "Sábado", "Sunday": "Domingo"}
                    _dow = _dfw.groupby("DOW").agg(Ords=("Orden", "count"), Fact=("Total ($)", "sum")).reset_index()
                    _dow["DOW"] = pd.Categorical(_dow["DOW"], categories=["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"], ordered=True)
                    for _, _r in _dow.sort_values("DOW").iterrows():
                        lines.append(f"  {_labels.get(_r['DOW'], _r['DOW'])}: {int(_r['Ords'])} órd — ${_r['Fact']:,.0f}")
                except Exception:
                    pass

                return "\n".join(lines)

            # ── Llamar a Claude API ──
            def _call_analyst(question, history=None):
                context = _build_analyst_context()
                system_prompt = f"""Sos el analista financiero de Market Gamer, un e-commerce argentino de consolas retro portátiles (Anbernic, Powkiddy, Trimui, Miyoo).

Tu trabajo es analizar los datos del negocio y dar insights accionables. Respondé en español argentino, directo y sin relleno.

REGLAS:
- Basate EXCLUSIVAMENTE en los datos proporcionados. No inventes números.
- Formato argentino para montos ($XXX.XXX).
- Detectá problemas y oportunidades con recomendaciones concretas.
- Si no hay datos suficientes, decilo.
- Priorizá insights accionables sobre descripciones obvias.
- Emojis para alertas: 🔴 problema, 🟡 atención, 🟢 bien, 💡 oportunidad.
- Si piden análisis general: performance, productos estrella, alertas, medios de pago, 3 recomendaciones.

DATOS DEL NEGOCIO:
{context}"""

                messages = []
                if history:
                    for msg in history:
                        messages.append({"role": msg["role"], "content": msg["content"]})
                messages.append({"role": "user", "content": question})

                try:
                    response = requests.post(
                        "https://api.anthropic.com/v1/messages",
                        headers={
                            "Content-Type": "application/json",
                            "x-api-key": ANTHROPIC_KEY,
                            "anthropic-version": "2023-06-01",
                        },
                        json={
                            "model": "claude-sonnet-4-6",
                            "max_tokens": 4096,
                            "system": system_prompt,
                            "messages": messages,
                        },
                        timeout=60,
                    )
                    if response.status_code == 200:
                        data = response.json()
                        return data["content"][0]["text"]
                    elif response.status_code == 401:
                        return "🔴 API key inválida. Revisá ANTHROPIC_KEY en secrets."
                    elif response.status_code == 429:
                        return "🟡 Rate limit. Esperá unos segundos."
                    else:
                        return f"🔴 Error ({response.status_code}): {response.text[:200]}"
                except requests.exceptions.Timeout:
                    return "🟡 Timeout. Probá con algo más específico."
                except Exception as e:
                    return f"🔴 Error de conexión: {str(e)}"

            # ── Preguntas rápidas predefinidas ──
            PREGUNTAS_RAPIDAS = {
                "📊 Resumen ejecutivo": (
                    "Haceme un resumen ejecutivo completo del período. Incluí: "
                    "1) Performance general (facturación, tendencia, ticket promedio), "
                    "2) Top 5 productos por revenue y por margen, "
                    "3) Alertas críticas (stock, márgenes negativos, productos que no se venden), "
                    "4) Análisis de medios de pago y su impacto en rentabilidad, "
                    "5) 3 recomendaciones priorizadas con impacto estimado."
                ),
                "🏆 Productos estrella": (
                    "Clasificá los productos en: ESTRELLAS (alto volumen + alto margen), "
                    "VACAS (alto volumen + bajo margen), OPORTUNIDADES (bajo volumen + alto margen), "
                    "PERROS (bajo volumen + bajo margen). Para cada uno, qué acción tomar."
                ),
                "💳 Impacto cuotas": (
                    "Analizá cómo los medios de pago impactan en rentabilidad: "
                    "% en cuotas vs débito, margen perdido por cuotas, productos con margen negativo en cuotas, "
                    "y si conviene incentivar transferencia/débito con descuento."
                ),
                "📦 Riesgo stock": (
                    "Analizá riesgo de stock: qué se agota primero, sobrestock, "
                    "qué pedir urgente con cantidades para 30 días."
                ),
                "📈 Tendencia": (
                    "Analizá tendencias: ventas subiendo/bajando, mejores días de semana, "
                    "productos ganando/perdiendo tracción, proyección mensual, "
                    "órdenes/día necesarias para cubrir gastos fijos."
                ),
                "💰 Punto equilibrio": (
                    "Calculá punto de equilibrio: facturación mínima mensual, "
                    "unidades necesarias, órdenes/día mínimas, "
                    "y si estoy arriba o abajo en este período."
                ),
                "🔍 Márgenes": (
                    "Diagnóstico de márgenes: margen promedio ponderado, productos con margen <10%, "
                    "desglose top 5 (precio/costo/comisión/IVA/envío/margen), "
                    "dónde se escapa la plata, impacto de subir precios 5-10%."
                ),
            }

            # ── Preview de datos ──
            with st.expander("📋 Ver datos que recibe el analista", expanded=False):
                _ctx_preview = _build_analyst_context()
                st.text(_ctx_preview)
                st.caption(f"📏 {len(_ctx_preview)} caracteres · ~{len(_ctx_preview)//4} tokens")

            # ── Botones de análisis rápido ──
            st.markdown("### ⚡ Análisis rápidos")

            _cols1 = st.columns(4)
            _cols2 = st.columns(4)
            _all_cols = _cols1 + _cols2
            _items = list(PREGUNTAS_RAPIDAS.items())

            for _i, (_label, _question) in enumerate(_items):
                if _i < len(_all_cols):
                    with _all_cols[_i]:
                        if st.button(_label, key=f"quick_{_i}", use_container_width=True):
                            st.session_state.analyst_messages.append(
                                {"role": "user", "content": _question}
                            )
                            with st.spinner("🤖 Analizando..."):
                                _resp = _call_analyst(
                                    _question,
                                    st.session_state.analyst_messages[:-1],
                                )
                            st.session_state.analyst_messages.append(
                                {"role": "assistant", "content": _resp}
                            )
                            st.rerun()

            st.divider()

            # ── Chat ──
            st.markdown("### 💬 Conversación")

            if not st.session_state.analyst_messages:
                st.info(
                    "👆 Usá un análisis rápido o escribí tu pregunta abajo.\n\n"
                    "**Ejemplos:** *¿Cuál es mi producto más rentable?* · "
                    "*¿Cuánto me cuesta Pago Nube en promedio?* · "
                    "*¿Qué consola debería dejar de vender?* · "
                    "*Si subo precios 10%, ¿cómo cambia el resultado?*"
                )

            for _msg in st.session_state.analyst_messages:
                with st.chat_message(_msg["role"], avatar="🧑‍💼" if _msg["role"] == "user" else "🤖"):
                    st.markdown(_msg["content"])

            _user_input = st.chat_input("Preguntale al analista...")

            if _user_input:
                st.session_state.analyst_messages.append(
                    {"role": "user", "content": _user_input}
                )
                with st.chat_message("user", avatar="🧑‍💼"):
                    st.markdown(_user_input)
                with st.chat_message("assistant", avatar="🤖"):
                    with st.spinner("🤖 Analizando datos..."):
                        _resp = _call_analyst(
                            _user_input,
                            st.session_state.analyst_messages[:-1],
                        )
                    st.markdown(_resp)
                st.session_state.analyst_messages.append(
                    {"role": "assistant", "content": _resp}
                )

            # ── Controles ──
            st.divider()
            _cc1, _cc2, _cc3 = st.columns(3)
            with _cc1:
                if st.button("🗑️ Limpiar conversación", use_container_width=True):
                    st.session_state.analyst_messages = []
                    st.rerun()
            with _cc2:
                if st.session_state.analyst_messages:
                    _chat_txt = "\n\n".join(
                        f"{'👤 Bruno' if m['role'] == 'user' else '🤖 Analista'}: {m['content']}"
                        for m in st.session_state.analyst_messages
                    )
                    st.download_button(
                        "⬇️ Exportar conversación",
                        _chat_txt.encode("utf-8"),
                        f"analisis_mg_{fecha_desde}_{fecha_hasta}.txt",
                        "text/plain",
                        use_container_width=True,
                    )
            with _cc3:
                if st.button("🔄 Refrescar datos", use_container_width=True):
                    st.session_state.analyst_context = ""
                    st.rerun()

else:
    st.info("⏳ Cargando datos...")
